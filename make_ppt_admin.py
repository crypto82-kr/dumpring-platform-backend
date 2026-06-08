import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# 컬러 팔레트 정의 (라이트 머스터드 옐로우 테마)
WHITE = RGBColor(255, 255, 255)
SAND_BG = RGBColor(249, 248, 246)
MUSTARD = RGBColor(245, 158, 11)     # F59E0B (주요 테마 색상)
AMBER = RGBColor(217, 119, 6)        # D97706 (포인트 색상)
CHARCOAL = RGBColor(31, 41, 55)      # 1F2937 (메인 텍스트)
LIGHT_GRAY = RGBColor(243, 244, 246)  # F3F4F6 (영역 구분 배경)
BORDER_GRAY = RGBColor(229, 231, 235) # E5E7EB (카드 테두리)
SLATE_GRAY = RGBColor(75, 85, 99)    # 4B5563 (서브 텍스트)
RED = RGBColor(239, 68, 68)          # EF4444 (반려/에러/취소/회차/중단)
GREEN = RGBColor(16, 185, 129)       # 10B981 (정산완료/하차완료/승인/운영중)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6] # 빈 레이아웃

# 유틸리티 함수: 텍스트 상자 추가
def add_textbox(slide, left, top, width, height, text, font_size=11, bold=False, color=CHARCOAL, align=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = 'Malgun Gothic'
    if align:
        p.alignment = align
    return txBox

# 유틸리티 함수: 직사각형 도형 추가
def add_rectangle(slide, left, top, width, height, fill_color, line_color=None, line_width=1):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

# 유틸리티 함수: 둥근 직사각형 도형 추가
def add_rounded_rect(slide, left, top, width, height, fill_color, line_color=None, line_width=1):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

# 유틸리티 함수: 표 형태의 마스터 레이아웃 (페이지별 제목 상단 한 줄, 좌 화면설계 : 우 설명 = 4:1 분할)
def add_master_table_layout(slide, title_text):
    # 1. 페이지별 제목 상단 한 줄 표시
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.6), title_text, 20, True, CHARCOAL)
    
    # 2. 표 외곽 테두리 (둥근 테두리 제거, 일반 직사각형)
    add_rectangle(slide, Inches(0.5), Inches(1.1), Inches(12.33), Inches(5.8), WHITE, BORDER_GRAY, 1)
    
    # 3. 좌측 '화면 설계' 헤더 배경
    add_rectangle(slide, Inches(0.5), Inches(1.1), Inches(9.8), Inches(0.4), LIGHT_GRAY, BORDER_GRAY, 1)
    add_textbox(slide, Inches(0.5), Inches(1.16), Inches(9.8), Inches(0.35), "화면 설계 (UI Mockup)", 11, True, CHARCOAL, PP_ALIGN.CENTER)
    
    # 4. 우측 '화면 설명' 헤더 배경
    add_rectangle(slide, Inches(10.3), Inches(1.1), Inches(2.53), Inches(0.4), LIGHT_GRAY, BORDER_GRAY, 1)
    add_textbox(slide, Inches(10.3), Inches(1.16), Inches(2.53), Inches(0.35), "화면 설명 및 상세 기능 명세", 11, True, CHARCOAL, PP_ALIGN.CENTER)
    
    # 5. 표 수직 구분선 (좌우 4:1 경계선)
    add_rectangle(slide, Inches(10.3), Inches(1.1), Inches(0.02), Inches(5.8), BORDER_GRAY)
    
    # 6. 표 헤더 수평 구분선
    add_rectangle(slide, Inches(0.5), Inches(1.5), Inches(12.33), Inches(0.02), BORDER_GRAY)


# 유틸리티 함수: 어드민 웹 콘솔 프레임 묘사 (상단 브라우저 헤더 및 빈 공백 제거)
def add_web_frame(slide, left, top, width, height, active_tab="대시보드"):
    # 전체 어드민 콘솔 외부 테두리
    add_rectangle(slide, left, top, width, height, WHITE, BORDER_GRAY, 1)
    
    # 좌측 사이드바 LNB
    sb_w = Inches(1.8)
    add_rectangle(slide, left, top, sb_w, height, LIGHT_GRAY, BORDER_GRAY)
    
    # 사이드바 메뉴 렌더링 (현장 관리 추가)
    menus = ["대시보드", "회원 심사", "현장 관리", "배차 제어", "분쟁 조정"]
    for i, menu in enumerate(menus):
        menu_y = top + Inches(0.3) + Inches(0.5) * i
        is_active = (menu == active_tab)
        if is_active:
            add_rectangle(slide, left, menu_y, sb_w, Inches(0.4), MUSTARD)
            add_textbox(slide, left + Inches(0.15), menu_y + Inches(0.08), sb_w - Inches(0.2), Inches(0.35), menu, 10, True, CHARCOAL)
        else:
            add_textbox(slide, left + Inches(0.15), menu_y + Inches(0.08), sb_w - Inches(0.2), Inches(0.35), menu, 10, False, SLATE_GRAY)
            
    return left + sb_w, top, width - sb_w, height


# =============================================================
# SLIDE 1: 커버 페이지 (표 형식 예외 슬라이드)
# =============================================================
slide1 = prs.slides.add_slide(blank_layout)
add_rectangle(slide1, Inches(0), Inches(0), Inches(13.33), Inches(7.5), SAND_BG)
add_rectangle(slide1, Inches(1.5), Inches(2.2), Inches(0.15), Inches(2.5), MUSTARD)
add_textbox(slide1, Inches(2.0), Inches(2.1), Inches(9.5), Inches(1.2), "덤프링 (Dumpring)", 44, True, CHARCOAL)
add_textbox(slide1, Inches(2.0), Inches(3.2), Inches(9.5), Inches(0.8), "07. 플랫폼 어드민 (Admin) 웹 화면 상세설계서", 28, True, SLATE_GRAY)
add_textbox(slide1, Inches(2.0), Inches(4.1), Inches(9.5), Inches(0.5), "회원/현장/하차지/배차 대행등록 및 CRUD 수명주기 통제판", 16, False, AMBER)
add_textbox(slide1, Inches(2.0), Inches(5.8), Inches(6.0), Inches(0.8), "작성일: 2026. 06. 05\n작성자: 덤프링 플랫폼 기획팀", 12, False, SLATE_GRAY)


# =============================================================
# SLIDE 2: 어드민 홈 (플랫폼 관제 대시보드)
# =============================================================
slide2 = prs.slides.add_slide(blank_layout)
add_rectangle(slide2, Inches(0), Inches(0), Inches(13.33), Inches(7.5), WHITE)

# 1. 마스터 테이블 레이아웃 추가
add_master_table_layout(slide2, "01. 플랫폼 종합 관제 대시보드 (Admin Home)")

# 2. 우측 설명 기입
add_textbox(slide2, Inches(10.4), Inches(1.6), Inches(2.3), Inches(5.0),
            "■ 대시보드 주요 기능\n"
            "- 플랫폼 내 실시간 운송/매칭 핵심 지표를 통합 모니터링합니다.\n"
            "- 긴급 조치가 필요한 회원가입 심사 및 미해결 분쟁 알림을 최우선 배치합니다.\n\n"
            "■ UI 핵심 명세\n"
            "1. 활성 콜 카드:\n"
            "   - 현재 매칭 중이거나 주행 상태인 콜 42건 표시.\n"
            "2. 미해결 분쟁 카드:\n"
            "   - 하차지 거부로 유보된 긴급 건수를 빨간색으로 시각화.\n"
            "3. 실시간 긴급 경고 배너:\n"
            "   - 티켓 #2043 건의 뻘흙 거부 이의제기를 헤드라인 노출.",
            9, False, CHARCOAL)

# 3. 좌측 화면설계 영역 내 웹 콘솔 배치 (x=0.8, y=1.7, w=9.2, h=5.0)
mc_l, mc_t, mc_w, mc_h = add_web_frame(slide2, Inches(0.8), Inches(1.7), Inches(9.2), Inches(5.0), "대시보드")

# 상단 카드 3종
card_w = Inches(2.1)
card_h = Inches(0.9)
card_y = mc_t + Inches(0.2)

# 카드 1
add_rectangle(slide2, mc_l + Inches(0.2), card_y, card_w, card_h, WHITE, BORDER_GRAY)
add_textbox(slide2, mc_l + Inches(0.3), card_y + Inches(0.1), card_w - Inches(0.2), Inches(0.25), "실시간 활성 콜", 9, False, SLATE_GRAY)
add_textbox(slide2, mc_l + Inches(0.3), card_y + Inches(0.35), card_w - Inches(0.2), Inches(0.4), "42 건", 14, True, MUSTARD)

# 카드 2
add_rectangle(slide2, mc_l + Inches(2.5), card_y, card_w, card_h, WHITE, BORDER_GRAY)
add_textbox(slide2, mc_l + Inches(2.6), card_y + Inches(0.1), card_w - Inches(0.2), Inches(0.25), "심사 대기 서류", 9, False, SLATE_GRAY)
add_textbox(slide2, mc_l + Inches(2.6), card_y + Inches(0.35), card_w - Inches(0.2), Inches(0.4), "8 건", 14, True, AMBER)

# 카드 3 (긴급 분쟁)
add_rectangle(slide2, mc_l + Inches(4.8), card_y, card_w, card_h, WHITE, BORDER_GRAY)
add_textbox(slide2, mc_l + Inches(4.9), card_y + Inches(0.1), card_w - Inches(0.2), Inches(0.25), "미해결 분쟁", 9, False, SLATE_GRAY)
add_textbox(slide2, mc_l + Inches(4.9), card_y + Inches(0.35), card_w - Inches(0.2), Inches(0.4), "3 건", 14, True, RED)

# 긴급 알림판 배너
banner_y = card_y + Inches(1.0)
banner_w = mc_w - Inches(0.4)
add_rectangle(slide2, mc_l + Inches(0.2), banner_y, banner_w, Inches(0.7), WHITE, RED, 1.5)
add_textbox(slide2, mc_l + Inches(0.35), banner_y + Inches(0.08), banner_w - Inches(0.3), Inches(0.3), "🚨 긴급 알림: 티켓 #2043 분쟁 접수 (김덤프 기사 ➔ 영종 사토장)", 10, True, RED)
add_textbox(slide2, mc_l + Inches(0.35), banner_y + Inches(0.35), banner_w - Inches(0.3), Inches(0.3), "반입 거부 사유: 뻘흙 의심 판정  |  기사 주장: 양질의 마사토", 8.5, False, SLATE_GRAY)

# 트렌드 차트 가상 뷰
chart_y = banner_y + Inches(0.85)
chart_w = mc_w - Inches(0.4)
add_rectangle(slide2, mc_l + Inches(0.2), chart_y, chart_w, Inches(1.8), WHITE, BORDER_GRAY)
add_textbox(slide2, mc_l + Inches(0.35), chart_y + Inches(0.1), chart_w - Inches(0.3), Inches(0.35), "일일 배차 매칭 추이 (통계)", 11, True, CHARCOAL)
add_rectangle(slide2, mc_l + Inches(1.5), chart_y + Inches(0.9), Inches(4.0), Inches(0.04), BORDER_GRAY)
add_rectangle(slide2, mc_l + Inches(2.2), chart_y + Inches(0.7), Inches(0.12), Inches(0.12), MUSTARD)
add_rectangle(slide2, mc_l + Inches(3.5), chart_y + Inches(0.5), Inches(0.12), Inches(0.12), AMBER)


# =============================================================
# SLIDE 3: 회원 수동 등록 및 자격 심사 콘솔
# =============================================================
slide3 = prs.slides.add_slide(blank_layout)
add_rectangle(slide3, Inches(0), Inches(0), Inches(13.33), Inches(7.5), WHITE)

# 1. 마스터 테이블 레이아웃 추가
add_master_table_layout(slide3, "02. 회원 조회 및 대행 등록/심사 콘솔")

# 2. 우측 설명 기입
add_textbox(slide3, Inches(10.4), Inches(1.6), Inches(2.3), Inches(5.0),
            "■ 회원 관리 주요 기능\n"
            "- 플랫폼 직원이 회원을 수동으로 직접 등록(대행 등록)할 수 있습니다.\n"
            "- 가입 서류 검증(면허증, 사업자사본) 승인/기각을 처리합니다.\n\n"
            "■ UI 핵심 명세\n"
            "1. 회원 수동 등록 버튼:\n"
            "   - 상단에 회원 등록 실행 단추 추가.\n"
            "2. 전체 회원 테이블 목록:\n"
            "   - 가입일, 구분, 연락처, 활동 정지 등 상태 표시.\n"
            "3. 우측 프로필 정보 수정:\n"
            "   - 계좌 변경 등 대행 편집 폼 구성.",
            9, False, CHARCOAL)

# 3. 좌측 화면설계 영역 내 웹 콘솔 배치
mc_l, mc_t, mc_w, mc_h = add_web_frame(slide3, Inches(0.8), Inches(1.7), Inches(9.2), Inches(5.0), "회원 심사")

# 상단 수동 등록 버튼 및 큐 영역
reg_btn = add_rectangle(slide3, mc_l + Inches(0.2), mc_t + Inches(0.2), Inches(1.5), Inches(0.4), MUSTARD)
add_textbox(slide3, mc_l + Inches(0.2), mc_t + Inches(0.28), Inches(1.5), Inches(0.3), "➕ 회원 수동 등록", 9.5, True, CHARCOAL, PP_ALIGN.CENTER)

# 좌측 테이블 목록
q_w = Inches(2.2)
add_rectangle(slide3, mc_l + Inches(0.2), mc_t + Inches(0.7), q_w, mc_h - Inches(0.9), WHITE, BORDER_GRAY)
add_textbox(slide3, mc_l + Inches(0.3), mc_t + Inches(0.8), q_w - Inches(0.2), Inches(0.3), "전체 가입 회원 (128)", 10, True, SLATE_GRAY)

# 대기 1 (활성 선택)
add_rectangle(slide3, mc_l + Inches(0.2), mc_t + Inches(1.15), q_w, Inches(0.65), MUSTARD)
add_textbox(slide3, mc_l + Inches(0.3), mc_t + Inches(1.2), q_w - Inches(0.2), Inches(0.35), "이차주 (차주)", 10, True, CHARCOAL)
add_textbox(slide3, mc_l + Inches(0.3), mc_t + Inches(1.48), q_w - Inches(0.2), Inches(0.25), "활동 상태: 정상 🟢", 8, False, SLATE_GRAY)

# 대기 2
add_rectangle(slide3, mc_l + Inches(0.2), mc_t + Inches(1.85), q_w, Inches(0.65), WHITE, BORDER_GRAY)
add_textbox(slide3, mc_l + Inches(0.3), mc_t + Inches(1.9), q_w - Inches(0.2), Inches(0.35), "김기사 (기사)", 10, False, CHARCOAL)

# 상세 콘텐츠 영역
det_l = mc_l + q_w + Inches(0.3)
det_w = mc_w - q_w - Inches(0.5)

add_textbox(slide3, det_l, mc_t + Inches(0.2), det_w, Inches(0.35), "회원 마스터 정보 편집 및 자격 서류", 11, True, CHARCOAL)

# 서류 이미지 뷰어
add_rectangle(slide3, det_l, mc_t + Inches(0.65), det_w - Inches(1.8), Inches(2.4), LIGHT_GRAY, BORDER_GRAY)
add_textbox(slide3, det_l + Inches(0.1), mc_t + Inches(1.6), det_w - Inches(2.0), Inches(0.5), "[제출 서류 사본 이미지 뷰어]", 9, False, SLATE_GRAY, PP_ALIGN.CENTER)

# 입력 정보
add_rectangle(slide3, det_l + det_w - Inches(1.7), mc_t + Inches(0.65), Inches(1.7), Inches(2.4), WHITE, BORDER_GRAY)
add_textbox(slide3, det_l + det_w - Inches(1.65), mc_t + Inches(0.72), Inches(1.6), Inches(2.2),
            "■ 상세 프로필\n\n"
            "• 이름: 이차주\n"
            "• 상호: 대박물류\n"
            "• 번호:\n"
            " 123-45-67890\n"
            "• 계좌: 국민은행\n"
            " 123456-78-901\n"
            "• 폰: 010-9999-8888", 8, False, CHARCOAL)

# 반려 사유 및 액션
btn_y = mc_t + Inches(3.2)
add_textbox(slide3, det_l, btn_y + Inches(0.1), Inches(2.0), Inches(0.3), "회원 상태 변경:", 9.5, False, SLATE_GRAY)
# 상태 변경 드롭다운 묘사
add_rectangle(slide3, det_l + Inches(1.1), btn_y, Inches(1.5), Inches(0.4), WHITE, BORDER_GRAY)
add_textbox(slide3, det_l + Inches(1.2), btn_y + Inches(0.08), Inches(1.3), Inches(0.3), "정상 활동 중  ▼", 10, False, CHARCOAL)

app_btn = add_rectangle(slide3, det_l + Inches(2.7), btn_y, Inches(0.8), Inches(0.4), GREEN)
add_textbox(slide3, det_l + Inches(2.7), btn_y + Inches(0.08), Inches(0.8), Inches(0.3), "수정저장", 10, True, WHITE, PP_ALIGN.CENTER)

rej_btn = add_rectangle(slide3, det_l + Inches(3.6), btn_y, Inches(0.8), Inches(0.4), RED)
add_textbox(slide3, det_l + Inches(3.6), btn_y + Inches(0.08), Inches(0.8), Inches(0.3), "강제정지", 10, True, WHITE, PP_ALIGN.CENTER)


# =============================================================
# SLIDE 4: 상차지 및 하차지 대행 등록 / 상태 관리
# =============================================================
slide4 = prs.slides.add_slide(blank_layout)
add_rectangle(slide4, Inches(0), Inches(0), Inches(13.33), Inches(7.5), WHITE)

# 1. 마스터 테이블 레이아웃 추가
add_master_table_layout(slide4, "03. 현장 및 하차지 대행 등록 및 수명 상태 관리")

# 2. 우측 설명 기입
add_textbox(slide4, Inches(10.4), Inches(1.6), Inches(2.3), Inches(5.0),
            "■ 현장/하차지 제어 기능\n"
            "- 상차 공사현장 및 하차 사토장 정보를 대행으로 직접 등록합니다.\n"
            "- 강우나 천재지변 시 현장 상태를 [일시중단]으로 제어하여 배차를 즉각 잠금합니다.\n\n"
            "■ UI 핵심 명세\n"
            "1. 대행 등록 버튼:\n"
            "   - 상차지 등록, 하차지 등록 실행.\n"
            "2. 현장 리스트 테이블:\n"
            "   - 현재 운영 상태, 주소, 매핑 사토장 표시.\n"
            "3. 상태 제어 스위치:\n"
            "   - 운영중/일시중단 토글 스위치 묘사.",
            9, False, CHARCOAL)

# 3. 좌측 화면설계 영역 내 웹 콘솔 배치
mc_l, mc_t, mc_w, mc_h = add_web_frame(slide4, Inches(0.8), Inches(1.7), Inches(9.2), Inches(5.0), "현장 관리")

# 상단 대행 등록 단추 세트
reg_site_btn = add_rectangle(slide4, mc_l + Inches(0.2), mc_t + Inches(0.2), Inches(1.4), Inches(0.4), MUSTARD)
add_textbox(slide4, mc_l + Inches(0.2), mc_t + Inches(0.28), Inches(1.4), Inches(0.3), "➕ 신규 현장 등록", 9.5, True, CHARCOAL, PP_ALIGN.CENTER)

reg_land_btn = add_rectangle(slide4, mc_l + Inches(1.7), mc_t + Inches(0.2), Inches(1.4), Inches(0.4), CHARCOAL)
add_textbox(slide4, mc_l + Inches(1.7), mc_t + Inches(0.28), Inches(1.4), Inches(0.3), "➕ 신규 사토장 등록", 9.5, True, WHITE, PP_ALIGN.CENTER)

# 좌측 테이블 (등록 현장 목록)
t_w = Inches(4.3)
add_rectangle(slide4, mc_l + Inches(0.2), mc_t + Inches(0.7), t_w, mc_h - Inches(0.9), WHITE, BORDER_GRAY)
add_textbox(slide4, mc_l + Inches(0.35), mc_t + Inches(0.8), t_w - Inches(0.3), Inches(0.3), "등록된 현장 및 하차지 목록", 10, True, SLATE_GRAY)

add_textbox(slide4, mc_l + Inches(0.35), mc_t + Inches(1.15), t_w - Inches(0.3), Inches(3.0),
            "현장명/구분           주소            상태\n"
            "------------------------------------------\n"
            "[상차] 송도 1공구     인천 연수구      운영중 🟢\n"
            "[하차] 영종 사토장     인천 중구       운영중 🟢\n"
            "[상차] 남동공단 현장   인천 남동구      일시중단 🟡\n"
            "[상차] 아라뱃길 2공구  인천 계양구      종료 🔴", 8.5, False, CHARCOAL)

# 우측 상세 상태 수정 폼
f_l = mc_l + t_w + Inches(0.4)
f_w = mc_w - t_w - Inches(0.6)

add_textbox(slide4, f_l, mc_t + Inches(0.2), f_w, Inches(0.3), "상세 정보 수정 및 지오펜싱", 11, True, CHARCOAL)

# 입력 폼 그리기
form_rect = add_rectangle(slide4, f_l, mc_t + Inches(0.65), f_w, Inches(2.8), LIGHT_GRAY, BORDER_GRAY)
add_textbox(slide4, f_l + Inches(0.15), mc_t + Inches(0.75), f_w - Inches(0.3), Inches(2.6),
            "• 현장 명칭: 송도 1공구 아파트 신축공사\n"
            "• 지오펜싱 인식 반경:  200 m\n"
            "• 지도 핀 매핑: 인천 연수구 송도동 11-1\n"
            "• 수용 제한 자재: 양질토(수용) | 뻘흙(차단)\n"
            "• 하차지 용량: 6,500㎥ / 10,000㎥ 수용완료\n"
            "• 정산 결제 정책: 1회당 150,000원 (수금)", 9, False, CHARCOAL)

# 하단 상태 제어 스위치 및 버튼
ctrl_y = mc_t + Inches(3.6)
add_textbox(slide4, f_l, ctrl_y + Inches(0.1), Inches(1.8), Inches(0.3), "현장 실시간 상태:", 9.5, False, SLATE_GRAY)
# 상태 변경 버튼
add_rectangle(slide4, f_l + Inches(1.3), ctrl_y, Inches(1.2), Inches(0.4), MUSTARD)
add_textbox(slide4, f_l + Inches(1.3), ctrl_y + Inches(0.08), Inches(1.2), Inches(0.3), "운영중 🟢", 9.5, True, CHARCOAL, PP_ALIGN.CENTER)

save_site_btn = add_rectangle(slide4, f_l + Inches(2.6), ctrl_y, Inches(1.1), Inches(0.4), GREEN)
add_textbox(slide4, f_l + Inches(2.6), ctrl_y + Inches(0.08), Inches(1.1), Inches(0.3), "수정 저장", 10, True, WHITE, PP_ALIGN.CENTER)


# =============================================================
# SLIDE 5: 배차 및 하차 요청 등록 / 수명 제어
# =============================================================
slide5 = prs.slides.add_slide(blank_layout)
add_rectangle(slide5, Inches(0), Inches(0), Inches(13.33), Inches(7.5), WHITE)

# 1. 마스터 테이블 레이아웃 추가
add_master_table_layout(slide5, "04. 배차 및 하차 오더 대행 등록 및 강제 제어")

# 2. 우측 설명 기입
add_textbox(slide5, Inches(10.4), Inches(1.6), Inches(2.3), Inches(5.0),
            "■ 오더 수명 제어 기능\n"
            "- 전화 접수된 배차 요청 및 사토장 하차 요청을 어드민이 대신 등록합니다.\n"
            "- 진행 중인 오더를 강제로 취소(Cancelled)하거나 마감(Closed) 처리할 수 있습니다.\n\n"
            "■ UI 핵심 명세\n"
            "1. 배차 대행 등록 버튼:\n"
            "   - 상차지 오더 등록 폼 로딩.\n"
            "2. 진행 오더 테이블 목록:\n"
            "   - 현장명, 단가, 현재 매칭 건수, 진행율 게이지.\n"
            "3. 강제 상태 변경 버튼:\n"
            "   - 강제 마감 및 취소 액션.",
            9, False, CHARCOAL)

# 3. 좌측 화면설계 영역 내 웹 콘솔 배치
mc_l, mc_t, mc_w, mc_h = add_web_frame(slide5, Inches(0.8), Inches(1.7), Inches(9.2), Inches(5.0), "배차 제어")

# 상단 오더 수동 등록 버튼
add_rectangle(slide5, mc_l + Inches(0.2), mc_t + Inches(0.2), Inches(1.5), Inches(0.4), MUSTARD)
add_textbox(slide5, mc_l + Inches(0.2), mc_t + Inches(0.28), Inches(1.5), Inches(0.3), "➕ 배차 요청 등록", 9.5, True, CHARCOAL, PP_ALIGN.CENTER)

add_rectangle(slide5, mc_l + Inches(1.8), mc_t + Inches(0.2), Inches(1.5), Inches(0.4), CHARCOAL)
add_textbox(slide5, mc_l + Inches(1.8), mc_t + Inches(0.28), Inches(1.5), Inches(0.3), "➕ 하차 요청 등록", 9.5, True, WHITE, PP_ALIGN.CENTER)

# 좌측 진행 오더 리스트 테이블
t_w = Inches(4.3)
add_rectangle(slide5, mc_l + Inches(0.2), mc_t + Inches(0.7), t_w, mc_h - Inches(0.9), WHITE, BORDER_GRAY)
add_textbox(slide5, mc_l + Inches(0.35), mc_t + Inches(0.8), t_w - Inches(0.3), Inches(0.3), "실시간 오더 및 계약 목록", 10, True, SLATE_GRAY)

add_textbox(slide5, mc_l + Inches(0.35), mc_t + Inches(1.15), t_w - Inches(0.3), Inches(3.0),
            "현장명/구분           단가        매칭진행     상태\n"
            "--------------------------------------------------\n"
            "[배차] 송도 1공구     150,000원   3대/10대     모집중 🟢\n"
            "[하차] 영종 사토장     150,000원   2대/5대      모집중 🟢\n"
            "[배차] 남동공단 현장   120,000원   5대/5대      완료 🔴\n"
            "[배차] 아라뱃길 현장   130,000원   0대/10대     취소 ⚪", 8.5, False, CHARCOAL)

# 우측 오더 상세 제어판
f_l = mc_l + t_w + Inches(0.4)
f_w = mc_w - t_w - Inches(0.6)

add_textbox(slide5, f_l, mc_t + Inches(0.2), f_w, Inches(0.3), "배차 오더 상세 제어", 11, True, CHARCOAL)

# 오더 세부 필드 폼
form_rect = add_rectangle(slide5, f_l, mc_t + Inches(0.65), f_w, Inches(2.8), LIGHT_GRAY, BORDER_GRAY)
add_textbox(slide5, f_l + Inches(0.15), mc_t + Inches(0.75), f_w - Inches(0.3), Inches(2.6),
            "• 오더 명칭: 송도 1공구 아파트 토사 반출\n"
            "• 매칭 상차지: 송도 1공구 현장 (소속 매핑)\n"
            "• 연동 하차지: 영종 남단 사토장 (단가수금)\n"
            "• 모집 덤프 톤수:  25 톤 (Good Soil 전용)\n"
            "• 목표 모집 대수:  10 대 (현재 3대 배차수락)\n"
            "• 운송 단가 요금:  150,000 원\n"
            "• 거래처 결제 주기: 월말 세금계산서 정산", 9, False, CHARCOAL)

# 하단 수명 제어 액션 세트
ctrl_y = mc_t + Inches(3.6)
add_textbox(slide5, f_l, ctrl_y + Inches(0.1), Inches(1.8), Inches(0.3), "오더 강제 제어:", 9.5, False, SLATE_GRAY)

close_btn = add_rectangle(slide5, f_l + Inches(1.2), ctrl_y, Inches(1.2), Inches(0.4), CHARCOAL)
add_textbox(slide5, f_l + Inches(1.2), ctrl_y + Inches(0.08), Inches(1.2), Inches(0.3), "강제 매칭마감", 9.5, True, WHITE, PP_ALIGN.CENTER)

cancel_btn = add_rectangle(slide5, f_l + Inches(2.5), ctrl_y, Inches(1.2), Inches(0.4), RED)
add_textbox(slide5, f_l + Inches(2.5), ctrl_y + Inches(0.08), Inches(1.2), Inches(0.3), "오더 강제취소", 9.5, True, WHITE, PP_ALIGN.CENTER)


# =============================================================
# SLIDE 6: 배차 강제 개입 및 실시간 콜 제어
# =============================================================
slide6 = prs.slides.add_slide(blank_layout)
add_rectangle(slide6, Inches(0), Inches(0), Inches(13.33), Inches(7.5), WHITE)

# 1. 마스터 테이블 레이아웃 추가
add_master_table_layout(slide6, "05. 배차 강제 개입 및 실시간 콜 제어")

# 2. 우측 설명 기입
add_textbox(slide6, Inches(10.4), Inches(1.6), Inches(2.3), Inches(5.0),
            "■ 배차 개입 주요 기능\n"
            "- 특정 공사현장의 매칭율이 극도로 저조하거나 긴급 지시가 있을 때 어드민 권한으로 개입\n"
            "- 인근 공차 상태의 기사 정보를 자동 조회하여 1:1로 직접 강제 강제 지정을 하거나 콜 취소 처리\n\n"
            "■ UI 핵심 명세\n"
            "1. 미매칭 현황판:\n"
            "   - 배차 지연 오더의 매칭률 및 대기시간 실시간 정렬 테이블.\n"
            "2. [배차 개입] 레이어 팝업 위젯:\n"
            "   - '강제 배차 지정 기사 검색' 필드\n"
            "   - 추천 리스트: '김기사 (25톤, 1.2km 전방)' 및 수동 강제 배차 지정 실행 버튼", 
            12, False, CHARCOAL)

# 3. 좌측 화면설계 영역 내 웹 콘솔 배치
mc_l, mc_t, mc_w, mc_h = add_web_frame(slide6, Inches(0.8), Inches(1.7), Inches(9.2), Inches(5.0), "배차 제어")

# 미매칭 리스트
add_textbox(slide6, mc_l + Inches(0.2), mc_t + Inches(0.2), mc_w - Inches(0.4), Inches(0.3), "매칭 지연 및 미배차 오더 현황", 11, True, CHARCOAL)
table_rect = add_rectangle(slide6, mc_l + Inches(0.2), mc_t + Inches(0.55), mc_w - Inches(0.4), Inches(1.6), WHITE, BORDER_GRAY)
add_textbox(slide6, mc_l + Inches(0.35), mc_t + Inches(0.65), mc_w - Inches(0.7), Inches(1.4),
            "현장 오더명            단가        매칭현황    대기시간    조작\n"
            "------------------------------------------------------------------\n"
            "송도 1공구 아파트      150,000원   2대/10대    45분        [개입하기]\n"
            "영종도 공항 3공구      120,000원   1대/5대     30분        [개입하기]\n"
            "인천 신항 매립지       130,000원   5대/5대     완료        [관제확인]", 9, False, CHARCOAL)

# 개입 레이어 모달 팝업
popup_rect = add_rectangle(slide6, mc_l + Inches(0.8), mc_t + Inches(1.8), mc_w - Inches(1.6), Inches(2.2), WHITE, CHARCOAL, 2)
# 팝업 헤더
add_rectangle(slide6, mc_l + Inches(0.8), mc_t + Inches(1.8), mc_w - Inches(1.6), Inches(0.35), LIGHT_GRAY, BORDER_GRAY)
add_textbox(slide6, mc_l + Inches(0.95), mc_t + Inches(1.85), mc_w - Inches(1.9), Inches(0.3), "🚨 배차 강제 개입 콘솔 - 송도 1공구", 10, True, CHARCOAL)

# 추천 기사
add_textbox(slide6, mc_l + Inches(1.0), mc_t + Inches(2.25), mc_w - Inches(2.0), Inches(0.3), "인근 공차 대기 기사 추천 목록", 9, True, SLATE_GRAY)
k_rect = add_rectangle(slide6, mc_l + Inches(1.0), mc_t + Inches(2.6), mc_w - Inches(2.0), Inches(0.65), LIGHT_GRAY, BORDER_GRAY)
add_textbox(slide6, mc_l + Inches(1.15), mc_t + Inches(2.65), mc_w - Inches(2.3), Inches(0.5),
            "김기사 (25톤 덤프)  |  대기 위치: 송도 2교 인근 (1.2km 전방)\n연락처: 010-8888-7777", 9, False, CHARCOAL)

# 강제 매칭 단추
force_btn = add_rectangle(slide6, mc_l + mc_w - Inches(2.8), mc_t + Inches(3.3), Inches(1.8), Inches(0.4), MUSTARD)
add_textbox(slide6, mc_l + mc_w - Inches(2.8), mc_t + Inches(3.38), Inches(1.8), Inches(0.3), "강제 배차 지정 실행", 10, True, CHARCOAL, PP_ALIGN.CENTER)


# =============================================================
# SLIDE 7: 분쟁 조정 센터 (GPS 궤적 분석 및 직권 정산)
# =============================================================
slide7 = prs.slides.add_slide(blank_layout)
add_rectangle(slide7, Inches(0), Inches(0), Inches(13.33), Inches(7.5), WHITE)

# 1. 마스터 테이블 레이아웃 추가
add_master_table_layout(slide7, "06. 분쟁 조정 센터 (GPS 중재 및 직권 판정)")

# 2. 우측 설명 기입
add_textbox(slide7, Inches(10.4), Inches(1.6), Inches(2.3), Inches(5.0),
            "■ 분쟁 조정 주요 기능\n"
            "- 토질 불합격 회차 시, 정산 분쟁을 최종 중재합니다.\n"
            "- 차량의 GPS 역사 데이터를 추적하여 실제 하차지 진입 궤적 여부를 증명합니다.\n\n"
            "■ UI 핵심 명세\n"
            "1. 분쟁 이력 요약:\n"
            "   - 하차지 기각 사유와 기사의 이의 소명 텍스트.\n"
            "2. GPS 분석 맵:\n"
            "   - 실제 정상 주행을 하였는지 나타내는 궤적 표시판.\n"
            "3. [직권 정산] 단추:\n"
            "   - 70% 정산 명령(SETTLE_ADJUSTED) 적용 기능.",
            9, False, CHARCOAL)

# 3. 좌측 화면설계 영역 내 웹 콘솔 배치
mc_l, mc_t, mc_w, mc_h = add_web_frame(slide7, Inches(0.8), Inches(1.7), Inches(9.2), Inches(5.0), "분쟁 조정")

# 분쟁 요약 내용
info_rect = add_rectangle(slide7, mc_l + Inches(0.2), mc_t + Inches(0.2), mc_w - Inches(0.4), Inches(0.95), WHITE, BORDER_GRAY)
add_textbox(slide7, mc_l + Inches(0.35), mc_t + Inches(0.25), mc_w - Inches(0.7), Inches(0.35), "분쟁 관리 티켓 #2043 (김덤프 기사 ➔ 영종 사토장)", 11, True, CHARCOAL)
add_textbox(slide7, mc_l + Inches(0.35), mc_t + Inches(0.55), mc_w - Inches(0.7), Inches(0.5), 
            "• 하차지 반려 사유: 뻘흙 유입 불가  |  • 기사 소명: 양질의 사토로 적합 주행 완료함\n• 소명 자료: 하차지 진입 인증 완료 (정상 주행 거리 18.2km)", 9, False, SLATE_GRAY)

# GPS 궤적 지도 영역
map_rect = add_rectangle(slide7, mc_l + Inches(0.2), mc_t + Inches(1.25), mc_w - Inches(2.4), Inches(2.2), LIGHT_GRAY, BORDER_GRAY)
add_textbox(slide7, mc_l + Inches(0.2), mc_t + Inches(2.1), mc_w - Inches(2.4), Inches(0.5), "[실시간 수집 GPS 주행 궤적 지도 영역]\n상차지 ➔ 하차지 진입 ➔ 회차(정상 궤적 수집 완료)", 9, False, SLATE_GRAY, PP_ALIGN.CENTER)

# 우측 증빙 사진 썸네일
pic_rect = add_rectangle(slide7, mc_l + mc_w - Inches(2.0), mc_t + Inches(1.25), Inches(1.8), Inches(2.2), LIGHT_GRAY, BORDER_GRAY)
add_textbox(slide7, mc_l + mc_w - Inches(2.0), mc_t + Inches(2.1), Inches(1.8), Inches(0.5), "[현장 증빙\n뻘흙 사진 썸네일]", 9, False, SLATE_GRAY, PP_ALIGN.CENTER)

# 하단 최종 판정 제어부
dec_y = mc_t + Inches(3.6)
add_textbox(slide7, mc_l + Inches(0.2), dec_y + Inches(0.1), Inches(2.0), Inches(0.3), "직권 중재 판정 실행:", 10, True, SLATE_GRAY)

ok_b = add_rectangle(slide7, mc_l + Inches(1.8), dec_y, Inches(2.3), Inches(0.45), MUSTARD)
add_textbox(slide7, mc_l + Inches(1.8), dec_y + Inches(0.08), Inches(2.3), Inches(0.35), "70% 직권 정산 (Settle Adjusted)", 10, True, CHARCOAL, PP_ALIGN.CENTER)

rej_b = add_rectangle(slide7, mc_l + Inches(4.3), dec_y, Inches(2.3), Inches(0.45), CHARCOAL)
add_textbox(slide7, mc_l + Inches(4.3), dec_y + Inches(0.08), Inches(2.3), Inches(0.35), "무단 회차 판정 (미지급)", 10, True, WHITE, PP_ALIGN.CENTER)


# 파워포인트 파일 저장
output_path = "D:\\Projects\\dumpring\\dumpring-platform-backend\\덤프링_어드민_화면상세설계_V4.pptx"
prs.save(output_path)
print(f"PPTX 파일 생성 성공: {output_path}")
