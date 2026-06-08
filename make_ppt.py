import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# 컬러 팔레트 정의 (라이트 머스터드 옐로우 테마)
WHITE = RGBColor(255, 255, 255)
SAND_BG = RGBColor(249, 248, 246)
MUSTARD = RGBColor(245, 158, 11)     # F59E0B (주요 테마 색상)
AMBER = RGBColor(217, 119, 6)        # D97706 (포인트 색상)
CHARCOAL = RGBColor(31, 41, 55)      # 1F2937 (메인 텍스트)
LIGHT_GRAY = RGBColor(243, 244, 246)  # F3F4F6 (영역 구분 배경)
BORDER_GRAY = RGBColor(229, 231, 235) # E5E7EB (카드 테두리)
SLATE_GRAY = RGBColor(75, 85, 99)    # 4B5563 (서브 텍스트)
RED = RGBColor(239, 68, 68)          # EF4444 (반려/취소)

prs = Presentation()
# 16:9 와이드스크린 비율 설정
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6] # 빈 레이아웃

# 유틸리티 함수: 텍스트 상자 추가
def add_textbox(slide, left, top, width, height, text, font_size=12, bold=False, color=CHARCOAL, align=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = 'Malgun Gothic' # 맑은 고딕
    if align:
        p.alignment = align
    return txBox

# 유틸리티 함수: 직사각형 도형 추가
def add_rectangle(slide, left, top, width, height, fill_color, line_color=None, line_width=1):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

# 유틸리티 함수: 둥근 직사각형 도형 추가
def add_rounded_rect(slide, left, top, width, height, fill_color, line_color=None, line_width=1):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

# 유틸리티 함수: 모바일 폰 외곽 프레임 추가
def add_phone_frame(slide, left, top, width, height, title_text="덤프링"):
    # 외곽 디바이스 베젤
    bezel = add_rounded_rect(slide, left, top, width, height, SAND_BG, CHARCOAL, 3)
    # 스크린 영역
    screen_left = left + Inches(0.15)
    screen_top = top + Inches(0.4)
    screen_width = width - Inches(0.3)
    screen_height = height - Inches(0.6)
    screen = add_rectangle(slide, screen_left, screen_top, screen_width, screen_height, WHITE, BORDER_GRAY, 1)
    
    # 상단 상태바/헤더 영역
    header = add_rectangle(slide, screen_left, screen_top, screen_width, Inches(0.5), SAND_BG, BORDER_GRAY, 1)
    add_textbox(slide, screen_left + Inches(0.2), screen_top + Inches(0.12), screen_width - Inches(0.4), Inches(0.3), title_text, 14, True, CHARCOAL)
    
    return screen_left, screen_top + Inches(0.5), screen_width, screen_height - Inches(0.5)


# =============================================================
# SLIDE 1: 커버 페이지 (Cover Page)
# =============================================================
slide1 = prs.slides.add_slide(blank_layout)
add_rectangle(slide1, Inches(0), Inches(0), Inches(13.33), Inches(7.5), SAND_BG)
add_rectangle(slide1, Inches(1.5), Inches(2.2), Inches(0.15), Inches(2.5), MUSTARD)
add_textbox(slide1, Inches(2.0), Inches(2.1), Inches(9.5), Inches(1.2), "덤프링 (Dumpring)", 44, True, CHARCOAL)
add_textbox(slide1, Inches(2.0), Inches(3.2), Inches(9.5), Inches(0.8), "사용자 역할별 메뉴 구조 및 전체 화면설계서", 28, True, SLATE_GRAY)
add_textbox(slide1, Inches(2.0), Inches(4.1), Inches(9.5), Inches(0.5), "최종 테마: 라이트 머스터드 옐로우 (Light Mustard Yellow)", 16, False, AMBER)
add_textbox(slide1, Inches(2.0), Inches(5.8), Inches(6.0), Inches(0.8), "작성일: 2026. 06. 05\n작성자: 덤프링 플랫폼 기획팀", 12, False, SLATE_GRAY)


# =============================================================
# SLIDE 2: 로그인 및 회원가입 / 역할 선택 화면
# =============================================================
slide2 = prs.slides.add_slide(blank_layout)
add_rectangle(slide2, Inches(0), Inches(0), Inches(13.33), Inches(7.5), WHITE)

# 좌측 설명
add_textbox(slide2, Inches(0.8), Inches(0.8), Inches(4.5), Inches(0.8), "01. 로그인 / 회원가입 및 역할 선택", 24, True, CHARCOAL)
add_textbox(slide2, Inches(0.8), Inches(1.8), Inches(4.5), Inches(5.0), 
            "■ 회원 인증 및 가입 흐름\n"
            "- 휴대폰 번호 입력을 통한 1회용 인증번호(SMS) 전송\n"
            "- 인증 완료 시 최초 가입자는 역할 선택 화면으로 이동\n"
            "- 역할에 따라 앱 대시보드 및 측면 메뉴 슬라이더가 동적으로 분기\n\n"
            "■ UI/UX 설계 요점\n"
            "- 로그인 시 번호 입력 필드 및 인증 버튼의 직관성 확보\n"
            "- 역할 선택(Role Selection) 시 각 역할을 굵고 큼직한 카드로 시각화하여 사용자의 오탑입 방지\n"
            "- 브랜드 시그니처 옐로우 컬러로 메인 버튼 및 선택 카드 강조", 
            13, False, CHARCOAL)

# 우측 폰 1: 로그인 화면
p1_l, p1_t, p1_w, p1_h = add_phone_frame(slide2, Inches(5.8), Inches(0.5), Inches(3.2), Inches(6.5), "덤프링 - 로그인")
add_textbox(slide2, p1_l + Inches(0.2), p1_t + Inches(0.5), p1_w - Inches(0.4), Inches(0.4), "반갑습니다!", 18, True, CHARCOAL)
add_textbox(slide2, p1_l + Inches(0.2), p1_t + Inches(0.9), p1_w - Inches(0.4), Inches(0.4), "휴대폰 번호로 시작해 보세요.", 11, False, SLATE_GRAY)

# 입력 필드
add_rounded_rect(slide2, p1_l + Inches(0.2), p1_t + Inches(1.8), p1_w - Inches(0.4), Inches(0.5), WHITE, BORDER_GRAY)
add_textbox(slide2, p1_l + Inches(0.35), p1_t + Inches(1.9), p1_w - Inches(0.7), Inches(0.3), "휴대폰 번호 입력 (- 제외)", 11, False, SLATE_GRAY)

# 인증번호 받기 버튼
btn_get_code = add_rounded_rect(slide2, p1_l + Inches(0.2), p1_t + Inches(2.5), p1_w - Inches(0.4), Inches(0.5), MUSTARD)
add_textbox(slide2, p1_l + Inches(0.2), p1_t + Inches(2.6), p1_w - Inches(0.4), Inches(0.3), "인증번호 받기", 12, True, CHARCOAL, PP_ALIGN.CENTER)


# 우측 폰 2: 역할 선택 화면
p2_l, p2_t, p2_w, p2_h = add_phone_frame(slide2, Inches(9.5), Inches(0.5), Inches(3.2), Inches(6.5), "덤프링 - 역할 선택")
add_textbox(slide2, p2_l + Inches(0.2), p2_t + Inches(0.3), p2_w - Inches(0.4), Inches(0.4), "어떤 역할로 가입하시겠습니까?", 13, True, CHARCOAL)

# 역할 리스트 5종
roles = ["🚜 덤프 기사", "🏢 덤프 차주", "🏗️ 현장 관리자 (소장)", "👷 현장 담당자 (직원)", "🏕️ 하차지 이용자 (지주)"]
for i, role in enumerate(roles):
    ry = p2_t + Inches(0.8) + Inches(i * 0.9)
    # 기사 카드만 강조 표시
    f_color = WHITE
    l_color = BORDER_GRAY
    if i == 0:
        f_color = RGBColor(254, 243, 199) # 살짝 옐로우 빛 면
        l_color = MUSTARD
    card = add_rounded_rect(slide2, p2_l + Inches(0.15), ry, p2_w - Inches(0.3), Inches(0.75), f_color, l_color)
    add_textbox(slide2, p2_l + Inches(0.3), ry + Inches(0.2), p2_w - Inches(0.6), Inches(0.4), role, 12, True, CHARCOAL)


# =============================================================
# SLIDE 3: 전체 역할별 메뉴 구조 (Menu Architecture Map)
# =============================================================
slide3 = prs.slides.add_slide(blank_layout)
add_rectangle(slide3, Inches(0), Inches(0), Inches(13.33), Inches(7.5), SAND_BG)

add_textbox(slide3, Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.6), "덤프링 전체 사용자 메뉴 및 권한 구조도", 24, True, CHARCOAL)

# 5개 열 생성 (기사, 차주, 현장관리자, 담당자, 하차지 이용자)
col_w = Inches(2.2)
col_y = Inches(1.3)
col_h = Inches(5.3)

role_menus = [
    {
        "title": "🚜 덤프 기사",
        "desc": "일감 수락 & 실시간 운행",
        "menus": [
            "① 배차 콜 탐색 (지도)",
            "② 실시간 GPS 미터기 주행",
            "③ 운행 전표/영수증 조회",
            "④ 기사 필수서류 제출\n  (면허증/자격증/통장)",
            "⑤ 내 정보 및 정산 계좌 관리"
        ]
    },
    {
        "title": "🏢 덤프 차주",
        "desc": "보유 차량 및 기사단 관리",
        "menus": [
            "① 소속 운전기사 관리\n  (번호 연동 및 초대)",
            "② 차량 목록 및 배정 등록",
            "③ 차주 전용 운행 실적 집계",
            "④ 세금계산서/정산 서류",
            "⑤ 차주 필수서류 제출\n  (사업자등록증/통장)"
        ]
    },
    {
        "title": "🏗️ 현장 관리자 (소장)",
        "desc": "현장 개설 및 오더 발행",
        "menus": [
            "① 신규 공사현장 개설/등록",
            "② 소속 현장 직원 승인/관리",
            "③ B2B 매칭 오더 등록\n  (톤수/단가/자재 지정)",
            "④ 현장 실시간 대시보드 관제",
            "⑤ 현장별 정산/청구 실적"
        ]
    },
    {
        "title": "👷 현장 담당자 (직원)",
        "desc": "현장 매핑 신청 및 차량 관제",
        "menus": [
            "① 소속 현장 매핑 신청",
            "② 진입 차량 게이트 출입 승인",
            "③ 실시간 운행/상차 모니터링",
            "④ 일일 상차 작업 일지 조회",
            "⑤ 서류 간편 조회 및 확인"
        ]
    },
    {
        "title": "🏕️ 하차지 이용자 (지주)",
        "desc": "사토장 허가 및 게이트 반입",
        "menus": [
            "① 하차지 사토장 개설 등록",
            "② 매립 수용 공고 발행\n  (수용 토종/잔여량 명시)",
            "③ B2B 매칭 오더 수신/승인",
            "④ 실시간 반입 게이트 승인/반려\n  (뻘흙 사진 및 톤수 검수)",
            "⑤ 하차 수수료 정산 및 지갑"
        ]
    }
]

for idx, r_data in enumerate(role_menus):
    cx = Inches(0.8) + Inches(idx * 2.4)
    # 열 배경 카드
    add_rounded_rect(slide3, cx, col_y, col_w, col_h, WHITE, BORDER_GRAY)
    
    # 열 헤더 (머스터드 옐로우 헤더)
    add_rounded_rect(slide3, cx, col_y, col_w, Inches(0.85), MUSTARD)
    add_textbox(slide3, cx, col_y + Inches(0.1), col_w, Inches(0.3), r_data["title"], 13, True, CHARCOAL, PP_ALIGN.CENTER)
    add_textbox(slide3, cx, col_y + Inches(0.45), col_w, Inches(0.35), r_data["desc"], 9, False, CHARCOAL, PP_ALIGN.CENTER)
    
    # 메뉴 리스트 텍스트 작성
    menu_txt = "\n\n".join(r_data["menus"])
    add_textbox(slide3, cx + Inches(0.12), col_y + Inches(1.0), col_w - Inches(0.24), col_h - Inches(1.1), menu_txt, 11, False, CHARCOAL)


# =============================================================
# SLIDE 4: 덤프 기사 (Driver) 화면 및 세부 메뉴
# =============================================================
slide4 = prs.slides.add_slide(blank_layout)
add_rectangle(slide4, Inches(0), Inches(0), Inches(13.33), Inches(7.5), WHITE)

# 좌측 설명 영역
add_textbox(slide4, Inches(0.8), Inches(0.8), Inches(5.5), Inches(0.8), "02. 덤프 기사 (Driver) 화면 설계", 24, True, CHARCOAL)
add_textbox(slide4, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0), 
            "■ 기사 권한 및 핵심 메뉴 구조\n"
            "1. 배차 콜 탐색 (홈)\n"
            "   - 상차지 주소, 하차지 주소, 운송 거리 및 단가 확인\n"
            "2. 실시간 GPS 미터기 주행\n"
            "   - Tmap/카카오내비 호출 및 백그라운드 운임 계산 연동\n"
            "3. 운행 이력 / 운송 전표 조회\n"
            "   - 일일 총 운반 회수 및 누적 정산금액 조회\n"
            "4. 서류 제출 탭\n"
            "   - 기사 필수 3대 서류(면허증/자격증/통장) 촬영 및 제출\n\n"
            "■ UI/UX 설계 요점\n"
            "- 주행 중에는 모노스페이스 골드 요금 텍스트와 대형 [하차지 도착 완료] 버튼 중심으로 레이아웃 단순화", 
            13, False, CHARCOAL)

# 우측 폰 프레임 생성
sc_l, sc_t, sc_w, sc_h = add_phone_frame(slide4, Inches(7.5), Inches(0.5), Inches(4.5), Inches(6.5), "덤프링 - 기사 실시간 주행")

# 폰 내부 UI 컴포넌트들 (지도 영역)
map_rect = add_rectangle(slide4, sc_l, sc_t, sc_w, Inches(2.6), LIGHT_GRAY, BORDER_GRAY)
add_textbox(slide4, sc_l + Inches(0.2), sc_t + Inches(0.2), sc_w - Inches(0.4), Inches(0.4), "지도/네비게이션 궤적 활성화 영역", 11, False, SLATE_GRAY, PP_ALIGN.CENTER)

# 미터기 요금 판넬
meter_rect = add_rounded_rect(slide4, sc_l + Inches(0.15), sc_t + Inches(2.8), sc_w - Inches(0.3), Inches(2.0), WHITE, BORDER_GRAY)
add_textbox(slide4, sc_l + Inches(0.3), sc_t + Inches(2.95), sc_w - Inches(0.6), Inches(0.3), "실시간 예상 운임", 12, False, SLATE_GRAY, PP_ALIGN.CENTER)
add_textbox(slide4, sc_l + Inches(0.3), sc_t + Inches(3.25), sc_w - Inches(0.6), Inches(0.6), "95,000 원", 28, True, AMBER, PP_ALIGN.CENTER)

# 계기판 그리드
add_textbox(slide4, sc_l + Inches(0.3), sc_t + Inches(4.1), sc_w - Inches(0.6), Inches(0.5), 
            "속도: 60 km/h   |   거리: 12.4 km   |   시간: 00:24:15", 10, False, CHARCOAL, PP_ALIGN.CENTER)

# 도착 완료 버튼
arrive_btn = add_rounded_rect(slide4, sc_l + Inches(0.2), sc_t + Inches(4.7), sc_w - Inches(0.4), Inches(0.6), MUSTARD)
add_textbox(slide4, sc_l + Inches(0.2), sc_t + Inches(4.85), sc_w - Inches(0.4), Inches(0.4), "하차지 도착 완료", 14, True, CHARCOAL, PP_ALIGN.CENTER)


# =============================================================
# SLIDE 5: 덤프 차주 (Owner) 화면 및 세부 메뉴
# =============================================================
slide5 = prs.slides.add_slide(blank_layout)
add_rectangle(slide5, Inches(0), Inches(0), Inches(13.33), Inches(7.5), WHITE)

# 좌측 설명 영역
add_textbox(slide5, Inches(0.8), Inches(0.8), Inches(5.5), Inches(0.8), "03. 덤프 차주 (Owner) 화면 설계", 24, True, CHARCOAL)
add_textbox(slide5, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0), 
            "■ 차주 권한 및 핵심 메뉴 구조\n"
            "1. 소속 운전기사 연동 및 초대\n"
            "   - 본인 소속 기사의 휴대폰 번호를 선등록하여 배정\n"
            "2. 보유 차량 관리\n"
            "   - 15톤, 25톤, 27톤 등 보유 덤프트럭 번호판 및 규격 목록 관리\n"
            "3. 차주 정산 대시보드\n"
            "   - 소속 기사들의 일일 총 운송 실적 합산 모니터링\n"
            "4. 세금계산서 청구 내역\n"
            "5. 차주 서류 등록 탭\n"
            "   - 차주 필수 서류(사업자등록증, 사업자 통장) 업로드\n\n"
            "■ UI/UX 설계 요점\n"
            "- 기사 및 차량을 동시에 조회하고 신규 등록(초대)을 손쉽게 수행할 수 있도록 상단 탭 뷰 설계", 
            13, False, CHARCOAL)

# 우측 폰 프레임 생성
sc_l, sc_t, sc_w, sc_h = add_phone_frame(slide5, Inches(7.5), Inches(0.5), Inches(4.5), Inches(6.5), "덤프링 - 차주 소속차량/기사 관리")

# 차주 대시보드 상단 탭
tab_rect = add_rectangle(slide5, sc_l, sc_t, sc_w, Inches(0.5), LIGHT_GRAY, BORDER_GRAY)
add_textbox(slide5, sc_l + Inches(0.15), sc_t + Inches(0.12), sc_w/2 - Inches(0.2), Inches(0.3), "기사 관리", 12, True, MUSTARD, PP_ALIGN.CENTER)
add_textbox(slide5, sc_l + sc_w/2 + Inches(0.05), sc_t + Inches(0.12), sc_w/2 - Inches(0.2), Inches(0.3), "차량 관리", 12, False, SLATE_GRAY, PP_ALIGN.CENTER)

# 등록 기사 수 카운터 카드
c_card = add_rounded_rect(slide5, sc_l + Inches(0.15), sc_t + Inches(0.7), sc_w - Inches(0.3), Inches(0.9), WHITE, BORDER_GRAY)
add_textbox(slide5, sc_l + Inches(0.35), sc_t + Inches(0.85), sc_w - Inches(0.7), Inches(0.3), "소속 운전기사: 총 3명 등록됨", 12, True, CHARCOAL)
add_textbox(slide5, sc_l + Inches(0.35), sc_t + Inches(1.15), sc_w - Inches(0.7), Inches(0.35), "차량 미배정 기사: 1명", 10, False, AMBER)

# 기사 리스트 목록
g_names = ["김덤프 (25톤, 12가3456)", "이배차 (15톤, 78나9012)", "박기사 (대기중, 차량 미지정)"]
for idx, name in enumerate(g_names):
    gy = sc_t + Inches(1.8) + Inches(idx * 0.75)
    g_card = add_rounded_rect(slide5, sc_l + Inches(0.15), gy, sc_w - Inches(0.3), Inches(0.65), WHITE, BORDER_GRAY)
    add_textbox(slide5, sc_l + Inches(0.35), gy + Inches(0.15), sc_w - Inches(0.7), Inches(0.3), name, 11, False, CHARCOAL)

# 신규 기사 초대 버튼
invite_btn = add_rounded_rect(slide5, sc_l + Inches(0.2), sc_t + Inches(4.7), sc_w - Inches(0.4), Inches(0.6), MUSTARD)
add_textbox(slide5, sc_l + Inches(0.2), sc_t + Inches(4.85), sc_w - Inches(0.4), Inches(0.4), "신규 기사 초대 (+)", 13, True, CHARCOAL, PP_ALIGN.CENTER)


# =============================================================
# SLIDE 6: 현장 관리자 (Site Manager) 화면 및 세부 메뉴
# =============================================================
slide6 = prs.slides.add_slide(blank_layout)
add_rectangle(slide6, Inches(0), Inches(0), Inches(13.33), Inches(7.5), WHITE)

# 좌측 설명 영역
add_textbox(slide6, Inches(0.8), Inches(0.8), Inches(4.0), Inches(0.8), "04. 현장 관리자 (소장) 화면 설계", 24, True, CHARCOAL)
add_textbox(slide6, Inches(0.8), Inches(1.8), Inches(4.0), Inches(5.0), 
            "■ 현장 관리자 권한 및 메뉴 구조\n"
            "1. 공사현장 개설 및 서류 등록\n"
            "   - 현장 위치 지오펜싱 좌표 입력 및 토사 관리\n"
            "2. 소속 현장 직원 승인/관리\n"
            "   - 가입한 일반 담당자를 수동 승인하여 권한 위임\n"
            "3. B2B 매칭 오더 등록\n"
            "   - 하차지와 기사 매칭을 위한 오더 신규 발행\n"
            "4. 현장 실시간 운송 대시보드 (Web)\n"
            "   - 기사 실시간 관제 및 누적 정산 현황 분석\n\n"
            "■ UI/UX 설계 요점\n"
            "- 웹 브라우저 대시보드 화면은 통계 지표와 시간대별 통계 그래프를 조합해 현대적인 SaaS 레이아웃으로 설계", 
            13, False, CHARCOAL)

# 우측 웹 브라우저 프레임 생성
web_l = Inches(5.2)
web_t = Inches(1.0)
web_w = Inches(7.3)
web_h = Inches(5.2)

# 브라우저 백그라운드 & 헤더
add_rounded_rect(slide6, web_l, web_t, web_w, web_h, SAND_BG, BORDER_GRAY, 2)
add_rectangle(slide6, web_l, web_t, web_w, Inches(0.4), LIGHT_GRAY, BORDER_GRAY)
add_textbox(slide6, web_l + Inches(0.2), web_t + Inches(0.08), web_w - Inches(0.4), Inches(0.3), "덤프링 관리자 콘솔 - 현장 관제", 11, True, SLATE_GRAY)

# 통계 카드 3종
card_w = Inches(2.1)
card_h = Inches(0.9)
card_y = web_t + Inches(0.6)

# 카드 1 (트럭 대수)
add_rounded_rect(slide6, web_l + Inches(0.3), card_y, card_w, card_h, WHITE, BORDER_GRAY)
add_textbox(slide6, web_l + Inches(0.45), card_y + Inches(0.1), card_w, Inches(0.3), "운행 트럭 수", 10, False, SLATE_GRAY)
add_textbox(slide6, web_l + Inches(0.45), card_y + Inches(0.35), card_w, Inches(0.4), "12 대", 16, True, MUSTARD)

# 카드 2 (운송 횟수)
add_rounded_rect(slide6, web_l + Inches(2.6), card_y, card_w, card_h, WHITE, BORDER_GRAY)
add_textbox(slide6, web_l + Inches(2.75), card_y + Inches(0.1), card_w, Inches(0.3), "총 운송 횟수", 10, False, SLATE_GRAY)
add_textbox(slide6, web_l + Inches(2.75), card_y + Inches(0.35), card_w, Inches(0.4), "48 회", 16, True, AMBER)

# 카드 3 (누적 정산액)
add_rounded_rect(slide6, web_l + Inches(4.9), card_y, card_w, card_h, WHITE, BORDER_GRAY)
add_textbox(slide6, web_l + Inches(5.05), card_y + Inches(0.1), card_w, Inches(0.3), "누적 발생 운임", 10, False, SLATE_GRAY)
add_textbox(slide6, web_l + Inches(5.05), card_y + Inches(0.35), card_w, Inches(0.4), "4,800,000 원", 16, True, CHARCOAL)

# 가상 차트 박스
chart_rect = add_rounded_rect(slide6, web_l + Inches(0.3), web_t + Inches(1.7), web_w - Inches(0.6), Inches(1.6), WHITE, BORDER_GRAY)
add_textbox(slide6, web_l + Inches(0.5), web_t + Inches(1.8), web_w - Inches(1.0), Inches(0.3), "시간대별 배차 물량 통계 (차트 영역)", 11, True, CHARCOAL)
add_rectangle(slide6, web_l + Inches(1.5), web_t + Inches(2.5), Inches(4.5), Inches(0.08), MUSTARD)

# 하단 데이터 리스트 테이블
table_rect = add_rounded_rect(slide6, web_l + Inches(0.3), web_t + Inches(3.5), web_w - Inches(0.6), Inches(1.4), WHITE, BORDER_GRAY)
add_textbox(slide6, web_l + Inches(0.4), web_t + Inches(3.6), web_w - Inches(0.8), Inches(1.2), 
            "차량번호      규격      기사명      상태\n"
            "12가3456      25톤      김덤프      하차 완료\n"
            "78나9012      25톤      이배차      상차 대기\n"
            "34다5678      15톤      박사토      운행 중", 10, False, CHARCOAL)


# =============================================================
# SLIDE 7: 현장 담당자 (Site Worker) 화면 및 세부 메뉴
# =============================================================
slide7 = prs.slides.add_slide(blank_layout)
add_rectangle(slide7, Inches(0), Inches(0), Inches(13.33), Inches(7.5), WHITE)

# 좌측 설명 영역
add_textbox(slide7, Inches(0.8), Inches(0.8), Inches(5.5), Inches(0.8), "05. 현장 담당자 (직원) 화면 설계", 24, True, CHARCOAL)
add_textbox(slide7, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0), 
            "■ 담당자 권한 및 핵심 메뉴 구조\n"
            "1. 소속 현장 매핑 신청\n"
            "   - 관리자(소장)가 승인해 준 현장과의 사후 연결 탭\n"
            "2. 게이트 출입 통제 및 관제\n"
            "   - 현장 상하차를 위해 진입하는 덤프트럭 체크인/아웃\n"
            "3. 일일 운행 및 상차 현황 실시간 조회\n"
            "4. 모바일 상차일지 기록\n\n"
            "■ UI/UX 설계 요점\n"
            "- 현장 한가운데서 작업 차량을 즉각적으로 출입 통제해야 하므로, 버튼식 간편 출입 기록 화면을 기획", 
            13, False, CHARCOAL)

# 우측 폰 프레임 생성
sc_l, sc_t, sc_w, sc_h = add_phone_frame(slide7, Inches(7.5), Inches(0.5), Inches(4.5), Inches(6.5), "덤프링 - 현장 담당자 게이트 통제")

# 현장 매핑 정보 카드
add_textbox(slide7, sc_l + Inches(0.2), sc_t + Inches(0.3), sc_w - Inches(0.4), Inches(0.4), "현재 근무 현장", 13, True, CHARCOAL)
site_info = add_rounded_rect(slide7, sc_l + Inches(0.15), sc_t + Inches(0.7), sc_w - Inches(0.3), Inches(0.8), WHITE, BORDER_GRAY)
add_textbox(slide7, sc_l + Inches(0.3), sc_t + Inches(0.85), sc_w - Inches(0.6), Inches(0.5), 
            "송도 현대 아파트 건설 현장 (소장 승인 완료)\n담당 직원: 홍길동 대리", 11, False, SLATE_GRAY)

# 출입 차량 현황 및 액션
add_textbox(slide7, sc_l + Inches(0.2), sc_t + Inches(1.8), sc_w - Inches(0.4), Inches(0.4), "상차 게이트 현황", 13, True, CHARCOAL)

# 차량 대기 리스트 박스
list_rect = add_rounded_rect(slide7, sc_l + Inches(0.15), sc_t + Inches(2.2), sc_w - Inches(0.3), Inches(1.8), WHITE, BORDER_GRAY)
add_textbox(slide7, sc_l + Inches(0.35), sc_t + Inches(2.35), sc_w - Inches(0.7), Inches(1.5), 
            "• 대기 차량 (상차 중)\n"
            "  - 12가3456 (25톤) - 흙 상차 중\n"
            "  - 78나9012 (25톤) - 흙 상차 대기\n\n"
            "• 금일 총 상차 완료 차량 수: 14대", 11, False, CHARCOAL)

# 현장 출입 등록 버튼
register_in_btn = add_rounded_rect(slide7, sc_l + Inches(0.2), sc_t + Inches(4.2), sc_w - Inches(0.4), Inches(0.5), MUSTARD)
add_textbox(slide7, sc_l + Inches(0.2), sc_t + Inches(4.3), sc_w - Inches(0.4), Inches(0.3), "신규 차량 진입 등록 (+)", 12, True, CHARCOAL, PP_ALIGN.CENTER)

# 상차 완료 및 출발 처리 버튼
load_complete_btn = add_rounded_rect(slide7, sc_l + Inches(0.2), sc_t + Inches(4.8), sc_w - Inches(0.4), Inches(0.5), CHARCOAL)
add_textbox(slide7, sc_l + Inches(0.2), sc_t + Inches(4.9), sc_w - Inches(0.4), Inches(0.3), "상차 완료 및 출발 처리", 12, True, WHITE, PP_ALIGN.CENTER)


# =============================================================
# SLIDE 8: 하차지 이용자 (Landowner) 화면 및 세부 메뉴
# =============================================================
slide8 = prs.slides.add_slide(blank_layout)
add_rectangle(slide8, Inches(0), Inches(0), Inches(13.33), Inches(7.5), WHITE)

# 좌측 설명 영역
add_textbox(slide8, Inches(0.8), Inches(0.8), Inches(5.5), Inches(0.8), "06. 하차지 이용자 (지주) 화면 설계", 24, True, CHARCOAL)
add_textbox(slide8, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0), 
            "■ 지주 권한 및 핵심 메뉴 구조\n"
            "1. 사토장 개설 등록\n"
            "   - 사토장 위치, 면적, 매립 허가 서류 및 계좌 등록\n"
            "2. 매립 수용 공고 발행 (일감 모집)\n"
            "   - 받아들일 토질(양질토 등), 단가, 총 수용량 설정 공고\n"
            "3. B2B 현장 매칭 오더 승인/반려\n"
            "   - 화주가 발행한 현장 계약 오더의 최종 사전 승인/반려\n"
            "4. 실시간 게이트 반입 승인/반려\n"
            "   - 게이트 도착 덤프트럭 자재 육안 검사 후 원터치 승인\n\n"
            "■ UI/UX 설계 요점\n"
            "- 기사가 도착하면 즉각적으로 반입 승인을 누를 수 있는 직관적인 승인 화면 제공", 
            13, False, CHARCOAL)

# 우측 폰 프레임 생성
sc_l, sc_t, sc_w, sc_h = add_phone_frame(slide8, Inches(7.5), Inches(0.5), Inches(4.5), Inches(6.5), "덤프링 - 하차지 반입 승인")

# 폰 내부 UI 컴포넌트들
add_textbox(slide8, sc_l + Inches(0.2), sc_t + Inches(0.2), sc_w - Inches(0.4), Inches(0.4), "실시간 게이트 반입 심사 대기", 13, True, CHARCOAL)

# 차량 카드
car_card = add_rounded_rect(slide8, sc_l + Inches(0.15), sc_t + Inches(0.6), sc_w - Inches(0.3), Inches(0.8), WHITE, BORDER_GRAY)
add_textbox(slide8, sc_l + Inches(0.3), sc_t + Inches(0.8), sc_w - Inches(0.6), Inches(0.5), "12가3456 (25톤) - 게이트 도착 완료", 12, True, CHARCOAL)

# 자재 검수용 썸네일 박스
photo_rect = add_rectangle(slide8, sc_l + Inches(0.15), sc_t + Inches(1.6), sc_w - Inches(0.3), Inches(2.2), LIGHT_GRAY, BORDER_GRAY)
add_textbox(slide8, sc_l + Inches(0.15), sc_t + Inches(2.5), sc_w - Inches(0.3), Inches(0.4), "[자재/토사 카메라 사진 분석 영역]", 11, False, SLATE_GRAY, PP_ALIGN.CENTER)

# 승인 버튼 (머스터드 옐로우)
approve_btn = add_rounded_rect(slide8, sc_l + Inches(0.2), sc_t + Inches(4.0), sc_w - Inches(0.4), Inches(0.5), MUSTARD)
add_textbox(slide8, sc_l + Inches(0.2), sc_t + Inches(4.1), sc_w - Inches(0.4), Inches(0.3), "반입 승인", 12, True, CHARCOAL, PP_ALIGN.CENTER)

# 반려 및 회차 버튼 (차콜 어두운 톤)
reject_btn = add_rounded_rect(slide8, sc_l + Inches(0.2), sc_t + Inches(4.6), sc_w - Inches(0.4), Inches(0.5), CHARCOAL)
add_textbox(slide8, sc_l + Inches(0.2), sc_t + Inches(4.7), sc_w - Inches(0.4), Inches(0.3), "반입 반려 및 회차", 12, True, WHITE, PP_ALIGN.CENTER)


# 파워포인트 파일 저장
output_path = "D:\\Projects\\dumpring\\dumpring-platform-backend\\덤프링_화면설계서_최종_V2.pptx"
prs.save(output_path)
print(f"PPTX 파일 생성 성공: {output_path}")
