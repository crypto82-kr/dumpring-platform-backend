import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# 컬러 팔레트 정의 (라이트 머스터드 옐로우 테마)
WHITE = RGBColor(255, 255, 255)
SAND_BG = RGBColor(249, 248, 246)
MUSTARD = RGBColor(245, 158, 11)     # F59E0B (주요 테마 색상)
AMBER = RGBColor(217, 119, 6)        # D97706 (포인트 색상)
CHARCOAL = RGBColor(31, 41, 55)      # 1F2937 (메인 텍스트)
LIGHT_GRAY = RGBColor(243, 244, 246)  # F3F4F6 (영역 구분 배경)
BORDER_GRAY = RGBColor(229, 231, 235) # E5E7EB (카드 테두리)
SLATE_GRAY = RGBColor(75, 85, 99)    # 4B5563 (서브 텍스트)
RED = RGBColor(239, 68, 68)          # EF4444 (반려/에러/취소)
GREEN = RGBColor(16, 185, 129)       # 10B981 (정산완료/정상)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6] # 빈 레이아웃

# 유틸리티 함수: 텍스트 상자 추가
def add_textbox(slide, left, top, width, height, text, font_size=11, bold=False, color=CHARCOAL, align=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = 'Malgun Gothic'
    if align:
        p.alignment = align
    return txBox

# 유틸리티 함수: 직사각형 도형 추가
def add_rectangle(slide, left, top, width, height, fill_color, line_color=None, line_width=1):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

# 유틸리티 함수: 둥근 직사각형 도형 추가
def add_rounded_rect(slide, left, top, width, height, fill_color, line_color=None, line_width=1):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

# 유틸리티 함수: 모바일 폰 프레임 생성
def add_phone_frame(slide, left, top, width, height, title_text="덤프링"):
    bezel = add_rounded_rect(slide, left, top, width, height, SAND_BG, CHARCOAL, 3)
    screen_left = left + Inches(0.15)
    screen_top = top + Inches(0.4)
    screen_width = width - Inches(0.3)
    screen_height = height - Inches(0.6)
    screen = add_rectangle(slide, screen_left, screen_top, screen_width, screen_height, WHITE, BORDER_GRAY, 1)
    
    # 상단 헤더
    header = add_rectangle(slide, screen_left, screen_top, screen_width, Inches(0.5), SAND_BG, BORDER_GRAY, 1)
    add_textbox(slide, screen_left + Inches(0.15), screen_top + Inches(0.12), screen_width - Inches(0.3), Inches(0.3), title_text, 13, True, CHARCOAL)
    
    return screen_left, screen_top + Inches(0.5), screen_width, screen_height - Inches(0.5)

# 유틸리티 함수: 문서 제출/사진 촬영 업로드용 가상 버튼 그리기
def add_upload_box(slide, left, top, width, height, label_text):
    add_rounded_rect(slide, left, top, width, height, LIGHT_GRAY, BORDER_GRAY, 1)
    add_rectangle(slide, left + (width/2) - Inches(0.15), top + (height/2) - Inches(0.2), Inches(0.3), Inches(0.05), SLATE_GRAY)
    add_rectangle(slide, left + (width/2) - Inches(0.025), top + (height/2) - Inches(0.3), Inches(0.05), Inches(0.25), SLATE_GRAY)
    add_textbox(slide, left, top + (height/2) + Inches(0.05), width, Inches(0.35), label_text, 8, False, SLATE_GRAY, PP_ALIGN.CENTER)


# =============================================================
# SLIDE 1: 커버 페이지 (Cover Page)
# =============================================================
slide1 = prs.slides.add_slide(blank_layout)
add_rectangle(slide1, Inches(0), Inches(0), Inches(13.33), Inches(7.5), SAND_BG)
add_rectangle(slide1, Inches(1.5), Inches(2.2), Inches(0.15), Inches(2.5), MUSTARD)
add_textbox(slide1, Inches(2.0), Inches(2.1), Inches(9.5), Inches(1.2), "덤프링 (Dumpring)", 44, True, CHARCOAL)
add_textbox(slide1, Inches(2.0), Inches(3.2), Inches(9.5), Inches(0.8), "04. 덤프 차주 (Owner) 업무 화면 상세설계서", 28, True, SLATE_GRAY)
add_textbox(slide1, Inches(2.0), Inches(4.1), Inches(9.5), Inches(0.5), "최종 테마: 라이트 머스터드 옐로우 (Light Mustard Yellow)", 16, False, AMBER)
add_textbox(slide1, Inches(2.0), Inches(5.8), Inches(6.0), Inches(0.8), "작성일: 2026. 06. 05\n작성자: 덤프링 플랫폼 기획팀", 12, False, SLATE_GRAY)


# =============================================================
# SLIDE 2: Screen 1. 차주 홈 화면 (Owner Home Dashboard)
# =============================================================
slide2 = prs.slides.add_slide(blank_layout)
add_rectangle(slide2, Inches(0), Inches(0), Inches(13.33), Inches(7.5), WHITE)

# 좌측 설명 영역
add_textbox(slide2, Inches(0.8), Inches(0.8), Inches(5.5), Inches(0.8), "01. 차주 홈 화면 (대시보드)", 24, True, CHARCOAL)
add_textbox(slide2, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0), 
            "■ 주요 기능\n"
            "- 차주가 고용하여 소속된 기사들의 당일 주행 관제 및 보유 덤프트럭 관리\n"
            "- 기사별/차량별 매핑 상태 및 실시간 운행 현황을 그리드로 모니터링\n\n"
            "■ 주요 UI 구성\n"
            "1. 차주 종합 지표:\n"
            "   - 소속 기사 5명 (4명 주행 중 / 1명 휴무)\n"
            "   - 보유 차량 4대 (3대 매핑 완료)\n"
            "   - 정산 대기 누적 요금: 1,800,000원\n"
            "2. 기사-차량 연동 매핑 리스트:\n"
            "   - 각 덤프트럭 번호판별로 매핑된 기사명 및 현재 주행/대기 단계 노출", 
            12, False, CHARCOAL)

# 우측 폰 프레임 생성
sc_l, sc_t, sc_w, sc_h = add_phone_frame(slide2, Inches(7.5), Inches(0.5), Inches(4.5), Inches(6.5), "덤프링 - 차주 대시보드")

# 차주 지표 판넬
summary_rect = add_rounded_rect(slide2, sc_l + Inches(0.15), sc_t + Inches(0.15), sc_w - Inches(0.3), Inches(1.6), WHITE, BORDER_GRAY)
add_textbox(slide2, sc_l + Inches(0.3), sc_t + Inches(0.25), sc_w - Inches(0.6), Inches(0.3), "금일 보유 차량 및 운송 현황", 12, True, CHARCOAL)
add_textbox(slide2, sc_l + Inches(0.3), sc_t + Inches(0.6), sc_w - Inches(0.6), Inches(0.7), 
            "• 소속 기사: 총 5명 (4명 운행 중  |  1명 대기)\n"
            "• 보유 차량: 총 4대 (3대 배정 완료  |  1대 대기중)", 10, False, SLATE_GRAY)
add_textbox(slide2, sc_l + Inches(0.3), sc_t + Inches(1.25), sc_w - Inches(0.6), Inches(0.4), "정산 대기 금액: 1,800,000 원", 13, True, AMBER)

# 매핑 관제 리스트
list_y = sc_t + Inches(1.9)
add_textbox(slide2, sc_l + Inches(0.2), list_y, sc_w - Inches(0.4), Inches(0.3), "실시간 차량/기사 운행 관제", 11, True, SLATE_GRAY)

# 리스트 아이템 1
i1_y = list_y + Inches(0.3)
add_rounded_rect(slide2, sc_l + Inches(0.15), i1_y, sc_w - Inches(0.3), Inches(0.7), WHITE, BORDER_GRAY)
add_textbox(slide2, sc_l + Inches(0.3), i1_y + Inches(0.1), sc_w - Inches(0.6), Inches(0.35), "경기 80사 1234 (25톤)", 11, True, CHARCOAL)
add_textbox(slide2, sc_l + Inches(0.3), i1_y + Inches(0.38), sc_w - Inches(0.6), Inches(0.3), "김덤프 기사  |  미터기 주행 중", 9, False, AMBER)

# 리스트 아이템 2
i2_y = list_y + Inches(1.1)
add_rounded_rect(slide2, sc_l + Inches(0.15), i2_y, sc_w - Inches(0.3), Inches(0.7), WHITE, BORDER_GRAY)
add_textbox(slide2, sc_l + Inches(0.3), i2_y + Inches(0.1), sc_w - Inches(0.6), Inches(0.35), "서울 80사 5678 (25톤)", 11, True, CHARCOAL)
add_textbox(slide2, sc_l + Inches(0.3), i2_y + Inches(0.38), sc_w - Inches(0.6), Inches(0.3), "이배차 기사  |  상차지 이동 중", 9, False, SLATE_GRAY)

# 리스트 아이템 3
i3_y = list_y + Inches(1.9)
add_rounded_rect(slide2, sc_l + Inches(0.15), i3_y, sc_w - Inches(0.3), Inches(0.7), WHITE, BORDER_GRAY)
add_textbox(slide2, sc_l + Inches(0.3), i3_y + Inches(0.1), sc_w - Inches(0.6), Inches(0.35), "인천 80사 9012 (15톤)", 11, True, CHARCOAL)
add_textbox(slide2, sc_l + Inches(0.3), i3_y + Inches(0.38), sc_w - Inches(0.6), Inches(0.3), "박기사 기사  |  운행 종료 (휴무)", 9, False, SLATE_GRAY)


# =============================================================
# SLIDE 3: Screen 2. 소속 운전기사 관리 및 초대 화면 (Driver Management)
# =============================================================
slide3 = prs.slides.add_slide(blank_layout)
add_rectangle(slide3, Inches(0), Inches(0), Inches(13.33), Inches(7.5), WHITE)

# 좌측 설명 영역
add_textbox(slide3, Inches(0.8), Inches(0.8), Inches(5.5), Inches(0.8), "02. 소속 운전기사 관리 및 초대", 24, True, CHARCOAL)
add_textbox(slide3, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0), 
            "■ 주요 기능\n"
            "- 차주 밑으로 기사를 신규 등록하기 위한 초대 링크 발송 기능\n"
            "- 기사 휴대폰 번호를 선등록하여 초대하면, 기사가 가입 완료 시 차주 소속으로 자동 배핑\n"
            "- 각 기사 카드에서 현재 배정된 트럭 번호판을 실시간으로 변경 지정하는 드롭다운 제공\n\n"
            "■ 주요 UI 구성\n"
            "1. 기사 초대 상자: 휴대폰 번호 기입 필드 및 [초대 링크 발송] 버튼\n"
            "2. 기사 리스트 카드:\n"
            "   - 기사 실명 및 면허 심사 승인 여부 표시\n"
            "   - 차량 배정 드롭다운 메뉴 (예: '경기 80사 1234 ▼')", 
            12, False, CHARCOAL)

# 우측 폰 프레임 생성
sc_l, sc_t, sc_w, sc_h = add_phone_frame(slide3, Inches(7.5), Inches(0.5), Inches(4.5), Inches(6.5), "소속 기사 관리")

# 기사 초대 상자
add_textbox(slide3, sc_l + Inches(0.2), sc_t + Inches(0.2), sc_w - Inches(0.4), Inches(0.3), "신규 운전기사 초대", 11, True, SLATE_GRAY)
invite_rect = add_rounded_rect(slide3, sc_l + Inches(0.15), sc_t + Inches(0.5), sc_w - Inches(0.3), Inches(1.2), WHITE, BORDER_GRAY)
# 번호 입력란
add_rounded_rect(slide3, sc_l + Inches(0.3), sc_t + Inches(0.65), sc_w - Inches(0.6), Inches(0.4), WHITE, BORDER_GRAY)
add_textbox(slide3, sc_l + Inches(0.45), sc_t + Inches(0.72), sc_w - Inches(0.9), Inches(0.3), "기사 휴대폰 번호 입력", 11, False, SLATE_GRAY)
# 초대 발송 버튼
inv_btn = add_rounded_rect(slide3, sc_l + Inches(0.3), sc_t + Inches(1.15), sc_w - Inches(0.6), Inches(0.4), MUSTARD)
add_textbox(slide3, sc_l + Inches(0.3), sc_t + Inches(1.22), sc_w - Inches(0.6), Inches(0.3), "초대 링크 발송 (SMS)", 11, True, CHARCOAL, PP_ALIGN.CENTER)

# 기사 리스트
list_y = sc_t + Inches(1.9)
add_textbox(slide3, sc_l + Inches(0.2), list_y, sc_w - Inches(0.4), Inches(0.3), "소속 기사 목록", 11, True, SLATE_GRAY)

# 기사 카드 1
k1_y = list_y + Inches(0.3)
add_rounded_rect(slide3, sc_l + Inches(0.15), k1_y, sc_w - Inches(0.3), Inches(1.0), WHITE, BORDER_GRAY)
add_textbox(slide3, sc_l + Inches(0.3), k1_y + Inches(0.1), sc_w - Inches(0.6), Inches(0.35), "김덤프 기사 (승인완료 🟢)", 12, True, CHARCOAL)
# 차량 배정 드롭다운 묘사
dd1 = add_rounded_rect(slide3, sc_l + Inches(0.3), k1_y + Inches(0.48), sc_w - Inches(0.6), Inches(0.4), LIGHT_GRAY, BORDER_GRAY)
add_textbox(slide3, sc_l + Inches(0.45), k1_y + Inches(0.53), sc_w - Inches(0.9), Inches(0.3), "차량 배정: 경기 80사 1234  ▼", 10, False, CHARCOAL)

# 기사 카드 2
k2_y = list_y + Inches(1.4)
add_rounded_rect(slide3, sc_l + Inches(0.15), k2_y, sc_w - Inches(0.3), Inches(1.0), WHITE, BORDER_GRAY)
add_textbox(slide3, sc_l + Inches(0.3), k2_y + Inches(0.1), sc_w - Inches(0.6), Inches(0.35), "이배차 기사 (승인완료 🟢)", 12, True, CHARCOAL)
# 차량 배정 드롭다운 묘사
dd2 = add_rounded_rect(slide3, sc_l + Inches(0.3), k2_y + Inches(0.48), sc_w - Inches(0.6), Inches(0.4), LIGHT_GRAY, BORDER_GRAY)
add_textbox(slide3, sc_l + Inches(0.45), k2_y + Inches(0.53), sc_w - Inches(0.9), Inches(0.3), "차량 배정: 미배정  ▼", 10, False, SLATE_GRAY)


# =============================================================
# SLIDE 4: Screen 3. 보유 차량 등록 및 서류 제출 화면 (Vehicle Registration)
# =============================================================
slide4 = prs.slides.add_slide(blank_layout)
add_rectangle(slide4, Inches(0), Inches(0), Inches(13.33), Inches(7.5), WHITE)

# 좌측 설명 영역
add_textbox(slide4, Inches(0.8), Inches(0.8), Inches(5.5), Inches(0.8), "03. 보유 차량 등록 및 서류 제출", 24, True, CHARCOAL)
add_textbox(slide4, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0), 
            "■ 주요 기능\n"
            "- 차주가 보유한 신규 덤프트럭 번호판 및 규격을 시스템에 등록하는 화면\n"
            "- 법적 운행 심사를 위한 필수 차량 서류(자동차등록증) 업로드 필수\n"
            "- 등록 신청 후 관리자가 번호판과 서류 진위 검사 완료 시 정식 배차 승인\n\n"
            "■ 주요 UI 구성\n"
            "1. 입력 폼: 차량 번호판 기입 필드, 덤프 톤수 토글 단추\n"
            "2. 증빙 첨부: 자동차등록증 사본 업로드 썸네일 박스\n"
            "3. [차량 등록 신청] 메인 옐로우 버튼", 
            12, False, CHARCOAL)

# 우측 폰 프레임 생성
sc_l, sc_t, sc_w, sc_h = add_phone_frame(slide4, Inches(7.5), Inches(0.5), Inches(4.5), Inches(6.5), "보유 차량 등록")

# 차량 번호 입력
add_textbox(slide4, sc_l + Inches(0.2), sc_t + Inches(0.2), sc_w - Inches(0.4), Inches(0.3), "차량 등록 번호판", 10, True, SLATE_GRAY)
add_rounded_rect(slide4, sc_l + Inches(0.2), sc_t + Inches(0.45), sc_w - Inches(0.4), Inches(0.45), WHITE, BORDER_GRAY)
add_textbox(slide4, sc_l + Inches(0.35), sc_t + Inches(0.53), sc_w - Inches(0.7), Inches(0.3), "서울 80사 7777", 12, True, CHARCOAL)

# 규격 선택 (25톤 활성 토글)
add_textbox(slide4, sc_l + Inches(0.2), sc_t + Inches(1.0), sc_w - Inches(0.4), Inches(0.3), "덤프 규격 (톤수)", 10, True, SLATE_GRAY)
add_rounded_rect(slide4, sc_l + Inches(0.2), sc_t + Inches(1.3), Inches(1.2), Inches(0.4), WHITE, BORDER_GRAY)
add_textbox(slide4, sc_l + Inches(0.2), sc_t + Inches(1.38), Inches(1.2), Inches(0.3), "15톤 덤프", 10, False, SLATE_GRAY, PP_ALIGN.CENTER)

add_rounded_rect(slide4, sc_l + Inches(1.5), sc_t + Inches(1.3), Inches(1.2), Inches(0.4), MUSTARD)
add_textbox(slide4, sc_l + Inches(1.5), sc_t + Inches(1.38), Inches(1.2), Inches(0.3), "25톤 덤프", 10, True, CHARCOAL, PP_ALIGN.CENTER)

# 자동차등록증 업로드 박스
add_textbox(slide4, sc_l + Inches(0.2), sc_t + Inches(1.9), sc_w - Inches(0.4), Inches(0.3), "필수 차량 서류 보관함", 10, True, SLATE_GRAY)
add_upload_box(slide4, sc_l + Inches(0.15), sc_t + Inches(2.2), sc_w - Inches(0.3), Inches(1.5), "자동차등록증 사본 첨부")

# 등록 신청 버튼
reg_btn = add_rounded_rect(slide4, sc_l + Inches(0.2), sc_t + Inches(4.7), sc_w - Inches(0.4), Inches(0.6), MUSTARD)
add_textbox(slide4, sc_l + Inches(0.2), sc_t + Inches(4.85), sc_w - Inches(0.4), Inches(0.4), "신규 차량 등록 신청", 14, True, CHARCOAL, PP_ALIGN.CENTER)


# =============================================================
# SLIDE 5: Screen 4. 차주 전용 운행 실적 및 정산 화면 (Owner Settlement)
# =============================================================
slide5 = prs.slides.add_slide(blank_layout)
add_rectangle(slide5, Inches(0), Inches(0), Inches(13.33), Inches(7.5), WHITE)

# 좌측 설명 영역
add_textbox(slide5, Inches(0.8), Inches(0.8), Inches(5.5), Inches(0.8), "04. 정산 및 기사별 매출 기여도 조회", 24, True, CHARCOAL)
add_textbox(slide5, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0), 
            "■ 주요 기능\n"
            "- 차주 소속 기사들이 운행을 통해 발생시킨 전체 매출 실적을 월별 합산 모니터링\n"
            "- 수수료 공제 시뮬레이션 및 최종 정산 예정액 산식 노출\n"
            "- 기사별 기여도 순위를 목록으로 보여주어 차주의 관리 효율 증대\n\n"
            "■ 주요 UI 구성\n"
            "1. 이번 달 종합 정산 카드:\n"
            "   - 총 발생 운임 4,920,000원, 플랫폼 수수료 -492,000원, 최종 차주 실수령 4,428,000원\n"
            "2. 기사별 매출 기여도 목록:\n"
            "   - 기사 실명, 담당 차량 톤수, 당월 총 운임 및 운송 횟수 집계 리스트", 
            12, False, CHARCOAL)

# 우측 폰 프레임 생성
sc_l, sc_t, sc_w, sc_h = add_phone_frame(slide5, Inches(7.5), Inches(0.5), Inches(4.5), Inches(6.5), "정산 및 운송 실적")

# 정산금 카드
fee_card = add_rounded_rect(slide5, sc_l + Inches(0.15), sc_t + Inches(0.15), sc_w - Inches(0.3), Inches(1.7), WHITE, BORDER_GRAY)
add_textbox(slide5, sc_l + Inches(0.3), sc_t + Inches(0.25), sc_w - Inches(0.6), Inches(0.3), "6월 청구/정산 예정 내역", 12, True, CHARCOAL)
add_textbox(slide5, sc_l + Inches(0.3), sc_t + Inches(0.58), sc_w - Inches(0.6), Inches(0.7), 
            "• 소속 총 발생 매출:  4,920,000 원\n"
            "• 플랫폼 수수료 (10%):  -492,000 원", 10, False, SLATE_GRAY)
add_textbox(slide5, sc_l + Inches(0.3), sc_t + Inches(1.35), sc_w - Inches(0.6), Inches(0.4), "최종 차주 실수령액: 4,428,000 원", 13, True, GREEN)

# 기여도 리스트
contrib_y = sc_t + Inches(2.0)
add_textbox(slide5, sc_l + Inches(0.2), contrib_y, sc_w - Inches(0.4), Inches(0.3), "기사별 매출 기여도", 11, True, SLATE_GRAY)

# 기사 1
ct1_y = contrib_y + Inches(0.3)
add_rounded_rect(slide5, sc_l + Inches(0.15), ct1_y, sc_w - Inches(0.3), Inches(0.75), WHITE, BORDER_GRAY)
add_textbox(slide5, sc_l + Inches(0.3), ct1_y + Inches(0.1), sc_w - Inches(0.6), Inches(0.35), "1. 김덤프 기사 (25톤)", 11, True, CHARCOAL)
add_textbox(slide5, sc_l + Inches(0.3), ct1_y + Inches(0.4), sc_w - Inches(0.6), Inches(0.3), "총 운임: 3,600,000 원  |  운행 횟수: 24회", 10, False, SLATE_GRAY)

# 기사 2
ct2_y = contrib_y + Inches(1.15)
add_rounded_rect(slide5, sc_l + Inches(0.15), ct2_y, sc_w - Inches(0.3), Inches(0.75), WHITE, BORDER_GRAY)
add_textbox(slide5, sc_l + Inches(0.3), ct2_y + Inches(0.1), sc_w - Inches(0.6), Inches(0.35), "2. 이배차 기사 (15톤)", 11, True, CHARCOAL)
add_textbox(slide5, sc_l + Inches(0.3), ct2_y + Inches(0.4), sc_w - Inches(0.6), Inches(0.3), "총 운임: 1,320,000 원  |  운행 횟수: 12회", 10, False, SLATE_GRAY)

# 엑셀 다운로드 버튼
excel_btn = add_rounded_rect(slide5, sc_l + Inches(0.2), sc_t + Inches(4.7), sc_w - Inches(0.4), Inches(0.6), CHARCOAL)
add_textbox(slide5, sc_l + Inches(0.2), sc_t + Inches(4.85), sc_w - Inches(0.4), Inches(0.4), "운송장 명세 Excel 전송", 14, True, WHITE, PP_ALIGN.CENTER)


# =============================================================
# SLIDE 6: Screen 5. 차주 필수 서류 관리 화면 (Owner Documents)
# =============================================================
slide6 = prs.slides.add_slide(blank_layout)
add_rectangle(slide6, Inches(0), Inches(0), Inches(13.33), Inches(7.5), WHITE)

# 좌측 설명 영역
add_textbox(slide6, Inches(0.8), Inches(0.8), Inches(5.5), Inches(0.8), "05. 차주 마스터 서류 보관함", 24, True, CHARCOAL)
add_textbox(slide6, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0), 
            "■ 주요 기능\n"
            "- 차주 가입 승인 및 최종 정산 계좌 연동을 위해 행정 서류를 업로드하고 상태를 검토하는 서류 관리 화면\n"
            "- 각 행정 서류의 심사 상태(승인완료/심사중/반려)를 실시간 피드백\n\n"
            "■ 주요 UI 구성\n"
            "1. 사업자 서류 업로드: 사업자등록증 사본 업로드 및 심사 상태 뱃지 표시\n"
            "2. 정산 통장 서류 업로드: 차주 본인/회사 통장 사본 업로드 및 상태 표시\n"
            "3. [제출 서류 수정 완료] 변경분 갱신 및 재심사 요청 버튼", 
            12, False, CHARCOAL)

# 우측 폰 프레임 생성
sc_l, sc_t, sc_w, sc_h = add_phone_frame(slide6, Inches(7.5), Inches(0.5), Inches(4.5), Inches(6.5), "차주 서류 보관함")

# 설명 가이드
add_textbox(slide6, sc_l + Inches(0.2), sc_t + Inches(0.2), sc_w - Inches(0.4), Inches(0.4), "차주 필수 서류 업로드 및 심사 상태", 11, True, SLATE_GRAY)

# 서류 1: 사업자등록증 사본 (승인 완료)
add_upload_box(slide6, sc_l + Inches(0.15), sc_t + Inches(0.7), sc_w - Inches(0.3), Inches(1.3), "사업자등록증 사본")
# 승인완료 초록색 뱃지
badge_s1 = add_rounded_rect(slide6, sc_l + sc_w - Inches(1.5), sc_t + Inches(0.8), Inches(1.2), Inches(0.32), GREEN)
add_textbox(slide6, sc_l + sc_w - Inches(1.5), sc_t + Inches(0.84), Inches(1.2), Inches(0.3), "승인 완료", 9, True, WHITE, PP_ALIGN.CENTER)

# 서류 2: 통장 사본 (심사 대기)
add_upload_box(slide6, sc_l + Inches(0.15), sc_t + Inches(2.2), sc_w - Inches(0.3), Inches(1.3), "통장 사본 (상호명 일치)")
# 심사대기 노란색 뱃지
badge_s2 = add_rounded_rect(slide6, sc_l + sc_w - Inches(1.5), sc_t + Inches(2.3), Inches(1.2), Inches(0.32), MUSTARD)
add_textbox(slide6, sc_l + sc_w - Inches(1.5), sc_t + Inches(2.34), Inches(1.2), Inches(0.3), "심사 대기", 9, True, CHARCOAL, PP_ALIGN.CENTER)

# 완료/갱신 버튼
save_btn = add_rounded_rect(slide6, sc_l + Inches(0.2), sc_t + Inches(4.7), sc_w - Inches(0.4), Inches(0.6), CHARCOAL)
add_textbox(slide6, sc_l + Inches(0.2), sc_t + Inches(4.85), sc_w - Inches(0.4), Inches(0.4), "제출 서류 수정 완료", 14, True, WHITE, PP_ALIGN.CENTER)


# 파워포인트 파일 저장
output_path = "D:\\Projects\\dumpring\\dumpring-platform-backend\\덤프링_차주_화면상세설계.pptx"
prs.save(output_path)
print(f"PPTX 파일 생성 성공: {output_path}")
