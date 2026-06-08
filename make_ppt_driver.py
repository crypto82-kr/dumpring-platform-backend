import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# 컬러 팔레트 정의 (라이트 머스터드 옐로우 테마)
WHITE = RGBColor(255, 255, 255)
SAND_BG = RGBColor(249, 248, 246)
MUSTARD = RGBColor(245, 158, 11)     # F59E0B (주요 테마 색상)
AMBER = RGBColor(217, 119, 6)        # D97706 (포인트 색상)
CHARCOAL = RGBColor(31, 41, 55)      # 1F2937 (메인 텍스트)
LIGHT_GRAY = RGBColor(243, 244, 246)  # F3F4F6 (영역 구분 배경)
BORDER_GRAY = RGBColor(229, 231, 235) # E5E7EB (카드 테두리)
SLATE_GRAY = RGBColor(75, 85, 99)    # 4B5563 (서브 텍스트)
RED = RGBColor(239, 68, 68)          # EF4444 (반려/에러/취소)
GREEN = RGBColor(16, 185, 129)       # 10B981 (정산완료/정상)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6] # 빈 레이아웃

# 유틸리티 함수: 텍스트 상자 추가
def add_textbox(slide, left, top, width, height, text, font_size=11, bold=False, color=CHARCOAL, align=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = 'Malgun Gothic'
    if align:
        p.alignment = align
    return txBox

# 유틸리티 함수: 직사각형 도형 추가
def add_rectangle(slide, left, top, width, height, fill_color, line_color=None, line_width=1):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

# 유틸리티 함수: 둥근 직사각형 도형 추가
def add_rounded_rect(slide, left, top, width, height, fill_color, line_color=None, line_width=1):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

# 유틸리티 함수: 모바일 폰 프레임 생성
def add_phone_frame(slide, left, top, width, height, title_text="덤프링"):
    bezel = add_rounded_rect(slide, left, top, width, height, SAND_BG, CHARCOAL, 3)
    screen_left = left + Inches(0.15)
    screen_top = top + Inches(0.4)
    screen_width = width - Inches(0.3)
    screen_height = height - Inches(0.6)
    screen = add_rectangle(slide, screen_left, screen_top, screen_width, screen_height, WHITE, BORDER_GRAY, 1)
    
    # 상단 헤더
    header = add_rectangle(slide, screen_left, screen_top, screen_width, Inches(0.5), SAND_BG, BORDER_GRAY, 1)
    add_textbox(slide, screen_left + Inches(0.15), screen_top + Inches(0.12), screen_width - Inches(0.3), Inches(0.3), title_text, 13, True, CHARCOAL)
    
    return screen_left, screen_top + Inches(0.5), screen_width, screen_height - Inches(0.5)

# 유틸리티 함수: 문서 제출/사진 촬영 업로드용 가상 버튼 그리기
def add_upload_box(slide, left, top, width, height, label_text):
    add_rounded_rect(slide, left, top, width, height, LIGHT_GRAY, BORDER_GRAY, 1)
    add_rectangle(slide, left + (width/2) - Inches(0.15), top + (height/2) - Inches(0.2), Inches(0.3), Inches(0.05), SLATE_GRAY)
    add_rectangle(slide, left + (width/2) - Inches(0.025), top + (height/2) - Inches(0.3), Inches(0.05), Inches(0.25), SLATE_GRAY)
    add_textbox(slide, left, top + (height/2) + Inches(0.05), width, Inches(0.35), label_text, 8, False, SLATE_GRAY, PP_ALIGN.CENTER)


# =============================================================
# SLIDE 1: 커버 페이지 (Cover Page)
# =============================================================
slide1 = prs.slides.add_slide(blank_layout)
add_rectangle(slide1, Inches(0), Inches(0), Inches(13.33), Inches(7.5), SAND_BG)
add_rectangle(slide1, Inches(1.5), Inches(2.2), Inches(0.15), Inches(2.5), MUSTARD)
add_textbox(slide1, Inches(2.0), Inches(2.1), Inches(9.5), Inches(1.2), "덤프링 (Dumpring)", 44, True, CHARCOAL)
add_textbox(slide1, Inches(2.0), Inches(3.2), Inches(9.5), Inches(0.8), "02. 덤프 기사(Driver) 업무 화면 상세설계서", 28, True, SLATE_GRAY)
add_textbox(slide1, Inches(2.0), Inches(4.1), Inches(9.5), Inches(0.5), "최종 테마: 라이트 머스터드 옐로우 (Light Mustard Yellow)", 16, False, AMBER)
add_textbox(slide1, Inches(2.0), Inches(5.8), Inches(6.0), Inches(0.8), "작성일: 2026. 06. 05\n작성자: 덤프링 플랫폼 기획팀", 12, False, SLATE_GRAY)


# =============================================================
# SLIDE 2: Screen 1A. 배차 목록 화면 (Order List)
# =============================================================
slide2 = prs.slides.add_slide(blank_layout)
add_rectangle(slide2, Inches(0), Inches(0), Inches(13.33), Inches(7.5), WHITE)

# 좌측 설명
add_textbox(slide2, Inches(0.8), Inches(0.8), Inches(5.5), Inches(0.8), "01. 배차 목록 화면 (홈)", 24, True, CHARCOAL)
add_textbox(slide2, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0), 
            "■ 주요 기능\n"
            "- 기사가 수락 가능한 신규 배차 일감 목록을 한눈에 확인하는 화면\n"
            "- 상단 검색바를 통해 원하는 현장명이나 자재명으로 필터링 기능 제공\n"
            "- 목록에서 오더 카드를 터치하면 상세 경로가 나오는 '상세조회(Screen 4)' 화면으로 이동\n\n"
            "■ UI 구성 요소\n"
            "1. 검색 입력창: 현장명, 지역, 자재명 검색 및 돋보기 아이콘\n"
            "2. 상단 토글 탭: [배차 목록 (선택)] / [완료 목록]\n"
            "3. 일감 카드 리스트:\n"
            "   - 출발지 ➔ 목적지, 자재(토사) 종류, 덤프 톤수 규격\n"
            "   - 배차 운임 및 지급 주체(상차지 지급 등) 정보 요약", 
            12, False, CHARCOAL)

# 우측 폰 프레임
sc_l, sc_t, sc_w, sc_h = add_phone_frame(slide2, Inches(7.5), Inches(0.5), Inches(4.5), Inches(6.5), "덤프링 - 배차 검색")

# 검색바
add_rounded_rect(slide2, sc_l + Inches(0.15), sc_t + Inches(0.15), sc_w - Inches(0.3), Inches(0.45), WHITE, BORDER_GRAY)
add_textbox(slide2, sc_l + Inches(0.3), sc_t + Inches(0.23), sc_w - Inches(0.6), Inches(0.3), "🔍 현장명, 지역, 자재명 검색", 11, False, SLATE_GRAY)

# 배차목록 / 완료목록 토글 탭
tab_y = sc_t + Inches(0.75)
add_rectangle(slide2, sc_l, tab_y, sc_w, Inches(0.45), LIGHT_GRAY, BORDER_GRAY)
# 배차목록 활성 탭
add_rectangle(slide2, sc_l, tab_y, sc_w/2, Inches(0.45), WHITE)
add_rectangle(slide2, sc_l, tab_y + Inches(0.4), sc_w/2, Inches(0.05), MUSTARD)
add_textbox(slide2, sc_l, tab_y + Inches(0.08), sc_w/2, Inches(0.3), "배차 목록", 11, True, CHARCOAL, PP_ALIGN.CENTER)
add_textbox(slide2, sc_l + sc_w/2, tab_y + Inches(0.08), sc_w/2, Inches(0.3), "완료 목록", 11, False, SLATE_GRAY, PP_ALIGN.CENTER)

# 일감 카드 1
card1_y = sc_t + Inches(1.35)
add_rounded_rect(slide2, sc_l + Inches(0.15), card1_y, sc_w - Inches(0.3), Inches(1.3), WHITE, BORDER_GRAY)
add_textbox(slide2, sc_l + Inches(0.25), card1_y + Inches(0.1), sc_w - Inches(0.5), Inches(0.3), "송도 1공구 아파트 ➔ 영종 사토장", 12, True, CHARCOAL)
add_textbox(slide2, sc_l + Inches(0.25), card1_y + Inches(0.38), sc_w - Inches(0.5), Inches(0.3), "• 자재: 양질토  |  규격: 25톤 덤프  |  이동: 18 km", 10, False, SLATE_GRAY)
add_textbox(slide2, sc_l + Inches(0.25), card1_y + Inches(0.85), sc_w - Inches(0.5), Inches(0.35), "150,000 원", 14, True, AMBER)
add_textbox(slide2, sc_l + sc_w - Inches(1.5), card1_y + Inches(0.88), Inches(1.2), Inches(0.3), "상차지 지급", 9, False, SLATE_GRAY, PP_ALIGN.RIGHT)

# 일감 카드 2
card2_y = sc_t + Inches(2.8)
add_rounded_rect(slide2, sc_l + Inches(0.15), card2_y, sc_w - Inches(0.3), Inches(1.3), WHITE, BORDER_GRAY)
add_textbox(slide2, sc_l + Inches(0.25), card2_y + Inches(0.1), sc_w - Inches(0.5), Inches(0.3), "인천대교 북단 현장 ➔ 송도 사토장", 12, True, CHARCOAL)
add_textbox(slide2, sc_l + Inches(0.25), card2_y + Inches(0.38), sc_w - Inches(0.5), Inches(0.3), "• 자재: 혼합토  |  규격: 15톤 덤프  |  이동: 10 km", 10, False, SLATE_GRAY)
add_textbox(slide2, sc_l + Inches(0.25), card2_y + Inches(0.85), sc_w - Inches(0.5), Inches(0.35), "110,000 원", 14, True, AMBER)
add_textbox(slide2, sc_l + sc_w - Inches(1.5), card2_y + Inches(0.88), Inches(1.2), Inches(0.3), "하차지 지급", 9, False, SLATE_GRAY, PP_ALIGN.RIGHT)


# =============================================================
# SLIDE 3: Screen 1B. 완료 목록 화면 (Completed History)
# =============================================================
slide3 = prs.slides.add_slide(blank_layout)
add_rectangle(slide3, Inches(0), Inches(0), Inches(13.33), Inches(7.5), WHITE)

# 좌측 설명
add_textbox(slide3, Inches(0.8), Inches(0.8), Inches(5.5), Inches(0.8), "02. 완료 목록 화면 (운행 이력)", 24, True, CHARCOAL)
add_textbox(slide3, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0), 
            "■ 주요 기능\n"
            "- 기사가 운행을 마치고 정산 처리된 과거 모든 오더 내역 확인\n"
            "- 정산 완료 여부에 따라 카드에 시각적인 뱃지(대기/완료) 표기\n"
            "- 기사가 스스로 일자별 총 수익과 누적 지급액을 대조해 볼 수 있는 투명한 이력 제공\n\n"
            "■ 주요 구성 요소\n"
            "1. 상단 토글 탭: [배차 목록] / [완료 목록 (선택)]\n"
            "2. 정산 여부 뱃지 명세:\n"
            "   - [정산완료 🟢]: 차주 및 기사 통장으로 입금 완료된 운행 건\n"
            "   - [정산대기 🟡]: 하차지 승인은 났으나 일괄 정산 주기 대기 건", 
            12, False, CHARCOAL)

# 우측 폰 프레임
sc_l, sc_t, sc_w, sc_h = add_phone_frame(slide3, Inches(7.5), Inches(0.5), Inches(4.5), Inches(6.5), "덤프링 - 운행 완료 이력")

# 검색바
add_rounded_rect(slide3, sc_l + Inches(0.15), sc_t + Inches(0.15), sc_w - Inches(0.3), Inches(0.45), WHITE, BORDER_GRAY)
add_textbox(slide3, sc_l + Inches(0.3), sc_t + Inches(0.23), sc_w - Inches(0.6), Inches(0.3), "🔍 날짜, 현장명 검색", 11, False, SLATE_GRAY)

# 배차목록 / 완료목록 토글 탭
tab_y = sc_t + Inches(0.75)
add_rectangle(slide3, sc_l, tab_y, sc_w, Inches(0.45), LIGHT_GRAY, BORDER_GRAY)
# 완료목록 활성 탭
add_rectangle(slide3, sc_l + sc_w/2, tab_y, sc_w/2, Inches(0.45), WHITE)
add_rectangle(slide3, sc_l + sc_w/2, tab_y + Inches(0.4), sc_w/2, Inches(0.05), MUSTARD)
add_textbox(slide3, sc_l, tab_y + Inches(0.08), sc_w/2, Inches(0.3), "배차 목록", 11, False, SLATE_GRAY, PP_ALIGN.CENTER)
add_textbox(slide3, sc_l + sc_w/2, tab_y + Inches(0.08), sc_w/2, Inches(0.3), "완료 목록", 11, True, CHARCOAL, PP_ALIGN.CENTER)

# 완료 카드 1 (정산완료)
card1_y = sc_t + Inches(1.35)
add_rounded_rect(slide3, sc_l + Inches(0.15), card1_y, sc_w - Inches(0.3), Inches(1.3), WHITE, BORDER_GRAY)
add_textbox(slide3, sc_l + Inches(0.25), card1_y + Inches(0.1), sc_w - Inches(0.5), Inches(0.3), "2026.06.05  |  티켓 #1094", 10, False, SLATE_GRAY)
add_textbox(slide3, sc_l + Inches(0.25), card1_y + Inches(0.35), sc_w - Inches(0.5), Inches(0.3), "송도 1공구 현장 ➔ 영종 사토장", 12, True, CHARCOAL)
add_textbox(slide3, sc_l + Inches(0.25), card1_y + Inches(0.85), sc_w - Inches(0.5), Inches(0.35), "150,000 원", 14, True, CHARCOAL)
# 정산완료 초록색 뱃지
badge1 = add_rounded_rect(slide3, sc_l + sc_w - Inches(1.4), card1_y + Inches(0.85), Inches(1.15), Inches(0.32), GREEN)
add_textbox(slide3, sc_l + sc_w - Inches(1.4), card1_y + Inches(0.89), Inches(1.15), Inches(0.3), "정산 완료", 9, True, WHITE, PP_ALIGN.CENTER)

# 완료 카드 2 (정산대기)
card2_y = sc_t + Inches(2.8)
add_rounded_rect(slide3, sc_l + Inches(0.15), card2_y, sc_w - Inches(0.3), Inches(1.3), WHITE, BORDER_GRAY)
add_textbox(slide3, sc_l + Inches(0.25), card2_y + Inches(0.1), sc_w - Inches(0.5), Inches(0.3), "2026.06.04  |  티켓 #1088", 10, False, SLATE_GRAY)
add_textbox(slide3, sc_l + Inches(0.25), card2_y + Inches(0.35), sc_w - Inches(0.5), Inches(0.3), "송도 1공구 현장 ➔ 영종 사토장", 12, True, CHARCOAL)
add_textbox(slide3, sc_l + Inches(0.25), card2_y + Inches(0.85), sc_w - Inches(0.5), Inches(0.35), "150,000 원", 14, True, CHARCOAL)
# 정산대기 노란색 뱃지
badge2 = add_rounded_rect(slide3, sc_l + sc_w - Inches(1.4), card2_y + Inches(0.85), Inches(1.15), Inches(0.32), MUSTARD)
add_textbox(slide3, sc_l + sc_w - Inches(1.4), card2_y + Inches(0.89), Inches(1.15), Inches(0.3), "정산 대기", 9, True, CHARCOAL, PP_ALIGN.CENTER)


# =============================================================
# SLIDE 4: Screen 2. 배차 상세 화면 (Order Detail)
# =============================================================
slide4 = prs.slides.add_slide(blank_layout)
add_rectangle(slide4, Inches(0), Inches(0), Inches(13.33), Inches(7.5), WHITE)

# 좌측 설명
add_textbox(slide4, Inches(0.8), Inches(0.8), Inches(5.5), Inches(0.8), "03. 배차 상세 화면 (수락 전)", 24, True, CHARCOAL)
add_textbox(slide4, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0), 
            "■ 주요 기능\n"
            "- 홈 목록에서 카드 선택 시 노출되며, 실제 일감의 상세 명세를 검토하는 단계\n"
            "- 지도로 상하차지 주행 경로선과 예상 거리를 먼저 검수하여 일의 강도 판단\n"
            "- 최종 확인 완료 시 하단 옐로우 버튼으로 배차 최종 계약 수락 처리\n\n"
            "■ 주요 정보 및 UI\n"
            "1. 경로 지도: 상차지(출발) ➔ 하차지(도착) 라우팅 라인\n"
            "2. 오더 상세 정보 패널:\n"
            "   - 상차 현장 및 사토지 주소, 운임 요금\n"
            "   - 토종(양질토), 차량 조건(25톤 덤프), 운송 계약 기본 사항\n"
            "3. [배차 수락] 고시인성 머스터드 황색 액션 버튼", 
            12, False, CHARCOAL)

# 우측 폰 프레임
sc_l, sc_t, sc_w, sc_h = add_phone_frame(slide4, Inches(7.5), Inches(0.5), Inches(4.5), Inches(6.5), "덤프링 - 배차 상세 조회")

# 지도 영역
map_rect = add_rectangle(slide4, sc_l, sc_t, sc_w, Inches(2.6), LIGHT_GRAY, BORDER_GRAY)
# 지도 위 출발지/도착지 텍스트 가이드
add_textbox(slide4, sc_l + Inches(0.2), sc_t + Inches(0.2), sc_w - Inches(0.4), Inches(0.3), "출발: 송도 현대 아파트 ➔ 도착: 영종 사토장", 10, True, CHARCOAL)
add_textbox(slide4, sc_l + Inches(0.2), sc_t + Inches(0.5), sc_w - Inches(0.4), Inches(0.3), "지도 시각화 및 예상 주행 거리: 18.2 km", 9, False, SLATE_GRAY)
# 지도 가로선 경로선 표현
add_rectangle(slide4, sc_l + Inches(1.0), sc_t + Inches(1.5), Inches(2.5), Inches(0.05), AMBER)

# 상세 패널
info_y = sc_t + Inches(2.7)
add_textbox(slide4, sc_l + Inches(0.2), info_y, sc_w - Inches(0.4), Inches(0.4), "상세 운송 요건", 13, True, CHARCOAL)
details_text = (
    "• 상차지: 송도 현대아파트 현장 (연락처: 010-***-****)\n"
    "• 하차지: 영종 남단 사토 매립지 (연락처: 010-***-****)\n"
    "• 토사 종류: 양질토 (성토/복토용 적합 토사)\n"
    "• 허용 규격: 25톤 덤프트럭 1대 모집\n"
    "• 지불 조건: 화주 상차지 지출  |  정산 방식: 당일 즉시 정산"
)
add_textbox(slide4, sc_l + Inches(0.2), info_y + Inches(0.4), sc_w - Inches(0.4), Inches(1.5), details_text, 10, False, CHARCOAL)

# 배차 금액 표시
add_textbox(slide4, sc_l + Inches(0.2), info_y + Inches(1.9), sc_w - Inches(0.4), Inches(0.4), "운송 단가: 150,000 원", 14, True, AMBER)

# 배차 수락 버튼
accept_btn = add_rounded_rect(slide4, sc_l + Inches(0.2), sc_t + Inches(4.7), sc_w - Inches(0.4), Inches(0.6), MUSTARD)
add_textbox(slide4, sc_l + Inches(0.2), sc_t + Inches(4.85), sc_w - Inches(0.4), Inches(0.4), "배차 수락", 14, True, CHARCOAL, PP_ALIGN.CENTER)


# =============================================================
# SLIDE 5: Screen 3. 실시간 GPS 미터기 주행 화면 (Meter Driving)
# =============================================================
slide5 = prs.slides.add_slide(blank_layout)
add_rectangle(slide5, Inches(0), Inches(0), Inches(13.33), Inches(7.5), WHITE)

# 좌측 설명
add_textbox(slide5, Inches(0.8), Inches(0.8), Inches(5.5), Inches(0.8), "04. 실시간 미터기 주행 화면", 24, True, CHARCOAL)
add_textbox(slide5, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0), 
            "■ 주요 기능\n"
            "- 운행 시작 클릭 시 백그라운드로 작동하는 실시간 위치 및 시간 요금 계산\n"
            "- 기사가 내비게이션 앱(Tmap 등)을 보더라도 화면 위에 플로팅 미터기로 금액 스트리밍\n"
            "- 하차지 지오펜싱(200m)에 도달하면 도착 완료 버튼 활성화\n\n"
            "■ 주요 UI\n"
            "1. 주행 요금 표시: 실시간 예상 금액을 대형 모노스페이스 골드 서체로 노출\n"
            "2. 가산 모드 뱃지: 거리 가산(>10km/h) 및 시간 가산(<10km/h) 정체 감지 자동 전환 표시\n"
            "3. 주행 계기판: 현재 속도, 이동한 실제 거리, 경과 시간 카운트", 
            12, False, CHARCOAL)

# 우측 폰 프레임
sc_l, sc_t, sc_w, sc_h = add_phone_frame(slide5, Inches(7.5), Inches(0.5), Inches(4.5), Inches(6.5), "덤프링 - 실시간 미터기")

# 지도 영역
map_rect = add_rectangle(slide5, sc_l, sc_t, sc_w, Inches(2.6), LIGHT_GRAY, BORDER_GRAY)
add_textbox(slide5, sc_l + Inches(0.2), sc_t + Inches(0.2), sc_w - Inches(0.4), Inches(0.4), "하차지 방향 실시간 경로 이동 중", 11, False, SLATE_GRAY, PP_ALIGN.CENTER)

# 미터기 요금 판넬
meter_rect = add_rounded_rect(slide5, sc_l + Inches(0.15), sc_t + Inches(2.8), sc_w - Inches(0.3), Inches(1.8), WHITE, BORDER_GRAY)
add_textbox(slide5, sc_l + Inches(0.3), sc_t + Inches(2.95), sc_w - Inches(0.6), Inches(0.3), "실시간 예상 운임", 12, False, SLATE_GRAY, PP_ALIGN.CENTER)
add_textbox(slide5, sc_l + Inches(0.3), sc_t + Inches(3.25), sc_w - Inches(0.6), Inches(0.6), "95,000 원", 28, True, AMBER, PP_ALIGN.CENTER)

# 계기판
add_textbox(slide5, sc_l + Inches(0.3), sc_t + Inches(4.0), sc_w - Inches(0.6), Inches(0.5), 
            "속도: 60 km/h   |   거리: 12.4 km   |   시간: 00:24:15", 10, False, CHARCOAL, PP_ALIGN.CENTER)

# 도착 완료 버튼
arrive_btn = add_rounded_rect(slide5, sc_l + Inches(0.2), sc_t + Inches(4.7), sc_w - Inches(0.4), Inches(0.6), MUSTARD)
add_textbox(slide5, sc_l + Inches(0.2), sc_t + Inches(4.85), sc_w - Inches(0.4), Inches(0.4), "하차지 도착 완료", 14, True, CHARCOAL, PP_ALIGN.CENTER)


# =============================================================
# SLIDE 6: Screen 4. 하차지 도착 및 지주 승인 대기 화면 (Gate Inflow)
# =============================================================
slide6 = prs.slides.add_slide(blank_layout)
add_rectangle(slide6, Inches(0), Inches(0), Inches(13.33), Inches(7.5), WHITE)

# 좌측 설명
add_textbox(slide6, Inches(0.8), Inches(0.8), Inches(5.5), Inches(0.8), "05. 하차지 도착 및 지주 승인 대기", 24, True, CHARCOAL)
add_textbox(slide6, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0), 
            "■ 주요 기능\n"
            "- 하차지 도착 신호 전송 후 사토장 지주가 반입 승인할 때까지 게이트 대기 단계\n"
            "- 지주 부재 시 강제 정산 접수를 위한 비상 사진 증빙 업로드 기능 탑재\n\n"
            "■ 비상 사진 증빙 프로세스\n"
            "- 지주가 자리에 없거나 통화 불가 시, 현장 상황을 증명하는 사진(계근대 수치, 차량 사토 모습)을 촬영해 업로드\n"
            "- 사진 업로드 시 70% 직권 정산 또는 추후 어드민 중재 정산 프로세스로 즉시 전환", 
            12, False, CHARCOAL)

# 우측 폰 프레임
sc_l, sc_t, sc_w, sc_h = add_phone_frame(slide6, Inches(7.5), Inches(0.5), Inches(4.5), Inches(6.5), "덤프링 - 지주 승인 대기")

# 헤더 텍스트
add_textbox(slide6, sc_l + Inches(0.2), sc_t + Inches(0.3), sc_w - Inches(0.4), Inches(0.4), "하차 반입 승인 대기 중", 15, True, CHARCOAL, PP_ALIGN.CENTER)

# 대기 스피너 박스
spinner_rect = add_rounded_rect(slide6, sc_l + Inches(0.15), sc_t + Inches(0.8), sc_w - Inches(0.3), Inches(1.8), WHITE, BORDER_GRAY)
add_rounded_rect(slide6, sc_l + Inches(1.8), sc_t + Inches(1.2), Inches(0.6), Inches(0.6), LIGHT_GRAY, AMBER, 2)
add_textbox(slide6, sc_l + Inches(0.3), sc_t + Inches(2.0), sc_w - Inches(0.6), Inches(0.5), "지주가 토사 검사 후 반입 승인을 진행 중입니다.", 10, False, CHARCOAL, PP_ALIGN.CENTER)

# 비상 사진 증빙 업로드 영역
add_textbox(slide6, sc_l + Inches(0.2), sc_t + Inches(2.8), sc_w - Inches(0.4), Inches(0.3), "지주 부재 시 비상 증빙 제출", 11, True, CHARCOAL)
add_upload_box(slide6, sc_l + Inches(0.15), sc_t + Inches(3.1), sc_w - Inches(0.3), Inches(1.3), "사토 현장 사진 촬영 및 첨부")

# 강제 완료 제출 버튼
submit_btn = add_rounded_rect(slide6, sc_l + Inches(0.2), sc_t + Inches(4.7), sc_w - Inches(0.4), Inches(0.6), MUSTARD)
add_textbox(slide6, sc_l + Inches(0.2), sc_t + Inches(4.85), sc_w - Inches(0.4), Inches(0.4), "증빙 사진 제출 및 완료", 14, True, CHARCOAL, PP_ALIGN.CENTER)


# =============================================================
# SLIDE 7: Screen 5. 운행 완료 전표 (Receipt/History)
# =============================================================
slide7 = prs.slides.add_slide(blank_layout)
add_rectangle(slide7, Inches(0), Inches(0), Inches(13.33), Inches(7.5), WHITE)

# 좌측 설명
add_textbox(slide7, Inches(0.8), Inches(0.8), Inches(5.5), Inches(0.8), "06. 운행 완료 전표 (영수증)", 24, True, CHARCOAL)
add_textbox(slide7, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0), 
            "■ 주요 기능\n"
            "- 운행이 최종 승인되어 즉시 지급 정산이 확정된 내역을 확인하는 영수증 화면\n"
            "- 수수료 공제 내역을 투명하게 안내하여 기사의 최종 실수령액 표시\n"
            "- 확인 터치 시 미터기 세션을 완전 종료하고 다시 홈 '배차 검색(Screen 2)'으로 복귀\n\n"
            "■ 주요 구성 요소\n"
            "1. 전자 전표 명세: 운임 발생 세부 내역 및 티켓 ID\n"
            "2. 정산금 공제 상세:\n"
            "   - 총 발생 운임 (예: 150,000원)\n"
            "   - 수수료 차감액 (예: -15,000원)\n"
            "   - 최종 지급 정산액 (예: 135,000원)", 
            12, False, CHARCOAL)

# 우측 폰 프레임
sc_l, sc_t, sc_w, sc_h = add_phone_frame(slide7, Inches(7.5), Inches(0.5), Inches(4.5), Inches(6.5), "덤프링 - 전자 전표")

# 영수증 카드
receipt_rect = add_rounded_rect(slide7, sc_l + Inches(0.15), sc_t + Inches(0.3), sc_w - Inches(0.3), Inches(4.2), WHITE, BORDER_GRAY)
add_textbox(slide7, sc_l + Inches(0.3), sc_t + Inches(0.5), sc_w - Inches(0.6), Inches(0.4), "운행 완료 영수증", 15, True, CHARCOAL, PP_ALIGN.CENTER)
add_textbox(slide7, sc_l + Inches(0.3), sc_t + Inches(0.9), sc_w - Inches(0.6), Inches(0.3), "티켓 번호: #1094  |  운행 완료", 9, False, SLATE_GRAY, PP_ALIGN.CENTER)

# 구분선
add_rectangle(slide7, sc_l + Inches(0.3), sc_t + Inches(1.3), sc_w - Inches(0.6), Inches(0.02), BORDER_GRAY)

# 영수증 본문 상세
receipt_txt = (
    "• 상차지: 송도 1공구 아파트 현장\n"
    "• 하차지: 영종 남단 사토 매립지\n"
    "• 운행 거리: 18.2 km\n"
    "• 소요 시간: 24분 15초\n"
    "• 운반 규격: 25톤 덤프 1회 운반"
)
add_textbox(slide7, sc_l + Inches(0.3), sc_t + Inches(1.5), sc_w - Inches(0.6), Inches(1.2), receipt_txt, 10, False, CHARCOAL)

# 정산 금액 계산 영역
add_rectangle(slide7, sc_l + Inches(0.3), sc_t + Inches(2.7), sc_w - Inches(0.6), Inches(0.02), BORDER_GRAY)
add_textbox(slide7, sc_l + Inches(0.3), sc_t + Inches(2.9), sc_w - Inches(0.6), Inches(0.3), "총 운임:  150,000 원", 11, False, CHARCOAL)
add_textbox(slide7, sc_l + Inches(0.3), sc_t + Inches(3.2), sc_w - Inches(0.6), Inches(0.3), "수수료 공제 (10%):  -15,000 원", 11, False, RED)
add_textbox(slide7, sc_l + Inches(0.3), sc_t + Inches(3.6), sc_w - Inches(0.6), Inches(0.4), "기사 최종 정산액: 135,000 원", 14, True, GREEN)

# 확인 및 홈으로 버튼
ok_btn = add_rounded_rect(slide7, sc_l + Inches(0.2), sc_t + Inches(4.7), sc_w - Inches(0.4), Inches(0.6), CHARCOAL)
add_textbox(slide7, sc_l + Inches(0.2), sc_t + Inches(4.85), sc_w - Inches(0.4), Inches(0.4), "확인 및 홈 복귀", 14, True, WHITE, PP_ALIGN.CENTER)


# 파워포인트 파일 저장
output_path = "D:\\Projects\\dumpring\\dumpring-platform-backend\\덤프링_기사_화면상세설계.pptx"
prs.save(output_path)
print(f"PPTX 파일 생성 성공: {output_path}")
