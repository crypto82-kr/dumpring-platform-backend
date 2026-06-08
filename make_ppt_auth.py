import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# 컬러 팔레트 정의 (라이트 머스터드 옐로우 테마)
WHITE = RGBColor(255, 255, 255)
SAND_BG = RGBColor(249, 248, 246)
MUSTARD = RGBColor(245, 158, 11)     # F59E0B (주요 테마 색상)
AMBER = RGBColor(217, 119, 6)        # D97706 (포인트 색상)
CHARCOAL = RGBColor(31, 41, 55)      # 1F2937 (메인 텍스트)
LIGHT_GRAY = RGBColor(243, 244, 246)  # F3F4F6 (영역 구분 배경)
BORDER_GRAY = RGBColor(229, 231, 235) # E5E7EB (카드 테두리)
SLATE_GRAY = RGBColor(75, 85, 99)    # 4B5563 (서브 텍스트)
RED = RGBColor(239, 68, 68)          # EF4444 (반려/취소/에러)
GREEN = RGBColor(16, 185, 129)       # 10B981 (승인/정상)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6] # 빈 레이아웃

# 유틸리티 함수: 텍스트 상자 추가
def add_textbox(slide, left, top, width, height, text, font_size=11, bold=False, color=CHARCOAL, align=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = 'Malgun Gothic'
    if align:
        p.alignment = align
    return txBox

# 유틸리티 함수: 직사각형 도형 추가
def add_rectangle(slide, left, top, width, height, fill_color, line_color=None, line_width=1):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

# 유틸리티 함수: 둥근 직사각형 도형 추가
def add_rounded_rect(slide, left, top, width, height, fill_color, line_color=None, line_width=1):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

# 유틸리티 함수: 모바일 폰 프레임 생성 (기본 사이즈)
def add_phone_frame(slide, left, top, width, height, title_text="덤프링"):
    # 외곽 디바이스 베젤
    bezel = add_rounded_rect(slide, left, top, width, height, SAND_BG, CHARCOAL, 3)
    # 스크린 영역
    screen_left = left + Inches(0.15)
    screen_top = top + Inches(0.4)
    screen_width = width - Inches(0.3)
    screen_height = height - Inches(0.6)
    screen = add_rectangle(slide, screen_left, screen_top, screen_width, screen_height, WHITE, BORDER_GRAY, 1)
    
    # 상단 상태바/헤더 영역
    header = add_rectangle(slide, screen_left, screen_top, screen_width, Inches(0.5), SAND_BG, BORDER_GRAY, 1)
    add_textbox(slide, screen_left + Inches(0.2), screen_top + Inches(0.12), screen_width - Inches(0.4), Inches(0.3), title_text, 14, True, CHARCOAL)
    
    return screen_left, screen_top + Inches(0.5), screen_width, screen_height - Inches(0.5)

# 유틸리티 함수: 모바일 폰 프레임 생성 (미니 사이즈 - 한 페이지에 여러 개 배치용)
def add_mini_phone(slide, left, top, width, height, title_text="덤프링"):
    # 디바이스 베젤
    add_rounded_rect(slide, left, top, width, height, SAND_BG, CHARCOAL, 2.5)
    # 스크린 영역
    scr_l = left + Inches(0.1)
    scr_t = top + Inches(0.3)
    scr_w = width - Inches(0.2)
    scr_h = height - Inches(0.4)
    add_rectangle(slide, scr_l, scr_t, scr_w, scr_h, WHITE, BORDER_GRAY, 1)
    
    # 상태바 및 헤더
    add_rectangle(slide, scr_l, scr_t, scr_w, Inches(0.35), SAND_BG, BORDER_GRAY, 1)
    add_textbox(slide, scr_l + Inches(0.1), scr_t + Inches(0.06), scr_w - Inches(0.2), Inches(0.25), title_text, 10, True, CHARCOAL)
    
    return scr_l, scr_t + Inches(0.35), scr_w, scr_h - Inches(0.35)

# 유틸리티 함수: 문서 제출/사진 촬영 업로드용 가상 버튼 그리기
def add_upload_box(slide, left, top, width, height, label_text):
    add_rounded_rect(slide, left, top, width, height, LIGHT_GRAY, BORDER_GRAY, 1)
    add_rectangle(slide, left + (width/2) - Inches(0.15), top + (height/2) - Inches(0.2), Inches(0.3), Inches(0.05), SLATE_GRAY)
    add_rectangle(slide, left + (width/2) - Inches(0.025), top + (height/2) - Inches(0.3), Inches(0.05), Inches(0.25), SLATE_GRAY)
    add_textbox(slide, left, top + (height/2) + Inches(0.05), width, Inches(0.35), label_text, 8, False, SLATE_GRAY, PP_ALIGN.CENTER)


# =============================================================
# SLIDE 1: 커버 페이지 (Cover Page)
# =============================================================
slide1 = prs.slides.add_slide(blank_layout)
add_rectangle(slide1, Inches(0), Inches(0), Inches(13.33), Inches(7.5), SAND_BG)
add_rectangle(slide1, Inches(1.5), Inches(2.2), Inches(0.15), Inches(2.5), MUSTARD)
add_textbox(slide1, Inches(2.0), Inches(2.1), Inches(9.5), Inches(1.2), "덤프링 (Dumpring) 상세설계", 44, True, CHARCOAL)
add_textbox(slide1, Inches(2.0), Inches(3.2), Inches(9.5), Inches(0.8), "01. 메인 로그인 및 역할별 회원가입 상세 설계서", 28, True, SLATE_GRAY)
add_textbox(slide1, Inches(2.0), Inches(4.1), Inches(9.5), Inches(0.5), "최종 테마: 라이트 머스터드 옐로우 (Light Mustard Yellow)", 16, False, AMBER)
add_textbox(slide1, Inches(2.0), Inches(5.8), Inches(6.0), Inches(0.8), "작성일: 2026. 06. 05\n작성자: 덤프링 플랫폼 기획팀", 12, False, SLATE_GRAY)


# =============================================================
# SLIDE 2: 메인 로그인 흐름 상세 (휴대폰 번호 입력 및 SMS 인증)
# =============================================================
slide2 = prs.slides.add_slide(blank_layout)
add_rectangle(slide2, Inches(0), Inches(0), Inches(13.33), Inches(7.5), WHITE)

add_textbox(slide2, Inches(0.8), Inches(0.8), Inches(4.5), Inches(0.8), "01. 메인 로그인 흐름 상세", 24, True, CHARCOAL)
add_textbox(slide2, Inches(0.8), Inches(1.8), Inches(4.5), Inches(5.0), 
            "■ 로그인 프로세스 상세 정의\n"
            "1단계: 휴대폰 번호 입력\n"
            "   - 가입 여부 관계없이 모든 사용자는 휴대폰 번호로 시작\n"
            "   - '010-1234-5678' 형식 자동 대시(-) 부여\n"
            "2단계: 인증번호 확인\n"
            "   - 문자(SMS)로 발송된 6자리 숫자 입력\n"
            "   - 3분(180초) 유효 시간 타이머 작동\n"
            "   - 번호 오기입 시 재전송 및 입력 초기화 지원\n\n"
            "■ 회원 상태 판정 흐름\n"
            "- 인증 완료 시 서버에서 기 가입 여부 검사\n"
            "  - [가존 회원]: 최종 선택 역할의 메인 홈 대시보드로 진입\n"
            "  - [신규 회원]: 회원가입 역할 선택 단계(Slide 3)로 이관", 
            12, False, CHARCOAL)

# 폰 1: 번호 입력 화면
p1_l, p1_t, p1_w, p1_h = add_mini_phone(slide2, Inches(5.8), Inches(0.5), Inches(3.2), Inches(6.5), "덤프링 - 본인 인증")
add_textbox(slide2, p1_l + Inches(0.15), p1_t + Inches(0.4), p1_w - Inches(0.3), Inches(0.4), "휴대폰 번호로 가입/로그인", 14, True, CHARCOAL)
add_textbox(slide2, p1_l + Inches(0.15), p1_t + Inches(0.75), p1_w - Inches(0.3), Inches(0.4), "안전한 운송 거래를 위해 번호 인증을 진행합니다.", 9, False, SLATE_GRAY)

# 번호 입력 필드
add_rounded_rect(slide2, p1_l + Inches(0.15), p1_t + Inches(1.6), p1_w - Inches(0.3), Inches(0.45), WHITE, BORDER_GRAY)
add_textbox(slide2, p1_l + Inches(0.25), p1_t + Inches(1.7), p1_w - Inches(0.5), Inches(0.35), "010-1234-5678", 12, True, CHARCOAL)

# 전송 버튼
btn_send = add_rounded_rect(slide2, p1_l + Inches(0.15), p1_t + Inches(2.2), p1_w - Inches(0.3), Inches(0.45), MUSTARD)
add_textbox(slide2, p1_l + Inches(0.15), p1_t + Inches(2.3), p1_w - Inches(0.3), Inches(0.3), "인증문자 발송", 11, True, CHARCOAL, PP_ALIGN.CENTER)


# 폰 2: 인증번호 입력 화면
p2_l, p2_t, p2_w, p2_h = add_mini_phone(slide2, Inches(9.4), Inches(0.5), Inches(3.2), Inches(6.5), "덤프링 - 인증번호 입력")
add_textbox(slide2, p2_l + Inches(0.15), p2_t + Inches(0.4), p2_w - Inches(0.3), Inches(0.4), "인증번호 입력", 14, True, CHARCOAL)
add_textbox(slide2, p2_l + Inches(0.15), p2_t + Inches(0.75), p2_w - Inches(0.3), Inches(0.4), "문자로 전송된 6자리 숫자를 기입해 주세요.", 9, False, SLATE_GRAY)

# 인증번호 입력 칸 (6칸 가상 표현)
for i in range(6):
    box_x = p2_l + Inches(0.15) + Inches(i * 0.45)
    add_rounded_rect(slide2, box_x, p2_t + Inches(1.6), Inches(0.38), Inches(0.45), WHITE, BORDER_GRAY)
    add_textbox(slide2, box_x, p2_t + Inches(1.68), Inches(0.38), Inches(0.35), str((3, 4, 5, 6, 7, 8)[i]), 14, True, CHARCOAL, PP_ALIGN.CENTER)

# 유효시간 표시
add_textbox(slide2, p2_l + Inches(0.15), p2_t + Inches(2.2), p2_w - Inches(0.3), Inches(0.3), "남은 시간 02:45  |  인증번호 재전송", 9, False, RED)

# 완료 버튼
btn_verify = add_rounded_rect(slide2, p2_l + Inches(0.15), p2_t + Inches(2.8), p2_w - Inches(0.3), Inches(0.45), MUSTARD)
add_textbox(slide2, p2_l + Inches(0.15), p2_t + Inches(2.9), p2_w - Inches(0.3), Inches(0.3), "인증 및 완료", 11, True, CHARCOAL, PP_ALIGN.CENTER)


# =============================================================
# SLIDE 3: 회원가입 역할 선택 화면
# =============================================================
slide3 = prs.slides.add_slide(blank_layout)
add_rectangle(slide3, Inches(0), Inches(0), Inches(13.33), Inches(7.5), WHITE)

add_textbox(slide3, Inches(0.8), Inches(0.8), Inches(5.0), Inches(0.8), "02. 회원가입 - 역할 선택", 24, True, CHARCOAL)
add_textbox(slide3, Inches(0.8), Inches(1.8), Inches(5.0), Inches(5.0), 
            "■ 최초 가입 시 역할 지정 상세\n"
            "- 본인 인증 완료 후 최초 가입 시 1회 필수 지정\n"
            "- 역할에 따라 업로드해야 하는 법적 서류 및 입력 정보 상이\n"
            "- 오지정을 방지하기 위해 각 역할 카드 밑에 명확한 한 줄 가이드라인 배치\n\n"
            "■ 역할 정의 및 가이드\n"
            "1. 덤프 기사: 현장에서 실시간 미터기를 켜고 덤프 운행을 하시는 기사님\n"
            "2. 덤프 차주: 트럭을 소유하고 소속 기사 및 차량을 총괄 관리하시는 사장님\n"
            "3. 현장 관리자: 공사현장 개설 및 배차 오더를 직접 끊어 발행하는 소장님\n"
            "4. 현장 담당자: 현장 게이트에서 차량 출입 통제 및 상차 완료를 체크하는 직원\n"
            "5. 하차지 이용자: 사토장을 소유하고 흙 반입 승인을 처리하는 지주/사장님", 
            12, False, CHARCOAL)

# 우측 폰 프레임
p3_l, p3_t, p3_w, p3_h = add_phone_frame(slide3, Inches(7.0), Inches(0.5), Inches(4.5), Inches(6.5), "덤프링 - 역할 지정")
add_textbox(slide3, p3_l + Inches(0.2), p3_t + Inches(0.3), p3_w - Inches(0.4), Inches(0.4), "회원 종류를 선택해 주세요", 15, True, CHARCOAL)
add_textbox(slide3, p3_l + Inches(0.2), p3_t + Inches(0.65), p3_w - Inches(0.4), Inches(0.35), "가입 후에는 역할 변경 시 서류 재심사가 필요합니다.", 9, False, SLATE_GRAY)

# 5개 카드 그리드 배치 (둥근 사각형 + 텍스트 가이드)
roles_details = [
    ("🚜 덤프 기사", "실시간 미터기 켜고 덤프트럭 운행"),
    ("🏢 덤프 차주", "기사 초대/차량 관리 및 세금 정산 총괄"),
    ("🏗️ 현장 관리자 (소장)", "공사현장 개설 및 매칭 오더 등록 발행"),
    ("👷 현장 담당자 (직원)", "현장 게이트 차량 출입 및 출발 승인"),
    ("🏕️ 하차지 이용자 (지주)", "사토장 등록 및 실시간 게이트 반입 승인")
]

for idx, (role_name, role_guide) in enumerate(roles_details):
    ry = p3_t + Inches(1.1) + Inches(idx * 0.88)
    is_active = (idx == 0) # 첫 번째 카드 활성 상태 표현
    f_color = WHITE
    l_color = BORDER_GRAY
    if is_active:
        f_color = RGBColor(254, 243, 199)
        l_color = MUSTARD
    card = add_rounded_rect(slide3, p3_l + Inches(0.15), ry, p3_w - Inches(0.3), Inches(0.8), f_color, l_color)
    add_textbox(slide3, p3_l + Inches(0.3), ry + Inches(0.1), p3_w - Inches(0.6), Inches(0.3), role_name, 11, True, CHARCOAL)
    add_textbox(slide3, p3_l + Inches(0.3), ry + Inches(0.42), p3_w - Inches(0.6), Inches(0.3), role_guide, 9, False, SLATE_GRAY)


# =============================================================
# SLIDE 4: 덤프 기사 (Driver) 회원가입 및 서류 제출 상세
# =============================================================
slide4 = prs.slides.add_slide(blank_layout)
add_rectangle(slide4, Inches(0), Inches(0), Inches(13.33), Inches(7.5), WHITE)

add_textbox(slide4, Inches(0.8), Inches(0.8), Inches(4.5), Inches(0.8), "03. 덤프 기사 (Driver) 회원가입", 24, True, CHARCOAL)
add_textbox(slide4, Inches(0.8), Inches(1.8), Inches(4.5), Inches(5.0), 
            "■ 기사 회원가입 요구사항\n"
            "- 기본 인적사항 입력 및 승인을 위한 3대 필수 서류 제출 단계 구성\n"
            "- 제출된 서류는 플랫폼 관리자 심사 완료 후 오더 수락 권한 활성화\n\n"
            "■ 상세 입력 항목 및 서류 명세\n"
            "1. 기사명 (실명 확인)\n"
            "2. 휴대폰 번호 (인증 번호 자동 매핑)\n"
            "3. 3대 필수 심사 서류 (카메라 촬영 첨부):\n"
            "   - 1종 대형 운전면허증\n"
            "   - 화물자동차 운송종사자 자격증\n"
            "   - 정산용 본인 명의 통장 사본\n\n"
            "■ 심사 상태 관리\n"
            "- 최초 제출 후 [대기 상태]로 진입하며 기사 대시보드 진입 시 '심사 중' 서브 화면 노출", 
            12, False, CHARCOAL)

# 폰 1: 기사 정보 입력
p4_l, p4_t, p4_w, p4_h = add_mini_phone(slide4, Inches(5.8), Inches(0.5), Inches(3.2), Inches(6.5), "기사 가입 - 정보 입력")
add_textbox(slide4, p4_l + Inches(0.15), p4_t + Inches(0.4), p4_w - Inches(0.3), Inches(0.4), "기사 정보 입력", 14, True, CHARCOAL)

# 입력 필드 2종 (이름, 면허 번호)
add_textbox(slide4, p4_l + Inches(0.15), p4_t + Inches(1.0), p4_w - Inches(0.3), Inches(0.3), "기사명 (실명)", 9, True, SLATE_GRAY)
add_rounded_rect(slide4, p4_l + Inches(0.15), p4_t + Inches(1.3), p4_w - Inches(0.3), Inches(0.4), WHITE, BORDER_GRAY)
add_textbox(slide4, p4_l + Inches(0.25), p4_t + Inches(1.38), p4_w - Inches(0.5), Inches(0.3), "홍길동", 11, True, CHARCOAL)

add_textbox(slide4, p4_l + Inches(0.15), p4_t + Inches(1.9), p4_w - Inches(0.3), Inches(0.3), "면허 등급", 9, True, SLATE_GRAY)
add_rounded_rect(slide4, p4_l + Inches(0.15), p4_t + Inches(2.2), p4_w - Inches(0.3), Inches(0.4), WHITE, BORDER_GRAY)
add_textbox(slide4, p4_l + Inches(0.25), p4_t + Inches(2.28), p4_w - Inches(0.5), Inches(0.3), "1종 대형 면허", 11, True, CHARCOAL)

btn_next = add_rounded_rect(slide4, p4_l + Inches(0.15), p4_t + Inches(4.8), p4_w - Inches(0.3), Inches(0.45), MUSTARD)
add_textbox(slide4, p4_l + Inches(0.15), p4_t + Inches(4.9), p4_w - Inches(0.3), Inches(0.3), "다음 단계로", 11, True, CHARCOAL, PP_ALIGN.CENTER)


# 폰 2: 기사 서류 제출
p5_l, p5_t, p5_w, p5_h = add_mini_phone(slide4, Inches(9.4), Inches(0.5), Inches(3.2), Inches(6.5), "기사 가입 - 서류 제출")
add_textbox(slide4, p5_l + Inches(0.15), p5_t + Inches(0.4), p5_w - Inches(0.3), Inches(0.4), "필수 서류 업로드", 14, True, CHARCOAL)
add_textbox(slide4, p5_l + Inches(0.15), p5_t + Inches(0.75), p5_w - Inches(0.3), Inches(0.4), "기사 승인을 위한 필수 파일 3종을 첨부해 주세요.", 9, False, SLATE_GRAY)

# 서류 업로드 영역 3개 가로 배치 또는 세로형 3개
for idx, doc in enumerate(["운전면허증 사본", "화물운송자격증 사본", "기사 본인 명의 통장 사본"]):
    dy = p5_t + Inches(1.2) + Inches(idx * 1.1)
    add_upload_box(slide4, p5_l + Inches(0.15), dy, p5_w - Inches(0.3), Inches(0.85), doc)

# 제출 버튼
btn_submit = add_rounded_rect(slide4, p5_l + Inches(0.15), p5_t + Inches(4.8), p5_w - Inches(0.3), Inches(0.45), MUSTARD)
add_textbox(slide4, p5_l + Inches(0.15), p5_t + Inches(4.9), p5_w - Inches(0.3), Inches(0.3), "제출 및 가입 신청", 11, True, CHARCOAL, PP_ALIGN.CENTER)


# =============================================================
# SLIDE 5: 덤프 차주 (Owner) 회원가입 및 서류 제출 상세
# =============================================================
slide5 = prs.slides.add_slide(blank_layout)
add_rectangle(slide5, Inches(0), Inches(0), Inches(13.33), Inches(7.5), WHITE)

# 설명
add_textbox(slide5, Inches(0.8), Inches(0.8), Inches(4.5), Inches(0.8), "04. 덤프 차주 (Owner) 회원가입", 24, True, CHARCOAL)
add_textbox(slide5, Inches(0.8), Inches(1.8), Inches(4.5), Inches(5.0), 
            "■ 차주 회원가입 요구사항\n"
            "- 소속 차량 배차 및 정산 청구 주체로서 사업체 검증 단계 구성\n"
            "- 차주 명의와 사업자 정보 검증 완료 후 기사 등록 및 배차 정산 기능 오픈\n\n"
            "■ 상세 입력 항목 및 서류 명세\n"
            "1. 상호명 (사업체 명칭)\n"
            "2. 대표자명\n"
            "3. 사업자등록번호\n"
            "4. 필수 증빙 서류 업로드:\n"
            "   - 사업자등록증 사본\n"
            "   - 정산 및 청구용 회사/차주 명의 통장 사본\n\n"
            "■ UI 설계 요점\n"
            "- 개인사업자와 법인사업자 구분을 위한 최상단 토글 스위치 도입", 
            12, False, CHARCOAL)

# 폰 1: 차주 사업자 정보 입력
p1_l, p1_t, p1_w, p1_h = add_mini_phone(slide5, Inches(5.8), Inches(0.5), Inches(3.2), Inches(6.5), "차주 가입 - 사업자 정보")
add_textbox(slide5, p1_l + Inches(0.15), p1_t + Inches(0.4), p1_w - Inches(0.3), Inches(0.4), "사업자 정보 입력", 14, True, CHARCOAL)

# 사업자 구분 탭 토글
toggle_rect = add_rounded_rect(slide5, p1_l + Inches(0.15), p1_t + Inches(0.9), p1_w - Inches(0.3), Inches(0.45), LIGHT_GRAY, BORDER_GRAY)
add_rounded_rect(slide5, p1_l + Inches(0.15), p1_t + Inches(0.9), p1_w/2 - Inches(0.15), Inches(0.45), WHITE, BORDER_GRAY)
add_textbox(slide5, p1_l + Inches(0.15), p1_t + Inches(1.0), p1_w/2 - Inches(0.15), Inches(0.35), "개인사업자", 10, True, CHARCOAL, PP_ALIGN.CENTER)
add_textbox(slide5, p1_l + p1_w/2, p1_t + Inches(1.0), p1_w/2 - Inches(0.15), Inches(0.35), "법인사업자", 10, False, SLATE_GRAY, PP_ALIGN.CENTER)

# 상호명 필드
add_textbox(slide5, p1_l + Inches(0.15), p1_t + Inches(1.5), p1_w - Inches(0.3), Inches(0.3), "상호명", 9, True, SLATE_GRAY)
add_rounded_rect(slide5, p1_l + Inches(0.15), p1_t + Inches(1.8), p1_w - Inches(0.3), Inches(0.4), WHITE, BORDER_GRAY)
add_textbox(slide5, p1_l + Inches(0.25), p1_t + Inches(1.88), p1_w - Inches(0.5), Inches(0.35), "(주) 대박 물류", 11, True, CHARCOAL)

# 사업자번호 필드
add_textbox(slide5, p1_l + Inches(0.15), p1_t + Inches(2.4), p1_w - Inches(0.3), Inches(0.3), "사업자등록번호", 9, True, SLATE_GRAY)
add_rounded_rect(slide5, p1_l + Inches(0.15), p1_t + Inches(2.7), p1_w - Inches(0.3), Inches(0.4), WHITE, BORDER_GRAY)
add_textbox(slide5, p1_l + Inches(0.25), p1_t + Inches(2.78), p1_w - Inches(0.5), Inches(0.35), "120-81-12345", 11, True, CHARCOAL)

btn_next = add_rounded_rect(slide5, p1_l + Inches(0.15), p1_t + Inches(4.8), p1_w - Inches(0.3), Inches(0.45), MUSTARD)
add_textbox(slide5, p1_l + Inches(0.15), p1_t + Inches(4.9), p1_w - Inches(0.3), Inches(0.3), "다음 단계로", 11, True, CHARCOAL, PP_ALIGN.CENTER)


# 폰 2: 차주 서류 제출
p2_l, p2_t, p2_w, p2_h = add_mini_phone(slide5, Inches(9.4), Inches(0.5), Inches(3.2), Inches(6.5), "차주 가입 - 서류 제출")
add_textbox(slide5, p2_l + Inches(0.15), p2_t + Inches(0.4), p2_w - Inches(0.3), Inches(0.4), "사업자 서류 제출", 14, True, CHARCOAL)
add_textbox(slide5, p2_l + Inches(0.15), p2_t + Inches(0.75), p2_w - Inches(0.3), Inches(0.4), "정산 및 세무 증빙 서류 2종을 등록해 주세요.", 9, False, SLATE_GRAY)

# 서류 2종
add_upload_box(slide5, p2_l + Inches(0.15), p2_t + Inches(1.3), p2_w - Inches(0.3), Inches(1.2), "사업자등록증 사본 업로드")
add_upload_box(slide5, p2_l + Inches(0.15), p2_t + Inches(2.8), p2_w - Inches(0.3), Inches(1.2), "통장 사본 업로드 (상호명/대표자 일치)")

btn_submit = add_rounded_rect(slide5, p2_l + Inches(0.15), p2_t + Inches(4.8), p2_w - Inches(0.3), Inches(0.45), MUSTARD)
add_textbox(slide5, p2_l + Inches(0.15), p2_t + Inches(4.9), p2_w - Inches(0.3), Inches(0.3), "제출 및 가입 완료", 11, True, CHARCOAL, PP_ALIGN.CENTER)


# =============================================================
# SLIDE 6: 현장 관리자 (Site Manager) 회원가입 및 현장 개설
# =============================================================
slide6 = prs.slides.add_slide(blank_layout)
add_rectangle(slide6, Inches(0), Inches(0), Inches(13.33), Inches(7.5), WHITE)

# 설명
add_textbox(slide6, Inches(0.8), Inches(0.8), Inches(4.5), Inches(0.8), "05. 현장 관리자 (소장) 회원가입", 24, True, CHARCOAL)
add_textbox(slide6, Inches(0.8), Inches(1.8), Inches(4.5), Inches(5.0), 
            "■ 현장 관리자 회원가입 요구사항\n"
            "- 공사현장을 개설하고 매칭 오더를 발행하는 주체\n"
            "- 건설사 정보 입력 및 최초 현장 1개를 가입 과정에서 즉시 등록\n\n"
            "■ 상세 입력 항목 및 프로세스\n"
            "1. 소속 건설사명\n"
            "2. 신규 개설할 현장 주소 및 명칭\n"
            "3. 차량 지오펜싱(출입 인식) 반경 설정 (기본 200m)\n"
            "4. 증빙 서류: 건설업 등록증 또는 계약 공문서\n\n"
            "■ UI/UX 설계 요점\n"
            "- 지도 기반의 현장 주소 검색 기능과 지오펜싱 출입 범위 선택 슬라이더 배치", 
            12, False, CHARCOAL)

# 폰 1: 관리자 인적 사항
p1_l, p1_t, p1_w, p1_h = add_mini_phone(slide6, Inches(5.8), Inches(0.5), Inches(3.2), Inches(6.5), "소장 가입 - 회사 정보")
add_textbox(slide6, p1_l + Inches(0.15), p1_t + Inches(0.4), p1_w - Inches(0.3), Inches(0.4), "현장 관리자 가입", 14, True, CHARCOAL)

# 건설사명
add_textbox(slide6, p1_l + Inches(0.15), p1_t + Inches(1.1), p1_w - Inches(0.3), Inches(0.3), "소속 건설사명", 9, True, SLATE_GRAY)
add_rounded_rect(slide6, p1_l + Inches(0.15), p1_t + Inches(1.4), p1_w - Inches(0.3), Inches(0.4), WHITE, BORDER_GRAY)
add_textbox(slide6, p1_l + Inches(0.25), p1_t + Inches(1.48), p1_w - Inches(0.5), Inches(0.3), "대우건설 (주)", 11, True, CHARCOAL)

# 직급
add_textbox(slide6, p1_l + Inches(0.15), p1_t + Inches(2.0), p1_w - Inches(0.3), Inches(0.3), "현장 직급/직책", 9, True, SLATE_GRAY)
add_rounded_rect(slide6, p1_l + Inches(0.15), p1_t + Inches(2.3), p1_w - Inches(0.3), Inches(0.4), WHITE, BORDER_GRAY)
add_textbox(slide6, p1_l + Inches(0.25), p1_t + Inches(2.38), p1_w - Inches(0.5), Inches(0.3), "현장 소장", 11, True, CHARCOAL)

btn_next = add_rounded_rect(slide6, p1_l + Inches(0.15), p1_t + Inches(4.8), p1_w - Inches(0.3), Inches(0.45), MUSTARD)
add_textbox(slide6, p1_l + Inches(0.15), p1_t + Inches(4.9), p1_w - Inches(0.3), Inches(0.3), "다음 (현장 등록)", 11, True, CHARCOAL, PP_ALIGN.CENTER)


# 폰 2: 현장 지도 등록
p2_l, p2_t, p2_w, p2_h = add_mini_phone(slide6, Inches(9.4), Inches(0.5), Inches(3.2), Inches(6.5), "소장 가입 - 현장 등록")
add_textbox(slide6, p2_l + Inches(0.15), p2_t + Inches(0.4), p2_w - Inches(0.3), Inches(0.4), "신규 공사현장 개설", 14, True, CHARCOAL)

# 현장명
add_textbox(slide6, p2_l + Inches(0.15), p2_t + Inches(0.9), p2_w - Inches(0.3), Inches(0.3), "공사 현장명", 9, True, SLATE_GRAY)
add_rounded_rect(slide6, p2_l + Inches(0.15), p2_t + Inches(1.15), p2_w - Inches(0.3), Inches(0.4), WHITE, BORDER_GRAY)
add_textbox(slide6, p2_l + Inches(0.25), p2_t + Inches(1.23), p2_w - Inches(0.5), Inches(0.3), "송도 1공구 아파트 신축 현장", 10, True, CHARCOAL)

# 가상 지도 영역
add_rectangle(slide6, p2_l + Inches(0.15), p2_t + Inches(1.7), p2_w - Inches(0.3), Inches(1.3), LIGHT_GRAY, BORDER_GRAY)
add_textbox(slide6, p2_l + Inches(0.2), p2_t + Inches(2.2), p2_w - Inches(0.4), Inches(0.4), "지도 주소 검색 및 지오펜싱 지정", 9, False, SLATE_GRAY, PP_ALIGN.CENTER)

# 지오펜싱 슬라이더 표현
add_textbox(slide6, p2_l + Inches(0.15), p2_t + Inches(3.1), p2_w - Inches(0.3), Inches(0.3), "출입 인식 반경 설정: 200m", 9, True, CHARCOAL)
add_rectangle(slide6, p2_l + Inches(0.15), p2_t + Inches(3.4), p2_w - Inches(0.3), Inches(0.05), BORDER_GRAY)
add_rounded_rect(slide6, p2_l + Inches(0.7), p2_t + Inches(3.32), Inches(0.2), Inches(0.2), MUSTARD)

# 문서 등록 박스
add_upload_box(slide6, p2_l + Inches(0.15), p2_t + Inches(3.7), p2_w - Inches(0.3), Inches(0.85), "공사현장 허가 공문 첨부")

btn_submit = add_rounded_rect(slide6, p2_l + Inches(0.15), p2_t + Inches(4.8), p2_w - Inches(0.3), Inches(0.45), MUSTARD)
add_textbox(slide6, p2_l + Inches(0.15), p2_t + Inches(4.9), p2_w - Inches(0.3), Inches(0.3), "가입 및 현장 개설 신청", 11, True, CHARCOAL, PP_ALIGN.CENTER)


# =============================================================
# SLIDE 7: 현장 담당자 (Site Worker) 회원가입 및 소장 매핑 신청
# =============================================================
slide7 = prs.slides.add_slide(blank_layout)
add_rectangle(slide7, Inches(0), Inches(0), Inches(13.33), Inches(7.5), WHITE)

# 설명
add_textbox(slide7, Inches(0.8), Inches(0.8), Inches(4.5), Inches(0.8), "06. 현장 담당자 (직원) 회원가입", 24, True, CHARCOAL)
add_textbox(slide7, Inches(0.8), Inches(1.8), Inches(4.5), Inches(5.0), 
            "■ 현장 담당자 회원가입 요구사항\n"
            "- 게이트에서 차량 입출차를 수동 기록하는 직원 역할\n"
            "- 가입 시 본인의 근무 현장을 검색하여 소속 신청을 넣음\n"
            "- 해당 현장의 관리자(소장)가 승인을 누르기 전까지 대기\n\n"
            "■ 상세 가입 프로세스\n"
            "1. 담당자 실명 입력\n"
            "2. 근무 현장 검색:\n"
            "   - 현장 명칭이나 건설사로 활성화된 현장 데이터 조회\n"
            "3. 소속 현장 매핑 요청\n"
            "4. 소장 승인 대기 화면 노출\n\n"
            "■ UI 설계 요점\n"
            "- 소속 현장 매핑 대기 중에는 로딩/대기 중 위젯을 중앙에 강조 표기", 
            12, False, CHARCOAL)

# 폰 1: 담당자 가입 및 현장 검색
p1_l, p1_t, p1_w, p1_h = add_mini_phone(slide7, Inches(5.8), Inches(0.5), Inches(3.2), Inches(6.5), "직원 가입 - 현장 지정")
add_textbox(slide7, p1_l + Inches(0.15), p1_t + Inches(0.4), p1_w - Inches(0.3), Inches(0.4), "현장 담당자 등록", 14, True, CHARCOAL)

# 실명 입력
add_textbox(slide7, p1_l + Inches(0.15), p1_t + Inches(1.0), p1_w - Inches(0.3), Inches(0.3), "직원 실명", 9, True, SLATE_GRAY)
add_rounded_rect(slide7, p1_l + Inches(0.15), p1_t + Inches(1.3), p1_w - Inches(0.3), Inches(0.4), WHITE, BORDER_GRAY)
add_textbox(slide7, p1_l + Inches(0.25), p1_t + Inches(1.38), p1_w - Inches(0.5), Inches(0.3), "홍길동 대리", 11, True, CHARCOAL)

# 현장 검색 필드
add_textbox(slide7, p1_l + Inches(0.15), p1_t + Inches(1.9), p1_w - Inches(0.3), Inches(0.3), "소속 공사현장 검색", 9, True, SLATE_GRAY)
add_rounded_rect(slide7, p1_l + Inches(0.15), p1_t + Inches(2.2), p1_w - Inches(0.3), Inches(0.4), WHITE, BORDER_GRAY)
add_textbox(slide7, p1_l + Inches(0.25), p1_t + Inches(2.28), p1_w - Inches(0.5), Inches(0.3), "송도 현대 아파트", 11, True, CHARCOAL)

# 검색 리스트 1개 노출
res_card = add_rounded_rect(slide7, p1_l + Inches(0.15), p1_t + Inches(2.7), p1_w - Inches(0.3), Inches(0.7), LIGHT_GRAY, MUSTARD)
add_textbox(slide7, p1_l + Inches(0.25), p1_t + Inches(2.8), p1_w - Inches(0.5), Inches(0.45), "송도 현대 아파트 신축공사 현장\n(소장: 김소장  |  대우건설)", 10, False, CHARCOAL)

btn_next = add_rounded_rect(slide7, p1_l + Inches(0.15), p1_t + Inches(4.8), p1_w - Inches(0.3), Inches(0.45), MUSTARD)
add_textbox(slide7, p1_l + Inches(0.15), p1_t + Inches(4.9), p1_w - Inches(0.3), Inches(0.3), "현장 매핑 신청", 11, True, CHARCOAL, PP_ALIGN.CENTER)


# 폰 2: 매핑 승인 대기
p2_l, p2_t, p2_w, p2_h = add_mini_phone(slide7, Inches(9.4), Inches(0.5), Inches(3.2), Inches(6.5), "직원 가입 - 승인 대기")
add_textbox(slide7, p2_l + Inches(0.15), p2_t + Inches(0.4), p2_w - Inches(0.3), Inches(0.4), "현장 승인 대기 중", 14, True, CHARCOAL, PP_ALIGN.CENTER)

# 대기 상태 애니메이션/아이콘 표현용 서클
add_rounded_rect(slide7, p2_l + Inches(1.1), p2_t + Inches(1.5), Inches(0.8), Inches(0.8), LIGHT_GRAY, AMBER, 2)
add_textbox(slide7, p2_l + Inches(0.2), p2_t + Inches(2.5), p2_w - Inches(0.4), Inches(1.2), 
            "현장 소장님의 승인을 대기 중입니다.\n\n승인이 완료되면 자동으로 대시보드로 진입합니다.", 11, False, CHARCOAL, PP_ALIGN.CENTER)

# 소장 정보 표시
info_card = add_rounded_rect(slide7, p2_l + Inches(0.15), p2_t + Inches(4.0), p2_w - Inches(0.3), Inches(1.0), WHITE, BORDER_GRAY)
add_textbox(slide7, p2_l + Inches(0.3), p2_t + Inches(4.15), p2_w - Inches(0.6), Inches(0.8), 
            "• 신청 현장: 송도 현대 아파트\n• 소장 연락처: 010-9999-8888\n  (빠른 승인은 소장님께 유선 통화 요청)", 9, False, SLATE_GRAY)


# =============================================================
# SLIDE 8: 하차지 이용자 (Landowner) 회원가입 및 사토장 등록
# =============================================================
slide8 = prs.slides.add_slide(blank_layout)
add_rectangle(slide8, Inches(0), Inches(0), Inches(13.33), Inches(7.5), WHITE)

# 설명
add_textbox(slide8, Inches(0.8), Inches(0.8), Inches(4.5), Inches(0.8), "07. 하차지 이용자 (지주) 회원가입", 24, True, CHARCOAL)
add_textbox(slide8, Inches(0.8), Inches(1.8), Inches(4.5), Inches(5.0), 
            "■ 하차지 이용자 회원가입 요구사항\n"
            "- 토사를 매립 수용하는 사토장의 지주/운영자 역할\n"
            "- 매립 수용이 법적으로 인가된 사토장 정보 및 허가 서류 검증\n\n"
            "■ 상세 입력 항목 및 서류 명세\n"
            "1. 사토장(하차지) 명칭\n"
            "2. 사토장 토양 지번 및 지목 주소\n"
            "3. 매립 가능 총 수용량 (루베 / 톤)\n"
            "4. 필수 증빙 서류 업로드:\n"
            "   - 개발행위 허가증 또는 사토장 신고 허가서 사본\n"
            "   - 하차비 정산 계좌 통장 사본\n\n"
            "■ UI/UX 설계 요점\n"
            "- 주소 입력을 편리하게 하기 위한 지번 주소 우편번호 API 연동 필드 배치", 
            12, False, CHARCOAL)

# 폰 1: 지주 인적 사항 및 지목 주소
p1_l, p1_t, p1_w, p1_h = add_mini_phone(slide8, Inches(5.8), Inches(0.5), Inches(3.2), Inches(6.5), "지주 가입 - 사토장 정보")
add_textbox(slide8, p1_l + Inches(0.15), p1_t + Inches(0.4), p1_w - Inches(0.3), Inches(0.4), "사토장 정보 입력", 14, True, CHARCOAL)

# 사토장명
add_textbox(slide8, p1_l + Inches(0.15), p1_t + Inches(1.0), p1_w - Inches(0.3), Inches(0.3), "사토장(하차지) 명칭", 9, True, SLATE_GRAY)
add_rounded_rect(slide8, p1_l + Inches(0.15), p1_t + Inches(1.3), p1_w - Inches(0.3), Inches(0.4), WHITE, BORDER_GRAY)
add_textbox(slide8, p1_l + Inches(0.25), p1_t + Inches(1.38), p1_w - Inches(0.5), Inches(0.3), "영종도 남단 사토장", 11, True, CHARCOAL)

# 주소
add_textbox(slide8, p1_l + Inches(0.15), p1_t + Inches(1.9), p1_w - Inches(0.3), Inches(0.3), "사토장 소재지 주소", 9, True, SLATE_GRAY)
add_rounded_rect(slide8, p1_l + Inches(0.15), p1_t + Inches(2.2), p1_w - Inches(0.3), Inches(0.4), WHITE, BORDER_GRAY)
add_textbox(slide8, p1_l + Inches(0.25), p1_t + Inches(2.28), p1_w - Inches(0.5), Inches(0.3), "인천 중구 영종동 산 12-34", 10, True, CHARCOAL)

# 수용량
add_textbox(slide8, p1_l + Inches(0.15), p1_t + Inches(2.8), p1_w - Inches(0.3), Inches(0.3), "최대 매립 수용량 (톤)", 9, True, SLATE_GRAY)
add_rounded_rect(slide8, p1_l + Inches(0.15), p1_t + Inches(3.1), p1_w - Inches(0.3), Inches(0.4), WHITE, BORDER_GRAY)
add_textbox(slide8, p1_l + Inches(0.25), p1_t + Inches(3.18), p1_w - Inches(0.5), Inches(0.3), "50,000 톤", 11, True, CHARCOAL)

btn_next = add_rounded_rect(slide8, p1_l + Inches(0.15), p1_t + Inches(4.8), p1_w - Inches(0.3), Inches(0.45), MUSTARD)
add_textbox(slide8, p1_l + Inches(0.15), p1_t + Inches(4.9), p1_w - Inches(0.3), Inches(0.3), "다음 단계로", 11, True, CHARCOAL, PP_ALIGN.CENTER)


# 폰 2: 지주 서류 제출
p2_l, p2_t, p2_w, p2_h = add_mini_phone(slide8, Inches(9.4), Inches(0.5), Inches(3.2), Inches(6.5), "지주 가입 - 서류 제출")
add_textbox(slide8, p2_l + Inches(0.15), p2_t + Inches(0.4), p2_w - Inches(0.3), Inches(0.4), "사토 허가 서류 제출", 14, True, CHARCOAL)
add_textbox(slide8, p2_l + Inches(0.15), p2_t + Inches(0.75), p2_w - Inches(0.3), Inches(0.4), "매립 및 정산 검증을 위한 필수 증빙을 첨부해 주세요.", 9, False, SLATE_GRAY)

# 서류 2종
add_upload_box(slide8, p2_l + Inches(0.15), p2_t + Inches(1.3), p2_w - Inches(0.3), Inches(1.2), "개발행위(사토장) 허가증 사본")
add_upload_box(slide8, p2_l + Inches(0.15), p2_t + Inches(2.8), p2_w - Inches(0.3), Inches(1.2), "정산 통장 사본 (사토장 운영인 명의)")

btn_submit = add_rounded_rect(slide8, p2_l + Inches(0.15), p2_t + Inches(4.8), p2_w - Inches(0.3), Inches(0.45), MUSTARD)
add_textbox(slide8, p2_l + Inches(0.15), p2_t + Inches(4.9), p2_w - Inches(0.3), Inches(0.3), "제출 및 가입 신청", 11, True, CHARCOAL, PP_ALIGN.CENTER)


# 파워포인트 파일 저장
output_path = "D:\\Projects\\dumpring\\dumpring-platform-backend\\덤프링_로그인_회원가입_상세설계.pptx"
prs.save(output_path)
print(f"PPTX 파일 생성 성공: {output_path}")
