import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# 컬러 팔레트 정의 (라이트 머스터드 옐로우 테마)
WHITE = RGBColor(255, 255, 255)
SAND_BG = RGBColor(249, 248, 246)
MUSTARD = RGBColor(245, 158, 11)     # F59E0B (주요 테마 색상)
AMBER = RGBColor(217, 119, 6)        # D97706 (포인트 색상)
CHARCOAL = RGBColor(31, 41, 55)      # 1F2937 (메인 텍스트)
LIGHT_GRAY = RGBColor(243, 244, 246)  # F3F4F6 (영역 구분 배경)
BORDER_GRAY = RGBColor(229, 231, 235) # E5E7EB (카드 테두리)
SLATE_GRAY = RGBColor(75, 85, 99)    # 4B5563 (서브 텍스트)
RED = RGBColor(239, 68, 68)          # EF4444 (반려/에러/취소/회차)
GREEN = RGBColor(16, 185, 129)       # 10B981 (정산완료/하차완료/정상)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6] # 빈 레이아웃

# 유틸리티 함수: 텍스트 상자 추가
def add_textbox(slide, left, top, width, height, text, font_size=11, bold=False, color=CHARCOAL, align=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = 'Malgun Gothic'
    if align:
        p.alignment = align
    return txBox

# 유틸리티 함수: 직사각형 도형 추가
def add_rectangle(slide, left, top, width, height, fill_color, line_color=None, line_width=1):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

# 유틸리티 함수: 둥근 직사각형 도형 추가
def add_rounded_rect(slide, left, top, width, height, fill_color, line_color=None, line_width=1):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

# 유틸리티 함수: 모바일 폰 프레임 생성
def add_phone_frame(slide, left, top, width, height, title_text="덤프링"):
    bezel = add_rounded_rect(slide, left, top, width, height, SAND_BG, CHARCOAL, 3)
    screen_left = left + Inches(0.15)
    screen_top = top + Inches(0.4)
    screen_width = width - Inches(0.3)
    screen_height = height - Inches(0.6)
    screen = add_rectangle(slide, screen_left, screen_top, screen_width, screen_height, WHITE, BORDER_GRAY, 1)
    
    header = add_rectangle(slide, screen_left, screen_top, screen_width, Inches(0.5), SAND_BG, BORDER_GRAY, 1)
    add_textbox(slide, screen_left + Inches(0.15), screen_top + Inches(0.12), screen_width - Inches(0.3), Inches(0.3), title_text, 13, True, CHARCOAL)
    
    return screen_left, screen_top + Inches(0.5), screen_width, screen_height - Inches(0.5)

# 유틸리티 함수: 문서 제출/사진 촬영 업로드용 가상 버튼 그리기
def add_upload_box(slide, left, top, width, height, label_text):
    add_rounded_rect(slide, left, top, width, height, LIGHT_GRAY, BORDER_GRAY, 1)
    add_rectangle(slide, left + (width/2) - Inches(0.15), top + (height/2) - Inches(0.2), Inches(0.3), Inches(0.05), SLATE_GRAY)
    add_rectangle(slide, left + (width/2) - Inches(0.025), top + (height/2) - Inches(0.3), Inches(0.05), Inches(0.25), SLATE_GRAY)
    add_textbox(slide, left, top + (height/2) + Inches(0.05), width, Inches(0.35), label_text, 8, False, SLATE_GRAY, PP_ALIGN.CENTER)

# 유틸리티 함수: 미니 폰 프레임 생성 (한 슬라이드에 다수 배치용)
def add_mini_phone(slide, left, top, width, height, title_text="덤프링"):
    bezel = add_rounded_rect(slide, left, top, width, height, SAND_BG, CHARCOAL, 2)
    screen_left = left + Inches(0.1)
    screen_top = top + Inches(0.35)
    screen_width = width - Inches(0.2)
    screen_height = height - Inches(0.5)
    screen = add_rectangle(slide, screen_left, screen_top, screen_width, screen_height, WHITE, BORDER_GRAY, 1)
    
    header = add_rectangle(slide, screen_left, screen_top, screen_width, Inches(0.4), SAND_BG, BORDER_GRAY, 1)
    add_textbox(slide, screen_left + Inches(0.1), screen_top + Inches(0.08), screen_width - Inches(0.2), Inches(0.25), title_text, 10, True, CHARCOAL)
    
    return screen_left, screen_top + Inches(0.4), screen_width, screen_height - Inches(0.4)


# =============================================================
# SLIDE 1: 커버 페이지
# =============================================================
slide1 = prs.slides.add_slide(blank_layout)
add_rectangle(slide1, Inches(0), Inches(0), Inches(13.33), Inches(7.5), SAND_BG)
add_rectangle(slide1, Inches(1.5), Inches(2.2), Inches(0.15), Inches(2.5), MUSTARD)
add_textbox(slide1, Inches(2.0), Inches(2.1), Inches(9.5), Inches(1.2), "덤프링 (Dumpring)", 44, True, CHARCOAL)
add_textbox(slide1, Inches(2.0), Inches(3.2), Inches(9.5), Inches(0.8), "플랫폼 전체 역할군 통합 화면 상세설계서", 28, True, SLATE_GRAY)
add_textbox(slide1, Inches(2.0), Inches(4.1), Inches(9.5), Inches(0.5), "최종 통합본 V3 (라이트 머스터드 옐로우 테마)", 16, False, AMBER)
add_textbox(slide1, Inches(2.0), Inches(5.8), Inches(6.0), Inches(0.8), "작성일: 2026. 06. 05\n작성자: 덤프링 플랫폼 기획팀", 12, False, SLATE_GRAY)


# =============================================================
# SLIDE 2: 통합 테이블 오브 콘텐츠 및 시스템 맵
# =============================================================
slide2 = prs.slides.add_slide(blank_layout)
add_rectangle(slide2, Inches(0), Inches(0), Inches(13.33), Inches(7.5), WHITE)
add_textbox(slide2, Inches(0.8), Inches(0.8), Inches(11.7), Inches(0.8), "덤프링 통합 화면설계서 목차 및 흐름 맵", 24, True, CHARCOAL)

# 목차 영역
add_textbox(slide2, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0),
            "■ 목차 일람 (Table of Contents)\n\n"
            "01. [🔑 로그인/가입] 로그인 및 회원가입 통합 흐름 (Slide 3 ~ 4)\n"
            "02. [🚖 덤프 기사] 배차 확인, GPS 미터기 및 실시간 정산 (Slide 5 ~ 6)\n"
            "03. [🚚 덤프 차주] 차량/기사 관제 대시보드 및 기사 관리 (Slide 7 ~ 8)\n"
            "04. [🏢 현장관리자] B2B 오더 발행 및 월별 계산서 청구 (Slide 9 ~ 11)\n"
            "05. [👷 현장담당자] 게이트 트럭 실시간 입출차 통제 (Slide 12 ~ 13)\n"
            "06. [🚜 하차지이용자] 사토 게이지 관리 및 최종 반입 승인 (Slide 14 ~ 15)",
            12, False, CHARCOAL)

# 시스템 맵 다이어그램 묘사
add_textbox(slide2, Inches(6.8), Inches(1.8), Inches(5.7), Inches(0.3), "■ 핵심 B2B 거래 트랜잭션 흐름도", 12, True, SLATE_GRAY)
diagram_rect = add_rounded_rect(slide2, Inches(6.8), Inches(2.2), Inches(5.7), Inches(4.3), LIGHT_GRAY, BORDER_GRAY)
add_textbox(slide2, Inches(7.0), Inches(2.4), Inches(5.3), Inches(4.0),
            "1. 오더 발행 (현장관리자)\n"
            "   ➔ 사토장 매칭 요청\n\n"
            "2. 오더 최종 승인 (하차지 지주)\n"
            "   ➔ 기사단말기에 실시간 배차 콜 개시\n\n"
            "3. 배차 수락 및 운행 (기사)\n"
            "   ➔ 상차 출발 및 GPS 택시미터기 작동\n\n"
            "4. 하차 게이트 도착 및 반입 검사 (지주)\n"
            "   ➔ 토질 육안 검사 및 [반입 승인] 처리\n\n"
            "5. 실시간 정산 및 영수증 발행\n"
            "   ➔ 차주/기사 매출 집계 완료 및 거래처 세금 청구",
            11, False, CHARCOAL)


# =============================================================
# SLIDE 3: 로그인 & 회원가입 - 통합 로그인 & 역할 선택
# =============================================================
slide3 = prs.slides.add_slide(blank_layout)
add_rectangle(slide3, Inches(0), Inches(0), Inches(13.33), Inches(7.5), WHITE)

# 좌측 설명
add_textbox(slide3, Inches(0.8), Inches(0.8), Inches(4.0), Inches(0.8), "01. 로그인 및 가입 역할 선택", 24, True, CHARCOAL)
add_textbox(slide3, Inches(0.8), Inches(1.8), Inches(4.0), Inches(5.0),
            "■ 기획 의도 및 흐름\n"
            "- 휴대폰 번호 입력을 통해 간편인증 및 회원 여부 조회를 원스톱 처리.\n"
            "- 신규 사용자일 경우 5가지 비즈니스 역할군 카드 중 하나를 택하여 가입 프로세스 진입.",
            12, False, CHARCOAL)

# 우측 폰 2대 (인증 / 역할선택)
# 폰 1: 번호 입력
p1_l, p1_t, p1_w, p1_h = add_mini_phone(slide3, Inches(5.2), Inches(0.5), Inches(3.4), Inches(6.5), "본인 인증")
add_textbox(slide3, p1_l + Inches(0.2), p1_t + Inches(0.3), p1_w - Inches(0.4), Inches(0.8), "휴대폰 번호로\n가입/로그인", 14, True, CHARCOAL)
add_rounded_rect(slide3, p1_l + Inches(0.2), p1_t + Inches(1.5), p1_w - Inches(0.4), Inches(0.45), WHITE, BORDER_GRAY)
add_textbox(slide3, p1_l + Inches(0.3), p1_t + Inches(1.6), p1_w - Inches(0.6), Inches(0.3), "010-1234-5678", 12, True, CHARCOAL)
btn1 = add_rounded_rect(slide3, p1_l + Inches(0.2), p1_t + Inches(2.1), p1_w - Inches(0.4), Inches(0.45), MUSTARD)
add_textbox(slide3, p1_l + Inches(0.2), p1_t + Inches(2.2), p1_w - Inches(0.4), Inches(0.3), "인증문자 발송", 11, True, CHARCOAL, PP_ALIGN.CENTER)

# 폰 2: 역할 선택
p2_l, p2_t, p2_w, p2_h = add_mini_phone(slide3, Inches(9.0), Inches(0.5), Inches(3.4), Inches(6.5), "역할 선택")
add_textbox(slide3, p2_l + Inches(0.2), p2_t + Inches(0.2), p2_w - Inches(0.4), Inches(0.5), "가입 유형 선택", 12, True, CHARCOAL)
# 카드 1: 기사 (선택 활성화)
r1 = add_rounded_rect(slide3, p2_l + Inches(0.2), p2_t + Inches(0.8), p2_w - Inches(0.4), Inches(0.7), WHITE, MUSTARD, 2)
add_textbox(slide3, p2_l + Inches(0.35), p2_t + Inches(0.9), p2_w - Inches(0.7), Inches(0.3), "덤프 기사 (Driver)", 11, True, CHARCOAL)
add_textbox(slide3, p2_l + Inches(0.35), p2_t + Inches(1.15), p2_w - Inches(0.7), Inches(0.3), "매칭 오더 수락 및 미터기 주행", 8, False, SLATE_GRAY)
# 카드 2: 차주
r2 = add_rounded_rect(slide3, p2_l + Inches(0.2), p2_t + Inches(1.6), p2_w - Inches(0.4), Inches(0.7), WHITE, BORDER_GRAY, 1)
add_textbox(slide3, p2_l + Inches(0.35), p2_t + Inches(1.7), p2_w - Inches(0.7), Inches(0.3), "덤프 차주 (Owner)", 11, True, CHARCOAL)
# 카드 3: 현장 관리자
r3 = add_rounded_rect(slide3, p2_l + Inches(0.2), p2_t + Inches(2.4), p2_w - Inches(0.4), Inches(0.7), WHITE, BORDER_GRAY, 1)
add_textbox(slide3, p2_l + Inches(0.35), p2_t + Inches(2.5), p2_w - Inches(0.7), Inches(0.3), "현장 관리자 (Manager)", 11, True, CHARCOAL)
# 다음 버튼
next_btn = add_rounded_rect(slide3, p2_l + Inches(0.2), p2_t + Inches(4.8), p2_w - Inches(0.4), Inches(0.5), MUSTARD)
add_textbox(slide3, p2_l + Inches(0.2), p2_t + Inches(4.92), p2_w - Inches(0.4), Inches(0.3), "다음 단계로", 11, True, CHARCOAL, PP_ALIGN.CENTER)


# =============================================================
# SLIDE 4: 로그인 & 회원가입 - 기사 / 차주 필수 서류 제출
# =============================================================
slide4 = prs.slides.add_slide(blank_layout)
add_rectangle(slide4, Inches(0), Inches(0), Inches(13.33), Inches(7.5), WHITE)

# 좌측 설명
add_textbox(slide4, Inches(0.8), Inches(0.8), Inches(4.0), Inches(0.8), "02. 회원가입 상세 서류 제출", 24, True, CHARCOAL)
add_textbox(slide4, Inches(0.8), Inches(1.8), Inches(4.0), Inches(5.0),
            "■ 서류 심사 파이프라인\n"
            "- 운송업 특성상 법적 자격 검토(면허증, 자동차등록증 등) 필수.\n"
            "- 가입 신청 시 심사 상태는 초기 PENDING으로 설정되며, 어드민 승인 시 정식 매칭 권한 부여.",
            12, False, CHARCOAL)

# 우측 폰 2대
# 폰 1: 기사 면허증 업로드
p1_l, p1_t, p1_w, p1_h = add_mini_phone(slide4, Inches(5.2), Inches(0.5), Inches(3.4), Inches(6.5), "기사 가입")
add_textbox(slide4, p1_l + Inches(0.2), p1_t + Inches(0.2), p1_w - Inches(0.4), Inches(0.4), "자격 서류 등록", 12, True, CHARCOAL)
add_upload_box(slide4, p1_l + Inches(0.2), p1_t + Inches(0.7), p1_w - Inches(0.4), Inches(1.6), "운전면허증 사본 촬영/업로드")
add_upload_box(slide4, p1_l + Inches(0.2), p1_t + Inches(2.5), p1_w - Inches(0.4), Inches(1.6), "화물운송 종사 자격증 첨부")
done_btn1 = add_rounded_rect(slide4, p1_l + Inches(0.2), p1_t + Inches(4.8), p1_w - Inches(0.4), Inches(0.5), MUSTARD)
add_textbox(slide4, p1_l + Inches(0.2), p1_t + Inches(4.92), p1_w - Inches(0.4), Inches(0.3), "기사 가입 신청 완료", 11, True, CHARCOAL, PP_ALIGN.CENTER)

# 폰 2: 차주 서류 업로드
p2_l, p2_t, p2_w, p2_h = add_mini_phone(slide4, Inches(9.0), Inches(0.5), Inches(3.4), Inches(6.5), "차주 가입")
add_textbox(slide4, p2_l + Inches(0.2), p2_t + Inches(0.2), p2_w - Inches(0.4), Inches(0.4), "차주 증빙 등록", 12, True, CHARCOAL)
add_upload_box(slide4, p2_l + Inches(0.2), p2_t + Inches(0.7), p2_w - Inches(0.4), Inches(1.6), "사업자등록증 사본 첨부")
add_upload_box(slide4, p2_l + Inches(0.2), p2_t + Inches(2.5), p2_w - Inches(0.4), Inches(1.6), "정산용 통장 사본 첨부")
done_btn2 = add_rounded_rect(slide4, p2_l + Inches(0.2), p2_t + Inches(4.8), p2_w - Inches(0.4), Inches(0.5), MUSTARD)
add_textbox(slide4, p2_l + Inches(0.2), p2_t + Inches(4.92), p2_w - Inches(0.4), Inches(0.3), "차주 가입 신청 완료", 11, True, CHARCOAL, PP_ALIGN.CENTER)


# =============================================================
# SLIDE 5: 기사 - 배차 목록 & 오더 상세 수락
# =============================================================
slide5 = prs.slides.add_slide(blank_layout)
add_rectangle(slide5, Inches(0), Inches(0), Inches(13.33), Inches(7.5), WHITE)

# 좌측 설명
add_textbox(slide5, Inches(0.8), Inches(0.8), Inches(4.0), Inches(0.8), "03. 기사 배차 검색 및 수락", 24, True, CHARCOAL)
add_textbox(slide5, Inches(0.8), Inches(1.8), Inches(4.0), Inches(5.0),
            "■ 기사용 메인 배차보드\n"
            "- 기사가 대기 중인 상태에서 즉각적인 배차 정보 확인 가능.\n"
            "- 출발/도착지 위치, 이동 거리, 거래 단가 노출.\n"
            "- 배차 수락 시 즉시 내비게이션 연동 및 미터기 대기 상태 돌입.",
            12, False, CHARCOAL)

# 우측 폰 2대
# 폰 1: 배차 리스트
p1_l, p1_t, p1_w, p1_h = add_mini_phone(slide5, Inches(5.2), Inches(0.5), Inches(3.4), Inches(6.5), "배차 리스트")
add_textbox(slide5, p1_l + Inches(0.2), p1_t + Inches(0.2), p1_w - Inches(0.4), Inches(0.3), "추천 운송 오더", 11, True, SLATE_GRAY)
# 오더 카드 1
o1 = add_rounded_rect(slide5, p1_l + Inches(0.15), p1_t + Inches(0.6), p1_w - Inches(0.3), Inches(1.1), WHITE, BORDER_GRAY)
add_textbox(slide5, p1_l + Inches(0.25), p1_t + Inches(0.65), p1_w - Inches(0.5), Inches(0.35), "송도 1공구 ➔ 영종 사토장", 11, True, CHARCOAL)
add_textbox(slide5, p1_l + Inches(0.25), p1_t + Inches(0.95), p1_w - Inches(0.5), Inches(0.3), "25톤  |  양질토  |  편도 18km", 9, False, SLATE_GRAY)
add_textbox(slide5, p1_l + Inches(0.25), p1_t + Inches(1.3), p1_w - Inches(0.5), Inches(0.3), "1회 운송 단가: 150,000 원", 10, True, AMBER)

# 오더 카드 2
o2 = add_rounded_rect(slide5, p1_l + Inches(0.15), p1_t + Inches(1.8), p1_w - Inches(0.3), Inches(1.1), WHITE, BORDER_GRAY)
add_textbox(slide5, p1_l + Inches(0.25), p1_t + Inches(1.85), p1_w - Inches(0.5), Inches(0.35), "인천 아라뱃길 ➔ 김포 사토장", 11, True, CHARCOAL)

# 폰 2: 오더 수락/지도
p2_l, p2_t, p2_w, p2_h = add_mini_phone(slide5, Inches(9.0), Inches(0.5), Inches(3.4), Inches(6.5), "오더 상세")
map_rect = add_rectangle(slide5, p2_l + Inches(0.15), p2_t + Inches(0.2), p2_w - Inches(0.3), Inches(1.8), LIGHT_GRAY, BORDER_GRAY)
add_textbox(slide5, p2_l + Inches(0.15), p2_t + Inches(0.9), p2_w - Inches(0.3), Inches(0.4), "[지리 정보 지도 맵 가상 뷰]", 10, False, SLATE_GRAY, PP_ALIGN.CENTER)
# 오더 기본 명세
spec_rect = add_rounded_rect(slide5, p2_l + Inches(0.15), p2_t + Inches(2.1), p2_w - Inches(0.3), Inches(1.3), WHITE, BORDER_GRAY)
add_textbox(slide5, p2_l + Inches(0.25), p2_t + Inches(2.18), p2_w - Inches(0.5), Inches(1.1), 
            "• 상차지: 송도 1공구 아파트 현장\n• 하차지: 영종 남단 사토장\n• 단가: 회당 150,000원\n• 자재: 양질토 (25톤 덤프)", 9, False, CHARCOAL)
ac_btn = add_rounded_rect(slide5, p2_l + Inches(0.2), p2_t + Inches(4.7), p2_w - Inches(0.4), Inches(0.6), MUSTARD)
add_textbox(slide5, p2_l + Inches(0.2), p2_t + Inches(4.85), p2_w - Inches(0.4), Inches(0.3), "배차 오더 수락", 13, True, CHARCOAL, PP_ALIGN.CENTER)


# =============================================================
# SLIDE 6: 기사 - 실시간 GPS 미터기 & 정산 영수증
# =============================================================
slide6 = prs.slides.add_slide(blank_layout)
add_rectangle(slide6, Inches(0), Inches(0), Inches(13.33), Inches(7.5), WHITE)

# 좌측 설명
add_textbox(slide6, Inches(0.8), Inches(0.8), Inches(4.0), Inches(0.8), "04. 실시간 미터기 및 정산 영수증", 24, True, CHARCOAL)
add_textbox(slide6, Inches(0.8), Inches(1.8), Inches(4.0), Inches(5.0),
            "■ GPS 미터기 및 실시간 관제\n"
            "- 주행 시작 시 백그라운드 GPS 트래킹을 통해 기사의 실시간 주행 궤적 및 속도 정보 수집.\n"
            "- 하차지 지오펜싱(200m) 접근 시 반입 심사 트리거가 지주용 화면에 자동 팝업.",
            12, False, CHARCOAL)

# 우측 폰 2대
# 폰 1: GPS 미터기
p1_l, p1_t, p1_w, p1_h = add_mini_phone(slide6, Inches(5.2), Inches(0.5), Inches(3.4), Inches(6.5), "GPS 미터기 운행")
add_textbox(slide6, p1_l + Inches(0.2), p1_t + Inches(0.2), p1_w - Inches(0.4), Inches(0.4), "미터기 주행 중 🟢", 12, True, GREEN)
# 미터기 원형 위젯 묘사
add_rounded_rect(slide6, p1_l + Inches(0.7), p1_t + Inches(0.8), Inches(2.0), Inches(2.0), SAND_BG, MUSTARD, 3)
add_textbox(slide6, p1_l + Inches(0.7), p1_t + Inches(1.2), Inches(2.0), Inches(0.4), "실시간 운임", 10, False, SLATE_GRAY, PP_ALIGN.CENTER)
add_textbox(slide6, p1_l + Inches(0.7), p1_t + Inches(1.6), Inches(2.0), Inches(0.5), "150,000원", 15, True, CHARCOAL, PP_ALIGN.CENTER)
add_textbox(slide6, p1_l + Inches(0.7), p1_t + Inches(2.1), Inches(2.0), Inches(0.3), "주행 거리: 12.4 km", 9, False, SLATE_GRAY, PP_ALIGN.CENTER)

add_textbox(slide6, p1_l + Inches(0.2), p1_t + Inches(3.1), p1_w - Inches(0.4), Inches(0.8), "• 경로: 송도 1공구 ➔ 영종 사토장\n• 하차지 진입 시 승인 요청이 자동 전달됩니다.", 9, False, SLATE_GRAY)
gate_btn = add_rounded_rect(slide6, p1_l + Inches(0.2), p1_t + Inches(4.7), p1_w - Inches(0.4), Inches(0.6), MUSTARD)
add_textbox(slide6, p1_l + Inches(0.2), p1_t + Inches(4.85), p1_w - Inches(0.4), Inches(0.3), "하차지 진입 완료 (대기)", 12, True, CHARCOAL, PP_ALIGN.CENTER)

# 폰 2: 전자 영수증
p2_l, p2_t, p2_w, p2_h = add_mini_phone(slide6, Inches(9.0), Inches(0.5), Inches(3.4), Inches(6.5), "전자 영수증")
add_textbox(slide6, p2_l + Inches(0.2), p2_t + Inches(0.2), p2_w - Inches(0.4), Inches(0.4), "운송 영수증 확인", 13, True, CHARCOAL)
# 영수증 전표 레이아웃
add_rounded_rect(slide6, p2_l + Inches(0.15), p2_t + Inches(0.7), p2_w - Inches(0.3), Inches(3.5), LIGHT_GRAY, BORDER_GRAY)
add_textbox(slide6, p2_l + Inches(0.3), p2_t + Inches(0.9), p2_w - Inches(0.6), Inches(3.0),
            "■ 영수증 번호: #TR-10992\n"
            "--------------------------\n"
            "• 차량번호: 경기 80사 1234\n"
            "• 운수회사: 대박운송\n"
            "• 기 사 명: 김덤프 기사\n"
            "• 상 차 지: 송도 1공구 아파트\n"
            "• 하 차 지: 영종도 사토장\n"
            "• 자재속성: 양질토\n"
            "• 출발시간: 14:20\n"
            "• 완료시간: 15:10\n"
            "--------------------------\n"
            "최종 운송료: 150,000 원 (정산완료 🟢)", 9, False, CHARCOAL)
close_btn = add_rounded_rect(slide6, p2_l + Inches(0.2), p2_t + Inches(4.8), p2_w - Inches(0.4), Inches(0.5), CHARCOAL)
add_textbox(slide6, p2_l + Inches(0.2), p2_t + Inches(4.92), p2_w - Inches(0.4), Inches(0.3), "영수증 닫기", 11, True, WHITE, PP_ALIGN.CENTER)


# =============================================================
# SLIDE 7: 차주 - 실시간 기사-차량 연동 대시보드
# =============================================================
slide7 = prs.slides.add_slide(blank_layout)
add_rectangle(slide7, Inches(0), Inches(0), Inches(13.33), Inches(7.5), WHITE)

# 좌측 설명
add_textbox(slide7, Inches(0.8), Inches(0.8), Inches(5.5), Inches(0.8), "05. 차주용 대시보드 (관제)", 24, True, CHARCOAL)
add_textbox(slide7, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0),
            "■ 주요 기능 및 UI 포인트\n"
            "- 차주는 자기가 보유한 전체 덤프트럭 차량 리스트와 고용한 기사의 매핑 상태를 일괄 제어.\n"
            "- 기사별 오늘 발생 운임 누계액 및 미터기 주행 단계 실시간 정보 노출.",
            12, False, CHARCOAL)

# 우측 폰 프레임 생성 (기본 단일 프레임)
sc_l, sc_t, sc_w, sc_h = add_phone_frame(slide7, Inches(7.5), Inches(0.5), Inches(4.5), Inches(6.5), "차주 관제 대시보드")

# 총 매출 정보 카드
tot_card = add_rounded_rect(slide7, sc_l + Inches(0.15), sc_t + Inches(0.15), sc_w - Inches(0.3), Inches(1.5), WHITE, BORDER_GRAY)
add_textbox(slide7, sc_l + Inches(0.3), sc_t + Inches(0.25), sc_w - Inches(0.6), Inches(0.3), "보유 차량 및 당일 지표", 11, True, CHARCOAL)
add_textbox(slide7, sc_l + Inches(0.3), sc_t + Inches(0.55), sc_w - Inches(0.6), Inches(0.6),
            "• 소속 기사: 총 5명 (4명 주행 중 / 1명 휴무)\n"
            "• 차량 매핑 완료: 3대 / 총 4대", 10, False, SLATE_GRAY)
add_textbox(slide7, sc_l + Inches(0.3), sc_t + Inches(1.15), sc_w - Inches(0.6), Inches(0.4), "정산 대기 금액: 1,800,000 원", 13, True, AMBER)

# 매핑 관제 목록
list_y = sc_t + Inches(1.8)
add_textbox(slide7, sc_l + Inches(0.2), list_y, sc_w - Inches(0.4), Inches(0.3), "차량별 실시간 매핑 정보", 11, True, SLATE_GRAY)

# 아이템 1
i1_y = list_y + Inches(0.3)
add_rounded_rect(slide7, sc_l + Inches(0.15), i1_y, sc_w - Inches(0.3), Inches(0.75), WHITE, BORDER_GRAY)
add_textbox(slide7, sc_l + Inches(0.3), i1_y + Inches(0.1), sc_w - Inches(0.6), Inches(0.35), "경기 80사 1234 (25톤)", 11, True, CHARCOAL)
add_textbox(slide7, sc_l + Inches(0.3), i1_y + Inches(0.4), sc_w - Inches(0.6), Inches(0.3), "매핑 기사: 김덤프 기사  |  미터기 주행 중", 9, False, AMBER)

# 아이템 2
i2_y = list_y + Inches(1.15)
add_rounded_rect(slide7, sc_l + Inches(0.15), i2_y, sc_w - Inches(0.3), Inches(0.75), WHITE, BORDER_GRAY)
add_textbox(slide7, sc_l + Inches(0.3), i2_y + Inches(0.1), sc_w - Inches(0.6), Inches(0.35), "서울 80사 5678 (25톤)", 11, True, CHARCOAL)
add_textbox(slide7, sc_l + Inches(0.3), i2_y + Inches(0.4), sc_w - Inches(0.6), Inches(0.3), "매핑 기사: 이배차 기사  |  상차지 이동 대기", 9, False, SLATE_GRAY)


# =============================================================
# SLIDE 8: 차주 - 기사 초대 관리 & 보유 차량 등록 서류 제출
# =============================================================
slide8 = prs.slides.add_slide(blank_layout)
add_rectangle(slide8, Inches(0), Inches(0), Inches(13.33), Inches(7.5), WHITE)

# 좌측 설명
add_textbox(slide8, Inches(0.8), Inches(0.8), Inches(4.0), Inches(0.8), "06. 기사 초대 및 보유 차량 등록", 24, True, CHARCOAL)
add_textbox(slide8, Inches(0.8), Inches(1.8), Inches(4.0), Inches(5.0),
            "■ 차주 행정 관리\n"
            "1. 신규 기사 초대:\n"
            "   - 전화번호만 입력하면 기사 앱 가입 링크가 자동 전송되며, 가입 시 차주 소속으로 즉시 연결.\n"
            "2. 보유 차량 등록:\n"
            "   - 트럭 등록과 필수 법적 서류(등록증) 업로드 폼 구성.",
            12, False, CHARCOAL)

# 우측 폰 2대
# 폰 1: 기사 초대 목록
p1_l, p1_t, p1_w, p1_h = add_mini_phone(slide8, Inches(5.2), Inches(0.5), Inches(3.4), Inches(6.5), "소속 기사 초대")
add_textbox(slide8, p1_l + Inches(0.2), p1_t + Inches(0.2), p1_w - Inches(0.4), Inches(0.3), "신규 기사 초대", 10, True, SLATE_GRAY)
add_rounded_rect(slide8, p1_l + Inches(0.2), p1_t + Inches(0.5), p1_w - Inches(0.4), Inches(0.4), WHITE, BORDER_GRAY)
add_textbox(slide8, p1_l + Inches(0.3), p1_t + Inches(0.56), p1_w - Inches(0.6), Inches(0.3), "기사 연락처 입력", 10, False, SLATE_GRAY)
inv_btn = add_rounded_rect(slide8, p1_l + Inches(0.2), p1_t + Inches(1.0), p1_w - Inches(0.4), Inches(0.4), MUSTARD)
add_textbox(slide8, p1_l + Inches(0.2), p1_t + Inches(1.06), p1_w - Inches(0.4), Inches(0.3), "초대 문자 전송", 10, True, CHARCOAL, PP_ALIGN.CENTER)

add_textbox(slide8, p1_l + Inches(0.2), p1_t + Inches(1.6), p1_w - Inches(0.4), Inches(0.3), "소속 기사 목록", 10, True, SLATE_GRAY)
k_card = add_rounded_rect(slide8, p1_l + Inches(0.2), p1_t + Inches(1.9), p1_w - Inches(0.4), Inches(0.85), WHITE, BORDER_GRAY)
add_textbox(slide8, p1_l + Inches(0.3), p1_t + Inches(1.95), p1_w - Inches(0.6), Inches(0.3), "김덤프 기사 (승인완료 🟢)", 11, True, CHARCOAL)
dd_rect = add_rounded_rect(slide8, p1_l + Inches(0.3), p1_t + Inches(2.3), p1_w - Inches(0.6), Inches(0.35), LIGHT_GRAY, BORDER_GRAY)
add_textbox(slide8, p1_l + Inches(0.4), p1_t + Inches(2.33), p1_w - Inches(0.8), Inches(0.3), "배정 차량: 경기 80사 1234  ▼", 9, False, CHARCOAL)

# 폰 2: 보유 차량 등록
p2_l, p2_t, p2_w, p2_h = add_mini_phone(slide8, Inches(9.0), Inches(0.5), Inches(3.4), Inches(6.5), "보유 차량 등록")
add_textbox(slide8, p2_l + Inches(0.2), p2_t + Inches(0.2), p2_w - Inches(0.4), Inches(0.3), "차량 등록 번호판", 10, True, SLATE_GRAY)
add_rounded_rect(slide8, p2_l + Inches(0.2), p2_t + Inches(0.45), p2_w - Inches(0.4), Inches(0.4), WHITE, BORDER_GRAY)
add_textbox(slide8, p2_l + Inches(0.3), p2_t + Inches(0.52), p2_w - Inches(0.6), Inches(0.3), "서울 80사 7777", 11, True, CHARCOAL)

add_textbox(slide8, p2_l + Inches(0.2), p2_t + Inches(1.0), p2_w - Inches(0.4), Inches(0.3), "차량 톤수 규격", 10, True, SLATE_GRAY)
add_rounded_rect(slide8, p2_l + Inches(0.2), p2_t + Inches(1.25), Inches(1.3), Inches(0.35), MUSTARD)
add_textbox(slide8, p2_l + Inches(0.2), p2_t + Inches(1.31), Inches(1.3), Inches(0.3), "25톤 덤프", 9, True, CHARCOAL, PP_ALIGN.CENTER)

add_upload_box(slide8, p2_l + Inches(0.2), p2_t + Inches(1.8), p2_w - Inches(0.4), Inches(1.4), "자동차등록증 사본 업로드")
reg_btn = add_rounded_rect(slide8, p2_l + Inches(0.2), p2_t + Inches(4.7), p2_w - Inches(0.4), Inches(0.6), MUSTARD)
add_textbox(slide8, p2_l + Inches(0.2), p2_t + Inches(4.85), p2_w - Inches(0.4), Inches(0.3), "차량 등록 신청", 13, True, CHARCOAL, PP_ALIGN.CENTER)


# =============================================================
# SLIDE 9: 현장관리자 - 웹 종합 관제 콘솔
# =============================================================
slide9 = prs.slides.add_slide(blank_layout)
add_rectangle(slide9, Inches(0), Inches(0), Inches(13.33), Inches(7.5), WHITE)

# 좌측 설명 영역
add_textbox(slide9, Inches(0.8), Inches(0.8), Inches(4.0), Inches(0.8), "07. 현장 관리자 웹 대시보드", 24, True, CHARCOAL)
add_textbox(slide9, Inches(0.8), Inches(1.8), Inches(4.0), Inches(5.0),
            "■ 건설 소장용 웹 서비스\n"
            "- 대화면 모니터링을 상정한 16:9 와이드 비율의 반응형 웹 콘솔 설계.\n"
            "- 금일 운행 트럭, 총 반출량, 누적 운임 등 지표 카드를 최상단에 고정.\n"
            "- 실시간 출입 차량 매핑 정보 및 상태 모니터링 테이블 제공.",
            12, False, CHARCOAL)

# 우측 웹 브라우저 프레임 생성
web_l = Inches(5.0)
web_t = Inches(1.0)
web_w = Inches(7.5)
web_h = Inches(5.2)

add_rounded_rect(slide9, web_l, web_t, web_w, web_h, SAND_BG, BORDER_GRAY, 2)
add_rectangle(slide9, web_l, web_t, web_w, Inches(0.4), LIGHT_GRAY, BORDER_GRAY)
add_textbox(slide9, web_l + Inches(0.2), web_t + Inches(0.08), web_w - Inches(0.4), Inches(0.3), "덤프링 관리자 콘솔 - 대시보드", 11, True, SLATE_GRAY)

# 카드 3종
c_w = Inches(2.1)
c_h = Inches(0.9)
c_y = web_t + Inches(0.6)

# 카드 1
add_rounded_rect(slide9, web_l + Inches(0.3), c_y, c_w, c_h, WHITE, BORDER_GRAY)
add_textbox(slide9, web_l + Inches(0.45), c_y + Inches(0.1), c_w, Inches(0.3), "금일 운행 트럭", 10, False, SLATE_GRAY)
add_textbox(slide9, web_l + Inches(0.45), c_y + Inches(0.35), c_w, Inches(0.4), "12 대", 15, True, MUSTARD)

# 카드 2
add_rounded_rect(slide9, web_l + Inches(2.6), c_y, c_w, c_h, WHITE, BORDER_GRAY)
add_textbox(slide9, web_l + Inches(2.75), c_y + Inches(0.1), c_w, Inches(0.3), "총 반출 완료", 10, False, SLATE_GRAY)
add_textbox(slide9, web_l + Inches(2.75), c_y + Inches(0.35), c_w, Inches(0.4), "48 회", 15, True, AMBER)

# 카드 3
add_rounded_rect(slide9, web_l + Inches(4.9), c_y, c_w, c_h, WHITE, BORDER_GRAY)
add_textbox(slide9, web_l + Inches(5.05), c_y + Inches(0.1), c_w, Inches(0.3), "누적 발생 운임", 10, False, SLATE_GRAY)
add_textbox(slide9, web_l + Inches(5.05), c_y + Inches(0.35), c_w, Inches(0.4), "4,800,000 원", 15, True, CHARCOAL)

# 테이블 뷰
tbl_rect = add_rounded_rect(slide9, web_l + Inches(0.3), web_t + Inches(1.7), web_w - Inches(0.6), Inches(3.1), WHITE, BORDER_GRAY)
add_textbox(slide9, web_l + Inches(0.45), web_t + Inches(1.85), web_w - Inches(0.9), Inches(2.8),
            "차량번호      규격      기사명      상태            상차완료시각\n"
            "----------------------------------------------------------------------\n"
            "경기 80사 1234  25톤      김덤프      하차 완료 🟢     15:10\n"
            "서울 80사 5678  25톤      이배차      미터기 주행 🟠   15:30\n"
            "인천 80사 9012  15톤      박사토      상차 완료 🟢     15:42\n"
            "경기 80사 5555  25톤      최주행      입차 대기 🟡     15:50", 10, False, CHARCOAL)


# =============================================================
# SLIDE 10: 현장관리자 - B2B 매칭 오더 등록 & 담당자 승인 큐
# =============================================================
slide10 = prs.slides.add_slide(blank_layout)
add_rectangle(slide10, Inches(0), Inches(0), Inches(13.33), Inches(7.5), WHITE)

# 좌측 설명
add_textbox(slide10, Inches(0.8), Inches(0.8), Inches(4.0), Inches(0.8), "08. 오더 생성 및 소속 직원 승인", 24, True, CHARCOAL)
add_textbox(slide10, Inches(0.8), Inches(1.8), Inches(4.0), Inches(5.0),
            "■ 오더 등록 및 직원 승인\n"
            "1. B2B 매칭 오더 등록:\n"
            "   - 현장 소장이 반출할 토량 및 매칭될 사토장 단가를 지정하여 일감 오더 등록.\n"
            "2. 가입 직원 승인:\n"
            "   - 현장 직원(Site Worker)이 소속 등록을 요청했을 때, 소장이 가입 진위 여부를 수동으로 최종 승인/반려하는 승인대기 대기 큐.",
            12, False, CHARCOAL)

# 우측 폰 2대
# 폰 1: 오더 등록 폼
p1_l, p1_t, p1_w, p1_h = add_mini_phone(slide10, Inches(5.2), Inches(0.5), Inches(3.4), Inches(6.5), "오더 등록")
add_textbox(slide10, p1_l + Inches(0.2), p1_t + Inches(0.2), p1_w - Inches(0.4), Inches(0.3), "오더 제목", 9, True, SLATE_GRAY)
add_rounded_rect(slide10, p1_l + Inches(0.2), p1_t + Inches(0.45), p1_w - Inches(0.4), Inches(0.4), WHITE, BORDER_GRAY)
add_textbox(slide10, p1_l + Inches(0.3), p1_t + Inches(0.51), p1_w - Inches(0.6), Inches(0.3), "송도 1공구 토사 반출", 10, True, CHARCOAL)

add_textbox(slide10, p1_l + Inches(0.2), p1_t + Inches(1.0), p1_w - Inches(0.4), Inches(0.3), "톤수 및 단가", 9, True, SLATE_GRAY)
add_rounded_rect(slide10, p1_l + Inches(0.2), p1_t + Inches(1.25), p1_w - Inches(0.4), Inches(0.4), WHITE, BORDER_GRAY)
add_textbox(slide10, p1_l + Inches(0.3), p1_t + Inches(1.31), p1_w - Inches(0.6), Inches(0.3), "25톤 덤프  |  150,000 원", 10, True, AMBER)

add_textbox(slide10, p1_l + Inches(0.2), p1_t + Inches(1.8), p1_w - Inches(0.4), Inches(0.3), "반출 토사 필터", 9, True, SLATE_GRAY)
add_rounded_rect(slide10, p1_l + Inches(0.2), p1_t + Inches(2.1), Inches(1.3), Inches(0.35), MUSTARD)
add_textbox(slide10, p1_l + Inches(0.2), p1_t + Inches(2.16), Inches(1.3), Inches(0.3), "양질토", 9, True, CHARCOAL, PP_ALIGN.CENTER)

done_btn1 = add_rounded_rect(slide10, p1_l + Inches(0.2), p1_t + Inches(4.7), p1_w - Inches(0.4), Inches(0.6), MUSTARD)
add_textbox(slide10, p1_l + Inches(0.2), p1_t + Inches(4.85), p1_w - Inches(0.4), Inches(0.3), "오더 공고 발행", 13, True, CHARCOAL, PP_ALIGN.CENTER)

# 폰 2: 직원 승인 큐
p2_l, p2_t, p2_w, p2_h = add_mini_phone(slide10, Inches(9.0), Inches(0.5), Inches(3.4), Inches(6.5), "직원 권한 승인")
add_textbox(slide10, p2_l + Inches(0.2), p2_t + Inches(0.2), p2_w - Inches(0.4), Inches(0.3), "소속 가입 대기 직원 (1)", 10, True, SLATE_GRAY)
p_card = add_rounded_rect(slide10, p2_l + Inches(0.15), p2_t + Inches(0.5), p2_w - Inches(0.3), Inches(1.4), WHITE, BORDER_GRAY)
add_textbox(slide10, p2_l + Inches(0.25), p2_t + Inches(0.6), p2_w - Inches(0.5), Inches(0.5), "홍길동 대리\n연락처: 010-1111-2222", 11, True, CHARCOAL)

ap_b = add_rounded_rect(slide10, p2_l + Inches(0.25), p2_t + Inches(1.15), Inches(1.2), Inches(0.35), MUSTARD)
add_textbox(slide10, p2_l + Inches(0.25), p2_t + Inches(1.2), Inches(1.2), Inches(0.25), "승인", 9, True, CHARCOAL, PP_ALIGN.CENTER)

re_b = add_rounded_rect(slide10, p2_l + Inches(1.55), p2_t + Inches(1.15), Inches(1.2), Inches(0.35), LIGHT_GRAY, BORDER_GRAY)
add_textbox(slide10, p2_l + Inches(1.55), p2_t + Inches(1.2), Inches(1.2), Inches(0.25), "반려", 9, False, SLATE_GRAY, PP_ALIGN.CENTER)


# =============================================================
# SLIDE 11: 현장관리자 - 정산 마감 및 세금계산서 발행
# =============================================================
slide11 = prs.slides.add_slide(blank_layout)
add_rectangle(slide11, Inches(0), Inches(0), Inches(13.33), Inches(7.5), WHITE)

# 좌측 설명 영역
add_textbox(slide11, Inches(0.8), Inches(0.8), Inches(4.0), Inches(0.8), "09. 정산 마감 및 국세청 발행", 24, True, CHARCOAL)
add_textbox(slide11, Inches(0.8), Inches(1.8), Inches(4.0), Inches(5.0),
            "■ 정산 마감 주기 콘솔\n"
            "- 일자별/차량별로 반출된 거래 금액을 월 단위로 일괄 취합하여 마감 처리.\n"
            "- 마감 완료된 내역은 세금계산서 청구 발행 프로세스로 연결.\n"
            "- Excel 세부 주행 상세 전표 내역 다운로드 제공.",
            12, False, CHARCOAL)

# 우측 웹 브라우저 프레임 생성
web_l = Inches(5.0)
web_t = Inches(1.0)
web_w = Inches(7.5)
web_h = Inches(5.2)

add_rounded_rect(slide11, web_l, web_t, web_w, web_h, SAND_BG, BORDER_GRAY, 2)
add_rectangle(slide11, web_l, web_t, web_w, Inches(0.4), LIGHT_GRAY, BORDER_GRAY)
add_textbox(slide11, web_l + Inches(0.2), web_t + Inches(0.08), web_w - Inches(0.4), Inches(0.3), "덤프링 관리자 콘솔 - 월간 정산마감", 11, True, SLATE_GRAY)

# 요약 카드
c_rect = add_rounded_rect(slide11, web_l + Inches(0.3), web_t + Inches(0.6), web_w - Inches(0.6), Inches(1.1), WHITE, BORDER_GRAY)
add_textbox(slide11, web_l + Inches(0.5), web_t + Inches(0.7), web_w - Inches(1.0), Inches(0.3), "2026년 06월 정산 대상 월마감 집계", 12, True, CHARCOAL)
add_textbox(slide11, web_l + Inches(0.5), web_t + Inches(1.05), web_w - Inches(1.0), Inches(0.5),
            "• 총 운송 완료 건수: 240 회  |  • 청구 금액 합산: 36,000,000 원 (VAT 별도)", 11, False, SLATE_GRAY)

# 마감 내역 그리드
g_rect = add_rounded_rect(slide11, web_l + Inches(0.3), web_t + Inches(1.8), web_w - Inches(0.6), Inches(2.2), WHITE, BORDER_GRAY)
add_textbox(slide11, web_l + Inches(0.4), web_t + Inches(1.9), web_w - Inches(0.8), Inches(2.0),
            "일자           티켓번호     차량번호     하차지           금액           마감여부\n"
            "----------------------------------------------------------------------\n"
            "2026.06.05     #1094        12가3456     영종 사토장      150,000원      마감 대기\n"
            "2026.06.05     #1095        78나9012     영종 사토장      150,000원      마감 대기\n"
            "2026.06.04     #1088        34다5678     영종 사토장      150,000원      마감 대기\n"
            "2026.06.03     #1072        90마1234     영종 사토장      150,000원      마감 완료", 9, False, CHARCOAL)

# 액션 단추
b1 = add_rounded_rect(slide11, web_l + Inches(3.2), web_t + Inches(4.3), Inches(1.8), Inches(0.45), MUSTARD)
add_textbox(slide11, web_l + Inches(3.2), web_t + Inches(4.38), Inches(1.8), Inches(0.35), "정산 마감 확정", 11, True, CHARCOAL, PP_ALIGN.CENTER)

b2 = add_rounded_rect(slide11, web_l + Inches(5.1), web_t + Inches(4.3), Inches(2.0), Inches(0.45), CHARCOAL)
add_textbox(slide11, web_l + Inches(5.1), web_t + Inches(4.38), Inches(2.0), Inches(0.35), "세금계산서 청구 요청", 11, True, WHITE, PP_ALIGN.CENTER)


# =============================================================
# SLIDE 12: 현장담당자 - 게이트 실시간 진입 대기 관제
# =============================================================
slide12 = prs.slides.add_slide(blank_layout)
add_rectangle(slide12, Inches(0), Inches(0), Inches(13.33), Inches(7.5), WHITE)

# 좌측 설명
add_textbox(slide12, Inches(0.8), Inches(0.8), Inches(5.5), Inches(0.8), "10. 현장 담당자(직원) 게이트 관제 홈", 24, True, CHARCOAL)
add_textbox(slide12, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0),
            "■ 게이트 상하차 통제 모바일 앱\n"
            "- 현장에 상주하는 게이트 직원이 진입 대기 중인 차량의 번호판을 실시간으로 확인하는 모바일 대시보드.\n"
            "- 지오펜싱(200m) 입차 감지 시 큐에 자동으로 적재되어 직원이 원클릭으로 상차 체크 가능.",
            12, False, CHARCOAL)

# 우측 폰 프레임 생성
sc_l, sc_t, sc_w, sc_h = add_phone_frame(slide12, Inches(7.5), Inches(0.5), Inches(4.5), Inches(6.5), "송도 1공구 - 게이트 관제")

# 3단 메트릭 카드
c_w = (sc_w - Inches(0.5)) / 3
c_h = Inches(0.8)

# 카드 1
add_rounded_rect(slide12, sc_l + Inches(0.15), sc_t + Inches(0.15), c_w, c_h, WHITE, BORDER_GRAY)
add_textbox(slide12, sc_l + Inches(0.15), sc_t + Inches(0.22), c_w, Inches(0.25), "입차대기", 9, False, SLATE_GRAY, PP_ALIGN.CENTER)
add_textbox(slide12, sc_l + Inches(0.15), sc_t + Inches(0.48), c_w, Inches(0.35), "2 대", 13, True, MUSTARD, PP_ALIGN.CENTER)

# 카드 2
add_rounded_rect(slide12, sc_l + Inches(0.15) + c_w + Inches(0.1), sc_t + Inches(0.15), c_w, c_h, WHITE, BORDER_GRAY)
add_textbox(slide12, sc_l + Inches(0.15) + c_w + Inches(0.1), sc_t + Inches(0.22), c_w, Inches(0.25), "상차중", 9, False, SLATE_GRAY, PP_ALIGN.CENTER)
add_textbox(slide12, sc_l + Inches(0.15) + c_w + Inches(0.1), sc_t + Inches(0.48), c_w, Inches(0.35), "3 대", 13, True, AMBER, PP_ALIGN.CENTER)

# 카드 3
add_rounded_rect(slide12, sc_l + Inches(0.15) + (c_w * 2) + Inches(0.2), sc_t + Inches(0.15), c_w, c_h, WHITE, BORDER_GRAY)
add_textbox(slide12, sc_l + Inches(0.15) + (c_w * 2) + Inches(0.2), sc_t + Inches(0.22), c_w, Inches(0.25), "상차완료", 9, False, SLATE_GRAY, PP_ALIGN.CENTER)
add_textbox(slide12, sc_l + Inches(0.15) + (c_w * 2) + Inches(0.2), sc_t + Inches(0.48), c_w, Inches(0.35), "18 대", 13, True, CHARCOAL, PP_ALIGN.CENTER)

# 대기 큐 리스트
list_y = sc_t + Inches(1.1)
add_textbox(slide12, sc_l + Inches(0.2), list_y, sc_w - Inches(0.4), Inches(0.3), "실시간 진입 대기 트럭 (2)", 11, True, SLATE_GRAY)

# 아이템 1
i1_y = list_y + Inches(0.3)
add_rounded_rect(slide12, sc_l + Inches(0.15), i1_y, sc_w - Inches(0.3), Inches(0.8), WHITE, BORDER_GRAY)
add_textbox(slide12, sc_l + Inches(0.3), i1_y + Inches(0.1), sc_w - Inches(0.6), Inches(0.35), "경기 80사 1234 (25톤)", 11, True, CHARCOAL)
add_textbox(slide12, sc_l + Inches(0.3), i1_y + Inches(0.42), sc_w - Inches(0.6), Inches(0.3), "김덤프 기사  |  지오펜싱 입차 (5분 전)", 9, False, GREEN)

# 하단 퀵 체크인 버튼
chk_btn = add_rounded_rect(slide12, sc_l + Inches(0.2), sc_t + Inches(4.7), sc_w - Inches(0.4), Inches(0.6), MUSTARD)
add_textbox(slide12, sc_l + Inches(0.2), sc_t + Inches(4.85), sc_w - Inches(0.4), Inches(0.4), "수동 체크인 / QR 스캔", 14, True, CHARCOAL, PP_ALIGN.CENTER)


# =============================================================
# SLIDE 13: 현장담당자 - 수동 입출차 상차 승인 & 당일 티켓 피드
# =============================================================
slide13 = prs.slides.add_slide(blank_layout)
add_rectangle(slide13, Inches(0), Inches(0), Inches(13.33), Inches(7.5), WHITE)

# 좌측 설명
add_textbox(slide13, Inches(0.8), Inches(0.8), Inches(4.0), Inches(0.8), "11. 수동 입출차 상차 승인", 24, True, CHARCOAL)
add_textbox(slide13, Inches(0.8), Inches(1.8), Inches(4.0), Inches(5.0),
            "■ 수동 개입 흐름\n"
            "- 기사 GPS 신호가 잡히지 않을 경우 게이트 직원이 수동으로 상차 출발 처리를 강제 제어.\n"
            "- 자재 종류 선택 및 상차 게이트 지정을 통해 잘못된 토사 반출 원천 차단.",
            12, False, CHARCOAL)

# 우측 폰 2대
# 폰 1: 수동 상차 승인
p1_l, p1_t, p1_w, p1_h = add_mini_phone(slide13, Inches(5.2), Inches(0.5), Inches(3.4), Inches(6.5), "상차 승인 통제")
add_textbox(slide13, p1_l + Inches(0.2), p1_t + Inches(0.2), p1_w - Inches(0.4), Inches(0.3), "승인 대상 차량", 10, True, SLATE_GRAY)
v_info = add_rounded_rect(slide13, p1_l + Inches(0.15), p1_t + Inches(0.5), p1_w - Inches(0.3), Inches(0.85), WHITE, BORDER_GRAY)
add_textbox(slide13, p1_l + Inches(0.25), p1_t + Inches(0.58), p1_w - Inches(0.5), Inches(0.3), "경기 80사 1234 (25톤)", 11, True, CHARCOAL)

add_textbox(slide13, p1_l + Inches(0.2), p1_t + Inches(1.5), p1_w - Inches(0.4), Inches(0.3), "반출 자재 종류", 10, True, SLATE_GRAY)
add_rounded_rect(slide13, p1_l + Inches(0.2), p1_t + Inches(1.8), Inches(1.4), Inches(0.4), MUSTARD)
add_textbox(slide13, p1_l + Inches(0.2), p1_t + Inches(1.88), Inches(1.4), Inches(0.3), "양질토", 9, True, CHARCOAL, PP_ALIGN.CENTER)

done_btn1 = add_rounded_rect(slide13, p1_l + Inches(0.2), p1_t + Inches(4.7), p1_w - Inches(0.4), Inches(0.6), MUSTARD)
add_textbox(slide13, p1_l + Inches(0.2), p1_t + Inches(4.85), p1_w - Inches(0.4), Inches(0.3), "상차 승인 (출발)", 13, True, CHARCOAL, PP_ALIGN.CENTER)

# 폰 2: 당일 티켓 피드
p2_l, p2_t, p2_w, p2_h = add_mini_phone(slide13, Inches(9.0), Inches(0.5), Inches(3.4), Inches(6.5), "당일 반출 티켓")
add_textbox(slide13, p2_l + Inches(0.2), p2_t + Inches(0.2), p2_w - Inches(0.4), Inches(0.3), "오늘 반출된 티켓 목록", 10, True, SLATE_GRAY)

# 티켓 1
t1_y = p2_t + Inches(0.6)
add_rounded_rect(slide13, p2_l + Inches(0.15), t1_y, p2_w - Inches(0.3), Inches(0.85), WHITE, BORDER_GRAY)
add_textbox(slide13, p2_l + Inches(0.25), t1_y + Inches(0.08), p2_w - Inches(0.5), Inches(0.35), "티켓 #2045  |  경기 80사 1234", 10, True, CHARCOAL)
add_textbox(slide13, p2_l + Inches(0.25), t1_y + Inches(0.43), p2_w - Inches(0.5), Inches(0.3), "양질토  |  15:42 상차완료", 8, False, SLATE_GRAY)
badge_c1 = add_rounded_rect(slide13, p2_l + p2_w - Inches(1.1), t1_y + Inches(0.1), Inches(0.95), Inches(0.28), MUSTARD)
add_textbox(slide13, p2_l + p2_w - Inches(1.1), t1_y + Inches(0.12), Inches(0.95), Inches(0.25), "운행중", 8, True, CHARCOAL, PP_ALIGN.CENTER)

# 티켓 2
t2_y = p2_t + Inches(1.55)
add_rounded_rect(slide13, p2_l + Inches(0.15), t2_y, p2_w - Inches(0.3), Inches(0.85), WHITE, BORDER_GRAY)
add_textbox(slide13, p2_l + Inches(0.25), t2_y + Inches(0.08), p2_w - Inches(0.5), Inches(0.35), "티켓 #2044  |  서울 80사 5678", 10, True, CHARCOAL)
add_textbox(slide13, p2_l + Inches(0.25), t2_y + Inches(0.43), p2_w - Inches(0.5), Inches(0.3), "양질토  |  15:30 상차완료", 8, False, SLATE_GRAY)
badge_c2 = add_rounded_rect(slide13, p2_l + p2_w - Inches(1.1), t2_y + Inches(0.1), Inches(0.95), Inches(0.28), GREEN)
add_textbox(slide13, p2_l + p2_w - Inches(1.1), t2_y + Inches(0.12), Inches(0.95), Inches(0.25), "하차완료", 8, True, WHITE, PP_ALIGN.CENTER)


# =============================================================
# SLIDE 14: 하차지이용자 - 사토장 잔여 수용 가이드 대시보드
# =============================================================
slide14 = prs.slides.add_slide(blank_layout)
add_rectangle(slide14, Inches(0), Inches(0), Inches(13.33), Inches(7.5), WHITE)

# 좌측 설명
add_textbox(slide14, Inches(0.8), Inches(0.8), Inches(5.5), Inches(0.8), "12. 하차지 지주 사토장 대시보드", 24, True, CHARCOAL)
add_textbox(slide14, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0),
            "■ 사토장 지주 모바일 대시보드\n"
            "- 사토장의 매립 잔여 허용 한도를 한눈에 확인할 수 있는 프로그레스 가이드 바 배치.\n"
            "- 하차지 200m 이내로 도달한 진입 예정 덤프트럭 리스트 실시간 자동 관제 피드.",
            12, False, CHARCOAL)

# 우측 폰 프레임 생성
sc_l, sc_t, sc_w, sc_h = add_phone_frame(slide14, Inches(7.5), Inches(0.5), Inches(4.5), Inches(6.5), "사토장 대시보드")

# 사토 수용 한도 게이지 카드
cap_card = add_rounded_rect(slide14, sc_l + Inches(0.15), sc_t + Inches(0.15), sc_w - Inches(0.3), Inches(1.1), WHITE, BORDER_GRAY)
add_textbox(slide14, sc_l + Inches(0.3), sc_t + Inches(0.22), sc_w - Inches(0.6), Inches(0.3), "사토장 매립량 게이지 (65%)", 10, True, CHARCOAL)
add_rounded_rect(slide14, sc_l + Inches(0.3), sc_t + Inches(0.55), sc_w - Inches(0.6), Inches(0.2), LIGHT_GRAY, BORDER_GRAY)
add_rounded_rect(slide14, sc_l + Inches(0.3), sc_t + Inches(0.55), (sc_w - Inches(0.6)) * 0.65, Inches(0.2), MUSTARD)
add_textbox(slide14, sc_l + Inches(0.3), sc_t + Inches(0.8), sc_w - Inches(0.6), Inches(0.35), "6,500㎥ 반입 완료 / 총 10,000㎥ 수용 한도", 9, False, SLATE_GRAY)

# 실시간 매출/반입 집계
sum_card = add_rounded_rect(slide14, sc_l + Inches(0.15), sc_t + Inches(1.4), sc_w - Inches(0.3), Inches(0.85), WHITE, BORDER_GRAY)
add_textbox(slide14, sc_l + Inches(0.3), sc_t + Inches(1.48), sc_w - Inches(0.6), Inches(0.35), "금일 반입 차량:  15 대", 11, True, CHARCOAL)
add_textbox(slide14, sc_l + Inches(0.3), sc_t + Inches(1.85), sc_w - Inches(0.6), Inches(0.35), "누적 정산액: 2,250,000 원", 11, True, AMBER)

# 진입 대기 리스트
add_textbox(slide14, sc_l + Inches(0.2), sc_t + Inches(2.4), sc_w - Inches(0.4), Inches(0.3), "실시간 진입 차량 (200m 내)", 11, True, SLATE_GRAY)
i1 = add_rounded_rect(slide14, sc_l + Inches(0.15), sc_t + Inches(2.7), sc_w - Inches(0.3), Inches(0.75), WHITE, BORDER_GRAY)
add_textbox(slide14, sc_l + Inches(0.3), sc_t + Inches(2.78), sc_w - Inches(0.6), Inches(0.35), "경기 80사 1234 (김덤프 기사)", 11, True, CHARCOAL)
add_textbox(slide14, sc_l + Inches(0.3), sc_t + Inches(3.08), sc_w - Inches(0.6), Inches(0.3), "양질토  |  진입 완료 (50m 전)", 9, False, GREEN)


# =============================================================
# SLIDE 15: 하차지이용자 - 최종 반입 승인 & 육안 토질 검증
# =============================================================
slide15 = prs.slides.add_slide(blank_layout)
add_rectangle(slide15, Inches(0), Inches(0), Inches(13.33), Inches(7.5), WHITE)

# 좌측 설명
add_textbox(slide15, Inches(0.8), Inches(0.8), Inches(4.0), Inches(0.8), "13. 토사 검증 및 최종 반입 승인", 24, True, CHARCOAL)
add_textbox(slide15, Inches(0.8), Inches(1.8), Inches(4.0), Inches(5.0),
            "■ 현장 최종 정산 트리거\n"
            "- 지주가 차량의 실제 하차 상태를 육안 검증한 뒤 [반입 승인] 버튼을 눌러야만 즉시 정산 요금이 이체되고 기사 전자 영수증이 발급됨.\n"
            "- 부적합 토사(뻘흙 등)일 경우 [회차 지시]를 내리고 증빙 사진을 통해 사후 분쟁 조정 센터로 이관 처리.",
            12, False, CHARCOAL)

# 우측 폰 2대
# 폰 1: 반입 승인 및 체크리스트
p1_l, p1_t, p1_w, p1_h = add_mini_phone(slide15, Inches(5.2), Inches(0.5), Inches(3.4), Inches(6.5), "반입 심사")
add_textbox(slide15, p1_l + Inches(0.2), p1_t + Inches(0.2), p1_w - Inches(0.4), Inches(0.3), "진입 차량 정보", 10, True, SLATE_GRAY)
vc = add_rounded_rect(slide15, p1_l + Inches(0.15), p1_t + Inches(0.5), p1_w - Inches(0.3), Inches(0.9), WHITE, BORDER_GRAY)
add_textbox(slide15, p1_l + Inches(0.25), p1_t + Inches(0.58), p1_w - Inches(0.5), Inches(0.7), "경기 80사 1234 (25톤)\n상차지: 송도 1공구 현장", 10, True, CHARCOAL)

add_textbox(slide15, p1_l + Inches(0.2), p1_t + Inches(1.5), p1_w - Inches(0.4), Inches(0.3), "육안 검증 체크", 10, True, SLATE_GRAY)
chk_c = add_rounded_rect(slide15, p1_l + Inches(0.15), p1_t + Inches(1.8), p1_w - Inches(0.3), Inches(0.8), WHITE, BORDER_GRAY)
add_textbox(slide15, p1_l + Inches(0.25), p1_t + Inches(1.88), p1_w - Inches(0.5), Inches(0.6), "☑ 양질의 토사 검토 확인\n☑ 비토사 및 부적합 폐기물 없음", 9, False, CHARCOAL)

ok_b = add_rounded_rect(slide15, p1_l + Inches(0.2), p1_t + Inches(4.3), p1_w - Inches(0.4), Inches(0.45), GREEN)
add_textbox(slide15, p1_l + Inches(0.2), p1_t + Inches(4.4), p1_w - Inches(0.4), Inches(0.3), "최종 반입 승인", 11, True, WHITE, PP_ALIGN.CENTER)

rej_b = add_rounded_rect(slide15, p1_l + Inches(0.2), p1_t + Inches(4.9), p1_w - Inches(0.4), Inches(0.45), RED)
add_textbox(slide15, p1_l + Inches(0.2), p1_t + Inches(5.0), p1_w - Inches(0.4), Inches(0.3), "반입 반려 및 회차 지시", 11, True, WHITE, PP_ALIGN.CENTER)

# 폰 2: 사토 수용 공고 등록
p2_l, p2_t, p2_w, p2_h = add_mini_phone(slide15, Inches(9.0), Inches(0.5), Inches(3.4), Inches(6.5), "사토 공고 등록")
add_textbox(slide15, p2_l + Inches(0.2), p2_t + Inches(0.2), p2_w - Inches(0.4), Inches(0.3), "사토장 비용 정책 설정", 10, True, SLATE_GRAY)
add_rounded_rect(slide15, p2_l + Inches(0.2), p2_t + Inches(0.5), p2_w - Inches(0.4), Inches(0.45), WHITE, BORDER_GRAY)
add_textbox(slide15, p2_l + Inches(0.3), p2_t + Inches(0.58), p2_w - Inches(0.6), Inches(0.3), "회당 단가 수금: 150,000 원", 10, True, AMBER)

add_textbox(slide15, p2_l + Inches(0.2), p2_t + Inches(1.1), p2_w - Inches(0.4), Inches(0.3), "수용 가능한 자재 속성", 10, True, SLATE_GRAY)
add_rounded_rect(slide15, p2_l + Inches(0.2), p2_t + Inches(1.4), Inches(1.3), Inches(0.4), MUSTARD)
add_textbox(slide15, p2_l + Inches(0.2), p2_t + Inches(1.48), Inches(1.3), Inches(0.3), "양질토 수용", 9, True, CHARCOAL, PP_ALIGN.CENTER)

add_rounded_rect(slide15, p2_l + Inches(1.6), p2_t + Inches(1.4), Inches(1.3), Inches(0.4), LIGHT_GRAY, BORDER_GRAY)
add_textbox(slide15, p2_l + Inches(1.6), p2_t + Inches(1.48), Inches(1.3), Inches(0.3), "뻘흙 제한", 9, False, SLATE_GRAY, PP_ALIGN.CENTER)

reg_ann_btn = add_rounded_rect(slide15, p2_l + Inches(0.2), p2_t + Inches(4.7), p2_w - Inches(0.4), Inches(0.6), CHARCOAL)
add_textbox(slide15, p2_l + Inches(0.2), p2_t + Inches(4.85), p2_w - Inches(0.4), Inches(0.3), "사토 수용 공고 배포", 13, True, WHITE, PP_ALIGN.CENTER)


# 파워포인트 파일 저장
output_path = "D:\\Projects\\dumpring\\dumpring-platform-backend\\덤프링_화면설계서_최종_통합본.pptx"
prs.save(output_path)
print(f"PPTX 파일 생성 성공: {output_path}")
