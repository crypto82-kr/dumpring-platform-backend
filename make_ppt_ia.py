import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# 컬러 팔레트 정의 (라이트 머스터드 옐로우 테마)
WHITE = RGBColor(255, 255, 255)
SAND_BG = RGBColor(249, 248, 246)
MUSTARD = RGBColor(245, 158, 11)     # F59E0B (주요 테마 색상)
AMBER = RGBColor(217, 119, 6)        # D97706 (포인트 색상)
CHARCOAL = RGBColor(31, 41, 55)      # 1F2937 (메인 텍스트)
LIGHT_GRAY = RGBColor(243, 244, 246)  # F3F4F6 (영역 구분 배경)
BORDER_GRAY = RGBColor(229, 231, 235) # E5E7EB (카드 테두리)
SLATE_GRAY = RGBColor(75, 85, 99)    # 4B5563 (서브 텍스트)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6] # 빈 레이아웃

# 유틸리티 함수: 텍스트 상자 추가
def add_textbox(slide, left, top, width, height, text, font_size=11, bold=False, color=CHARCOAL, align=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = 'Malgun Gothic'
    if align:
        p.alignment = align
    return txBox

# 유틸리티 함수: 직사각형 도형 추가
def add_rectangle(slide, left, top, width, height, fill_color, line_color=None, line_width=1):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

# =============================================================
# SLIDE 1: 커버 페이지
# =============================================================
slide1 = prs.slides.add_slide(blank_layout)
add_rectangle(slide1, Inches(0), Inches(0), Inches(13.33), Inches(7.5), SAND_BG)
add_rectangle(slide1, Inches(1.5), Inches(2.2), Inches(0.15), Inches(2.5), MUSTARD)
add_textbox(slide1, Inches(2.0), Inches(2.1), Inches(9.5), Inches(1.2), "덤프링 (Dumpring)", 44, True, CHARCOAL)
add_textbox(slide1, Inches(2.0), Inches(3.2), Inches(9.5), Inches(0.8), "모바일 앱 통합 메뉴 구조도 (IA)", 28, True, SLATE_GRAY)
add_textbox(slide1, Inches(2.0), Inches(4.1), Inches(9.5), Inches(0.5), "5대 사용자 역할군별 수직형 흐름 매핑", 16, False, AMBER)
add_textbox(slide1, Inches(2.0), Inches(5.8), Inches(6.0), Inches(0.8), "작성일: 2026. 06. 05\n작성자: 덤프링 플랫폼 기획팀", 12, False, SLATE_GRAY)


# =============================================================
# SLIDE 2: 5분할 세로형 메뉴 구조도 (IA)
# =============================================================
slide2 = prs.slides.add_slide(blank_layout)
add_rectangle(slide2, Inches(0), Inches(0), Inches(13.33), Inches(7.5), WHITE)

# 슬라이드 타이틀
add_textbox(slide2, Inches(0.54), Inches(0.4), Inches(12.25), Inches(0.6), "덤프링 모바일 앱 역할군별 통합 정보 구조도 (Information Architecture)", 20, True, CHARCOAL)

# 5개 컬럼의 가로 좌표 계산
col_w = Inches(2.25)
col_gap = Inches(0.25)
start_x = Inches(0.54)

# 공통 메뉴 데이터 정의
columns_data = [
    {
        "title": "🔑 공통 / 설정",
        "text": "1. 로그인/본인인증\n   - 휴대폰 번호 입력\n   - SMS 6자리 인증\n   - 가입 여부 자동 판정\n\n2. 역할별 회원가입\n   - 기본 정보 등록\n   - 필수 자격증 제출\n\n3. 마이페이지/설정\n   - 내 프로필 관리\n   - 정산 은행 계좌 연동\n   - 푸시 알림 수신 설정"
    },
    {
        "title": "🚖 덤프 기사",
        "text": "1. 배차 검색\n   - 실시간 콜 목록 피드\n   - 톤수/단가 필터링\n   - 지도 상세 및 오더 수락\n\n2. 실시간 운행\n   - GPS 연동 택시 미터기\n   - 게이트 200m 접근 알림\n   - 회차 반려 소명 업로드\n\n3. 정산 및 운행이력\n   - 오늘 전표 내역 피드\n   - 완료 티켓 전자 영수증"
    },
    {
        "title": "🚚 덤프 차주",
        "text": "1. 관제 대시보드\n   - 기사 주행 상태 관제\n   - 오늘 정산대기 누적액\n\n2. 소속 기사 관리\n   - 번호로 SMS 초대 전송\n   - 기사별 차량 매핑 지정\n\n3. 보유 차량 등록\n   - 번호판/톤수 신규 신청\n   - 자동차등록증 첨부함\n\n4. 매출 및 기여도\n   - 월별 기사 매출 순위\n   - 정산서 Excel 내보내기"
    },
    {
        "title": "👷 현장 담당자",
        "text": "1. 게이트 입출 통제\n   - 진입 대기 트럭 큐\n   - 수동 입차 체크인\n   - 바코드/QR코드 스캔\n\n2. 상차 출발 승인\n   - 반출 자재 종류 선택\n   - 출발 승인 (미터기시작)\n\n3. 반출 전표 이력\n   - 오늘 반출 티켓 완료 목록\n   - 게이트 전표 일일 마감"
    },
    {
        "title": "🚜 하차지 지주",
        "text": "1. 사토장 대시보드\n   - 잔여 수용량 게이지\n   - 200m 이내 진입 목록\n\n2. 반입 승인/회차\n   - 육안 토사 상태 검증\n   - 최종 승인 (정산확정)\n   - 반려 및 회차 사진 업로드\n\n3. 공고 및 정산 설정\n   - 단가 정책 수금/지급/무상\n   - 수용 가능 토사 토글\n   - 월말 마감 세금계산서"
    }
]

# 5개 열을 가로로 배치하여 그리기
for i, col in enumerate(columns_data):
    col_x = start_x + (col_w + col_gap) * i
    
    # 1. 헤더 그리기 (머스터드 옐로우)
    add_rectangle(slide2, col_x, Inches(1.2), col_w, Inches(0.45), MUSTARD)
    add_textbox(slide2, col_x, Inches(1.26), col_w, Inches(0.35), col["title"], 11, True, CHARCOAL, PP_ALIGN.CENTER)
    
    # 2. 본문 박스 그리기 (하단 흰색 배경 및 연한 외곽선)
    add_rectangle(slide2, col_x, Inches(1.7), col_w, Inches(5.1), WHITE, BORDER_GRAY, 1)
    
    # 3. 본문 텍스트 기입 (글씨 크기 9.5로 모바일 화면처럼 컴팩트하게 노출)
    add_textbox(slide2, col_x + Inches(0.1), Inches(1.8), col_w - Inches(0.2), Inches(4.9), col["text"], 9.5, False, CHARCOAL)


# 파워포인트 파일 저장
output_path = "D:\\Projects\\dumpring\\dumpring-platform-backend\\덤프링_메뉴구조도.pptx"
prs.save(output_path)
print(f"PPTX 파일 생성 성공: {output_path}")
