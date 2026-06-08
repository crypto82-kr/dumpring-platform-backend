import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# 컬러 팔레트 정의 (라이트 머스터드 옐로우 테마)
WHITE = RGBColor(255, 255, 255)
SAND_BG = RGBColor(249, 248, 246)
MUSTARD = RGBColor(245, 158, 11)     # F59E0B (주요 테마 색상)
AMBER = RGBColor(217, 119, 6)        # D97706 (포인트 색상)
CHARCOAL = RGBColor(31, 41, 55)      # 1F2937 (메인 텍스트)
LIGHT_GRAY = RGBColor(243, 244, 246)  # F3F4F6 (영역 구분 배경)
BORDER_GRAY = RGBColor(229, 231, 235) # E5E7EB (카드 테두리)
SLATE_GRAY = RGBColor(75, 85, 99)    # 4B5563 (서브 텍스트)
RED = RGBColor(239, 68, 68)          # EF4444 (반려/에러/취소)
GREEN = RGBColor(16, 185, 129)       # 10B981 (정산완료/정상)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6] # 빈 레이아웃

# 유틸리티 함수: 텍스트 상자 추가
def add_textbox(slide, left, top, width, height, text, font_size=11, bold=False, color=CHARCOAL, align=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = 'Malgun Gothic'
    if align:
        p.alignment = align
    return txBox

# 유틸리티 함수: 직사각형 도형 추가
def add_rectangle(slide, left, top, width, height, fill_color, line_color=None, line_width=1):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

# 유틸리티 함수: 둥근 직사각형 도형 추가
def add_rounded_rect(slide, left, top, width, height, fill_color, line_color=None, line_width=1):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

# 유틸리티 함수: 모바일 폰 프레임 생성
def add_phone_frame(slide, left, top, width, height, title_text="덤프링"):
    bezel = add_rounded_rect(slide, left, top, width, height, SAND_BG, CHARCOAL, 3)
    screen_left = left + Inches(0.15)
    screen_top = top + Inches(0.4)
    screen_width = width - Inches(0.3)
    screen_height = height - Inches(0.6)
    screen = add_rectangle(slide, screen_left, screen_top, screen_width, screen_height, WHITE, BORDER_GRAY, 1)
    
    # 상단 헤더
    header = add_rectangle(slide, screen_left, screen_top, screen_width, Inches(0.5), SAND_BG, BORDER_GRAY, 1)
    add_textbox(slide, screen_left + Inches(0.15), screen_top + Inches(0.12), screen_width - Inches(0.3), Inches(0.3), title_text, 13, True, CHARCOAL)
    
    return screen_left, screen_top + Inches(0.5), screen_width, screen_height - Inches(0.5)

# 유틸리티 함수: 문서 제출/사진 촬영 업로드용 가상 버튼 그리기
def add_upload_box(slide, left, top, width, height, label_text):
    add_rounded_rect(slide, left, top, width, height, LIGHT_GRAY, BORDER_GRAY, 1)
    add_rectangle(slide, left + (width/2) - Inches(0.15), top + (height/2) - Inches(0.2), Inches(0.3), Inches(0.05), SLATE_GRAY)
    add_rectangle(slide, left + (width/2) - Inches(0.025), top + (height/2) - Inches(0.3), Inches(0.05), Inches(0.25), SLATE_GRAY)
    add_textbox(slide, left, top + (height/2) + Inches(0.05), width, Inches(0.35), label_text, 8, False, SLATE_GRAY, PP_ALIGN.CENTER)


# =============================================================
# SLIDE 1: 커버 페이지 (Cover Page)
# =============================================================
slide1 = prs.slides.add_slide(blank_layout)
add_rectangle(slide1, Inches(0), Inches(0), Inches(13.33), Inches(7.5), SAND_BG)
add_rectangle(slide1, Inches(1.5), Inches(2.2), Inches(0.15), Inches(2.5), MUSTARD)
add_textbox(slide1, Inches(2.0), Inches(2.1), Inches(9.5), Inches(1.2), "덤프링 (Dumpring)", 44, True, CHARCOAL)
add_textbox(slide1, Inches(2.0), Inches(3.2), Inches(9.5), Inches(0.8), "03. 현장 관리자 (소장) 업무 화면 상세설계서", 28, True, SLATE_GRAY)
add_textbox(slide1, Inches(2.0), Inches(4.1), Inches(9.5), Inches(0.5), "최종 테마: 라이트 머스터드 옐로우 (Light Mustard Yellow)", 16, False, AMBER)
add_textbox(slide1, Inches(2.0), Inches(5.8), Inches(6.0), Inches(0.8), "작성일: 2026. 06. 05\n작성자: 덤프링 플랫폼 기획팀", 12, False, SLATE_GRAY)


# =============================================================
# SLIDE 2: Screen 1. 현장 종합 관제 대시보드 (Web Console)
# =============================================================
slide2 = prs.slides.add_slide(blank_layout)
add_rectangle(slide2, Inches(0), Inches(0), Inches(13.33), Inches(7.5), WHITE)

# 좌측 설명 영역
add_textbox(slide2, Inches(0.8), Inches(0.8), Inches(4.0), Inches(0.8), "01. 현장 종합 관제 대시보드", 24, True, CHARCOAL)
add_textbox(slide2, Inches(0.8), Inches(1.8), Inches(4.0), Inches(5.0), 
            "■ 주요 기능\n"
            "- 현장 투입 차량 관리 및 일일 운송 현황을 실시간 분석하는 메인 웹 콘솔 화면\n"
            "- 금일 요약 지표 카드를 상단에 배치하여 빠른 의사결정 지원\n"
            "- 시간대별 주행 물량을 나타내는 통계 차트 영역 탑재\n\n"
            "■ UI 구성 요소\n"
            "1. 통계 지표 카드: 운행 트럭(12대), 총 운송(48회), 누적 발생 운임(4.8백만원) 합계 표시\n"
            "2. 그래프 위젯: 당일 주행 물량 변동 차트\n"
            "3. 관제 테이블: 현재 운행 중인 전체 트럭 번호판, 기사명 및 실시간 운송 상태 명세표", 
            12, False, CHARCOAL)

# 우측 웹 브라우저 프레임 생성
web_l = Inches(5.2)
web_t = Inches(1.0)
web_w = Inches(7.3)
web_h = Inches(5.2)

# 브라우저 백그라운드 & 헤더
add_rounded_rect(slide2, web_l, web_t, web_w, web_h, SAND_BG, BORDER_GRAY, 2)
add_rectangle(slide2, web_l, web_t, web_w, Inches(0.4), LIGHT_GRAY, BORDER_GRAY)
add_textbox(slide2, web_l + Inches(0.2), web_t + Inches(0.08), web_w - Inches(0.4), Inches(0.3), "덤프링 관리자 콘솔 - 대시보드", 11, True, SLATE_GRAY)

# 통계 카드 3종
card_w = Inches(2.1)
card_h = Inches(0.9)
card_y = web_t + Inches(0.6)

# 카드 1 (트럭 대수)
add_rounded_rect(slide2, web_l + Inches(0.3), card_y, card_w, card_h, WHITE, BORDER_GRAY)
add_textbox(slide2, web_l + Inches(0.45), card_y + Inches(0.1), card_w, Inches(0.3), "운행 트럭 수", 10, False, SLATE_GRAY)
add_textbox(slide2, web_l + Inches(0.45), card_y + Inches(0.35), card_w, Inches(0.4), "12 대", 16, True, MUSTARD)

# 카드 2 (운송 횟수)
add_rounded_rect(slide2, web_l + Inches(2.6), card_y, card_w, card_h, WHITE, BORDER_GRAY)
add_textbox(slide2, web_l + Inches(2.75), card_y + Inches(0.1), card_w, Inches(0.3), "총 운송 횟수", 10, False, SLATE_GRAY)
add_textbox(slide2, web_l + Inches(2.75), card_y + Inches(0.35), card_w, Inches(0.4), "48 회", 16, True, AMBER)

# 카드 3 (누적 정산액)
add_rounded_rect(slide2, web_l + Inches(4.9), card_y, card_w, card_h, WHITE, BORDER_GRAY)
add_textbox(slide2, web_l + Inches(5.05), card_y + Inches(0.1), card_w, Inches(0.3), "누적 발생 운임", 10, False, SLATE_GRAY)
add_textbox(slide2, web_l + Inches(5.05), card_y + Inches(0.35), card_w, Inches(0.4), "4,800,000 원", 16, True, CHARCOAL)

# 가상 차트 박스
chart_rect = add_rounded_rect(slide2, web_l + Inches(0.3), web_t + Inches(1.7), web_w - Inches(0.6), Inches(1.6), WHITE, BORDER_GRAY)
add_textbox(slide2, web_l + Inches(0.5), web_t + Inches(1.8), web_w - Inches(1.0), Inches(0.3), "시간대별 배차 물량 통계 (차트 영역)", 11, True, CHARCOAL)
add_rectangle(slide2, web_l + Inches(1.5), web_t + Inches(2.5), Inches(4.5), Inches(0.08), MUSTARD)

# 하단 데이터 리스트 테이블
table_rect = add_rounded_rect(slide2, web_l + Inches(0.3), web_t + Inches(3.5), web_w - Inches(0.6), Inches(1.4), WHITE, BORDER_GRAY)
add_textbox(slide2, web_l + Inches(0.4), web_t + Inches(3.6), web_w - Inches(0.8), Inches(1.2), 
            "차량번호      규격      기사명      상태\n"
            "12가3456      25톤      김덤프      하차 완료\n"
            "78나9012      25톤      이배차      상차 대기\n"
            "34다5678      15톤      박사토      운행 중", 10, False, CHARCOAL)


# =============================================================
# SLIDE 3: Screen 2. B2B 매칭 오더 발행 화면 (오더 등록)
# =============================================================
slide3 = prs.slides.add_slide(blank_layout)
add_rectangle(slide3, Inches(0), Inches(0), Inches(13.33), Inches(7.5), WHITE)

# 좌측 설명 영역
add_textbox(slide3, Inches(0.8), Inches(0.8), Inches(5.5), Inches(0.8), "02. B2B 매칭 오더 발행 화면", 24, True, CHARCOAL)
add_textbox(slide3, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0), 
            "■ 주요 기능\n"
            "- 현장 내 반출할 토사가 생겼을 때 배차 모집을 위해 오더를 기안 및 공포하는 기능\n"
            "- 오더 발행 시 연동된 하차지(사토장) 지주에게 1차 사전 승인 요청 발송\n"
            "- 지주가 승인을 완료하면, 기사 홈 '배차 검색 리스트'에 실시간 노출 개시\n\n"
            "■ 주요 UI 구성\n"
            "1. 기본 필드: 오더 명칭 및 모집할 목표 차량 대수 기입\n"
            "2. 규격/자재 토글: 덤프트럭 톤수(15T/25T) 및 토사 속성(양질토/뻘흙 등) 선택\n"
            "3. 거래 조건: 연동할 사토장 검색/선택 및 1회당 정산 요금(단가) 입력", 
            12, False, CHARCOAL)

# 우측 폰 프레임 생성
sc_l, sc_t, sc_w, sc_h = add_phone_frame(slide3, Inches(7.5), Inches(0.5), Inches(4.5), Inches(6.5), "오더 등록")

# 제목 필드
add_textbox(slide3, sc_l + Inches(0.2), sc_t + Inches(0.2), sc_w - Inches(0.4), Inches(0.3), "오더 제목", 9, True, SLATE_GRAY)
add_rounded_rect(slide3, sc_l + Inches(0.2), sc_t + Inches(0.45), sc_w - Inches(0.4), Inches(0.4), WHITE, BORDER_GRAY)
add_textbox(slide3, sc_l + Inches(0.3), sc_t + Inches(0.52), sc_w - Inches(0.6), Inches(0.3), "송도 1공구 아파트 토사 반출", 11, True, CHARCOAL)

# 대수 및 차량
add_textbox(slide3, sc_l + Inches(0.2), sc_t + Inches(1.0), sc_w - Inches(0.4), Inches(0.3), "모집 차량 대수 및 규격", 9, True, SLATE_GRAY)
add_rounded_rect(slide3, sc_l + Inches(0.2), sc_t + Inches(1.25), sc_w - Inches(0.4), Inches(0.4), WHITE, BORDER_GRAY)
add_textbox(slide3, sc_l + Inches(0.3), sc_t + Inches(1.32), sc_w - Inches(0.6), Inches(0.3), "25톤 덤프 10 대 모집", 11, True, CHARCOAL)

# 토사 선택 (토글 뱃지)
add_textbox(slide3, sc_l + Inches(0.2), sc_t + Inches(1.8), sc_w - Inches(0.4), Inches(0.3), "반출 토사 속성", 9, True, SLATE_GRAY)
add_rounded_rect(slide3, sc_l + Inches(0.2), sc_t + Inches(2.1), Inches(1.1), Inches(0.4), MUSTARD)
add_textbox(slide3, sc_l + Inches(0.2), sc_t + Inches(2.18), Inches(1.1), Inches(0.3), "양질토", 10, True, CHARCOAL, PP_ALIGN.CENTER)

add_rounded_rect(slide3, sc_l + Inches(1.4), sc_t + Inches(2.1), Inches(1.1), Inches(0.4), WHITE, BORDER_GRAY)
add_textbox(slide3, sc_l + Inches(1.4), sc_t + Inches(2.18), Inches(1.1), Inches(0.3), "뻘흙", 10, False, SLATE_GRAY, PP_ALIGN.CENTER)

# 사토장 연동
add_textbox(slide3, sc_l + Inches(0.2), sc_t + Inches(2.7), sc_w - Inches(0.4), Inches(0.3), "연동 하차지(사토장)", 9, True, SLATE_GRAY)
add_rounded_rect(slide3, sc_l + Inches(0.2), sc_t + Inches(2.95), sc_w - Inches(0.4), Inches(0.4), WHITE, BORDER_GRAY)
add_textbox(slide3, sc_l + Inches(0.3), sc_t + Inches(3.02), sc_w - Inches(0.6), Inches(0.3), "영종도 남단 사토 매립지", 11, True, CHARCOAL)

# 단가 설정
add_textbox(slide3, sc_l + Inches(0.2), sc_t + Inches(3.5), sc_w - Inches(0.4), Inches(0.3), "1회당 운송 요금 (단가)", 9, True, SLATE_GRAY)
add_rounded_rect(slide3, sc_l + Inches(0.2), sc_t + Inches(3.75), sc_w - Inches(0.4), Inches(0.4), WHITE, BORDER_GRAY)
add_textbox(slide3, sc_l + Inches(0.3), sc_t + Inches(3.82), sc_w - Inches(0.6), Inches(0.3), "150,000 원", 12, True, AMBER)

# 발행 버튼
post_btn = add_rounded_rect(slide3, sc_l + Inches(0.2), sc_t + Inches(4.7), sc_w - Inches(0.4), Inches(0.6), MUSTARD)
add_textbox(slide3, sc_l + Inches(0.2), sc_t + Inches(4.85), sc_w - Inches(0.4), Inches(0.4), "B2B 오더 발행 신청", 14, True, CHARCOAL, PP_ALIGN.CENTER)


# =============================================================
# SLIDE 4: Screen 3. 소속 현장 및 서류 관리 화면
# =============================================================
slide4 = prs.slides.add_slide(blank_layout)
add_rectangle(slide4, Inches(0), Inches(0), Inches(13.33), Inches(7.5), WHITE)

# 좌측 설명 영역
add_textbox(slide4, Inches(0.8), Inches(0.8), Inches(5.5), Inches(0.8), "03. 공사현장 마스터 정보 관리", 24, True, CHARCOAL)
add_textbox(slide4, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0), 
            "■ 주요 기능\n"
            "- 현장 정보 검증 및 지오펜싱(차량 자동 출입 감지) 영역 설정 화면\n"
            "- 소장이 등록한 지오펜싱 좌표 반경 안에 트럭이 들어오면 기사 앱과 연동되어 출입 자동 확인\n"
            "- 현장 안전 증빙 서류의 심사 이력 관리 기능 탑재\n\n"
            "■ 주요 UI 구성\n"
            "1. 현장 개요: 건설사명, 현장명, 공사 총 기간 등 출력\n"
            "2. 지오펜싱 슬라이더: 차량 입차 인식을 위한 반경(m) 설정 바\n"
            "3. 증빙 서류함: 건설업 면허 및 계약 서류 업로드 썸네일 박스", 
            12, False, CHARCOAL)

# 우측 폰 프레임 생성
sc_l, sc_t, sc_w, sc_h = add_phone_frame(slide4, Inches(7.5), Inches(0.5), Inches(4.5), Inches(6.5), "현장 정보 관리")

# 현장명 정보
add_textbox(slide4, sc_l + Inches(0.2), sc_t + Inches(0.2), sc_w - Inches(0.4), Inches(0.3), "소속 현장 개요", 10, True, SLATE_GRAY)
info_card = add_rounded_rect(slide4, sc_l + Inches(0.15), sc_t + Inches(0.5), sc_w - Inches(0.3), Inches(1.1), WHITE, BORDER_GRAY)
add_textbox(slide4, sc_l + Inches(0.3), sc_t + Inches(0.6), sc_w - Inches(0.6), Inches(0.9), 
            "• 건설사: 대우건설 (주)\n• 현장명: 송도 1공구 아파트 신축공사 현장\n• 공사 기간: 2026.06.01 ~ 2028.05.31", 11, False, CHARCOAL)

# 지오펜싱 반경 설정 슬라이더
add_textbox(slide4, sc_l + Inches(0.2), sc_t + Inches(1.8), sc_w - Inches(0.4), Inches(0.3), "자동 입차 지오펜싱 반경 설정: 200m", 10, True, SLATE_GRAY)
# 슬라이더 라인
add_rectangle(slide4, sc_l + Inches(0.2), sc_t + Inches(2.2), sc_w - Inches(0.4), Inches(0.06), BORDER_GRAY)
# 슬라이더 조절자 (골드 핸들)
add_rounded_rect(slide4, sc_l + Inches(1.5), sc_t + Inches(2.1), Inches(0.25), Inches(0.25), MUSTARD)

# 서류 업로드 영역
add_textbox(slide4, sc_l + Inches(0.2), sc_t + Inches(2.6), sc_w - Inches(0.4), Inches(0.3), "현장 증빙 서류 보관함", 10, True, SLATE_GRAY)
add_upload_box(slide4, sc_l + Inches(0.15), sc_t + Inches(2.9), sc_w - Inches(0.3), Inches(1.3), "건설업 등록증 사본 첨부 (승인완료)")

# 수정 완료 버튼
save_btn = add_rounded_rect(slide4, sc_l + Inches(0.2), sc_t + Inches(4.7), sc_w - Inches(0.4), Inches(0.6), CHARCOAL)
add_textbox(slide4, sc_l + Inches(0.2), sc_t + Inches(4.85), sc_w - Inches(0.4), Inches(0.4), "현장 정보 수정 완료", 14, True, WHITE, PP_ALIGN.CENTER)


# =============================================================
# SLIDE 5: Screen 4. 현장 담당자(직원) 승인 및 관리 화면
# =============================================================
slide5 = prs.slides.add_slide(blank_layout)
add_rectangle(slide5, Inches(0), Inches(0), Inches(13.33), Inches(7.5), WHITE)

# 좌측 설명 영역
add_textbox(slide5, Inches(0.8), Inches(0.8), Inches(5.5), Inches(0.8), "04. 현장 직원 권한 관리 및 승인", 24, True, CHARCOAL)
add_textbox(slide5, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0), 
            "■ 주요 기능\n"
            "- 현장 게이트 통제 권한을 획득하고자 가입 신청한 직원을 수동 승인/반려하는 화면\n"
            "- 소장(관리자)의 승인을 득한 직원만 본 현장의 차량 출입 체크인/아웃 권한 활성화\n\n"
            "■ 주요 UI 구성\n"
            "1. 승인 대기 목록 (Pending Queue):\n"
            "   - 가입한 직원의 실명, 연락처, 기입 직급 노출\n"
            "   - 즉시 판단할 수 있는 [소속 승인] 및 [반려] 제어 버튼 제공\n"
            "2. 등록 직원 리스트:\n"
            "   - 이미 매핑되어 실시간 게이트 근무 중인 소속 직원 일람표", 
            12, False, CHARCOAL)

# 우측 폰 프레임 생성
sc_l, sc_t, sc_w, sc_h = add_phone_frame(slide5, Inches(7.5), Inches(0.5), Inches(4.5), Inches(6.5), "현장 직원 권한 관리")

# 대기 큐 타이틀
add_textbox(slide5, sc_l + Inches(0.2), sc_t + Inches(0.2), sc_w - Inches(0.4), Inches(0.3), "승인 대기 직원", 11, True, SLATE_GRAY)

# 대기 직원 카드
pending_card = add_rounded_rect(slide5, sc_l + Inches(0.15), sc_t + Inches(0.5), sc_w - Inches(0.3), Inches(1.3), WHITE, BORDER_GRAY)
add_textbox(slide5, sc_l + Inches(0.3), sc_t + Inches(0.62), sc_w - Inches(0.6), Inches(0.35), "홍길동 대리 (연락처: 010-1111-2222)", 11, True, CHARCOAL)
add_textbox(slide5, sc_l + Inches(0.3), sc_t + Inches(0.9), sc_w - Inches(0.6), Inches(0.35), "가입 신청일: 2026.06.05  |  역할: 현장담당자", 9, False, SLATE_GRAY)

# 승인 / 반려 버튼 세트 (카드 내 배치)
app_btn = add_rounded_rect(slide5, sc_l + Inches(0.3), sc_t + Inches(1.3), Inches(1.6), Inches(0.38), MUSTARD)
add_textbox(slide5, sc_l + Inches(0.3), sc_t + Inches(1.36), Inches(1.6), Inches(0.3), "소속 승인", 10, True, CHARCOAL, PP_ALIGN.CENTER)

rej_btn = add_rounded_rect(slide5, sc_l + Inches(2.0), sc_t + Inches(1.3), Inches(1.6), Inches(0.38), LIGHT_GRAY, BORDER_GRAY)
add_textbox(slide5, sc_l + Inches(2.0), sc_t + Inches(1.36), Inches(1.6), Inches(0.3), "반려", 10, False, SLATE_GRAY, PP_ALIGN.CENTER)

# 소속 직원 목록
add_textbox(slide5, sc_l + Inches(0.2), sc_t + Inches(2.0), sc_w - Inches(0.4), Inches(0.3), "등록된 소속 직원 목록", 11, True, SLATE_GRAY)

worker1 = add_rounded_rect(slide5, sc_l + Inches(0.15), sc_t + Inches(2.3), sc_w - Inches(0.3), Inches(0.65), WHITE, BORDER_GRAY)
add_textbox(slide5, sc_l + Inches(0.3), sc_t + Inches(2.45), sc_w - Inches(0.6), Inches(0.35), "이순신 과장  |  상차 게이트 담당 (활성)", 11, False, CHARCOAL)

worker2 = add_rounded_rect(slide5, sc_l + Inches(0.15), sc_t + Inches(3.05), sc_w - Inches(0.3), Inches(0.65), WHITE, BORDER_GRAY)
add_textbox(slide5, sc_l + Inches(0.3), sc_t + Inches(3.2), sc_w - Inches(0.6), Inches(0.35), "임꺽정 사원  |  미배정 (휴무)", 11, False, SLATE_GRAY)


# =============================================================
# SLIDE 6: Screen 5. 현장 정산 / 청구서 발행 콘솔 (Web Console)
# =============================================================
slide6 = prs.slides.add_slide(blank_layout)
add_rectangle(slide6, Inches(0), Inches(0), Inches(13.33), Inches(7.5), WHITE)

# 좌측 설명 영역
add_textbox(slide6, Inches(0.8), Inches(0.8), Inches(4.0), Inches(0.8), "05. 정산 마감 및 세금계산서 청구", 24, True, CHARCOAL)
add_textbox(slide6, Inches(0.8), Inches(1.8), Inches(4.0), Inches(5.0), 
            "■ 주요 기능\n"
            "- 완료된 운행 내역에 대해 마감 처리를 진행하고 세금계산서 발행을 요청하는 청구 콘솔\n"
            "- 일자별/차량별 상세 주행 전표 합산 내역을 한눈에 조회\n"
            "- 마감 확정 완료 시 국세청 연동 세금계산서 정산 발행 프로세스 트리거\n\n"
            "■ 주요 UI 구성\n"
            "1. 마감 지표: 운송 마감 대상 월, 총 운송 횟수, 마감 합산 세금 청구 금액 요약\n"
            "2. 하단 상세 그리드: 마감 대상 세부 티켓 리스트\n"
            "3. 마감 처리 전용 우측 정산 액션 버튼 세트", 
            12, False, CHARCOAL)

# 우측 웹 브라우저 프레임 생성
web_l = Inches(5.2)
web_t = Inches(1.0)
web_w = Inches(7.3)
web_h = Inches(5.2)

# 브라우저 백그라운드 & 헤더
add_rounded_rect(slide6, web_l, web_t, web_w, web_h, SAND_BG, BORDER_GRAY, 2)
add_rectangle(slide6, web_l, web_t, web_w, Inches(0.4), LIGHT_GRAY, BORDER_GRAY)
add_textbox(slide6, web_l + Inches(0.2), web_t + Inches(0.08), web_w - Inches(0.4), Inches(0.3), "덤프링 관리자 콘솔 - 정산/마감 실적", 11, True, SLATE_GRAY)

# 정산 개요 카드
summary_rect = add_rounded_rect(slide6, web_l + Inches(0.3), web_t + Inches(0.6), web_w - Inches(0.6), Inches(1.1), WHITE, BORDER_GRAY)
add_textbox(slide6, web_l + Inches(0.5), web_t + Inches(0.72), web_w - Inches(1.0), Inches(0.35), "2026년 06월 정산 운송 마감 집계", 12, True, CHARCOAL)
add_textbox(slide6, web_l + Inches(0.5), web_t + Inches(1.05), web_w - Inches(1.0), Inches(0.5), 
            "• 총 운송 완료 건수: 240 회  |  • 총 누적 청구 운임: 36,000,000 원 (계산서 발행대기)", 11, False, SLATE_GRAY)

# 데이터 테이블 그리드 영역
table_rect = add_rounded_rect(slide6, web_l + Inches(0.3), web_t + Inches(1.9), web_w - Inches(0.6), Inches(2.2), WHITE, BORDER_GRAY)
add_textbox(slide6, web_l + Inches(0.4), web_t + Inches(2.0), web_w - Inches(0.8), Inches(2.0), 
            "일자           티켓번호     차량번호     하차지           금액           마감상태\n"
            "2026.06.05     #1094        12가3456     영종 사토장      150,000원      마감 대기\n"
            "2026.06.05     #1095        78나9012     영종 사토장      150,000원      마감 대기\n"
            "2026.06.04     #1088        34다5678     영종 사토장      150,000원      마감 대기\n"
            "2026.06.04     #1089        56라7890     영종 사토장      150,000원      마감 대기\n"
            "2026.06.03     #1072        90마1234     영종 사토장      150,000원      마감 완료", 10, False, CHARCOAL)

# 정산 처리 액션 버튼 세트
close_btn = add_rounded_rect(slide6, web_l + Inches(3.2), web_t + Inches(4.3), Inches(1.8), Inches(0.45), MUSTARD)
add_textbox(slide6, web_l + Inches(3.2), web_t + Inches(4.38), Inches(1.8), Inches(0.35), "정산 마감 확정", 11, True, CHARCOAL, PP_ALIGN.CENTER)

bill_btn = add_rounded_rect(slide6, web_l + Inches(5.1), web_t + Inches(4.3), Inches(1.9), Inches(0.45), CHARCOAL)
add_textbox(slide6, web_l + Inches(5.1), web_t + Inches(4.38), Inches(1.9), Inches(0.35), "세금계산서 발행 신청", 11, True, WHITE, PP_ALIGN.CENTER)


# 파워포인트 파일 저장
output_path = "D:\\Projects\\dumpring\\dumpring-platform-backend\\덤프링_현장관리자_화면상세설계.pptx"
prs.save(output_path)
print(f"PPTX 파일 생성 성공: {output_path}")
