import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# 컬러 팔레트 정의 (라이트 머스터드 옐로우 테마)
WHITE = RGBColor(255, 255, 255)
SAND_BG = RGBColor(249, 248, 246)
MUSTARD = RGBColor(245, 158, 11)     # F59E0B (주요 테마 색상)
AMBER = RGBColor(217, 119, 6)        # D97706 (포인트 색상)
CHARCOAL = RGBColor(31, 41, 55)      # 1F2937 (메인 텍스트)
LIGHT_GRAY = RGBColor(243, 244, 246)  # F3F4F6 (영역 구분 배경)
BORDER_GRAY = RGBColor(229, 231, 235) # E5E7EB (카드 테두리)
SLATE_GRAY = RGBColor(75, 85, 99)    # 4B5563 (서브 텍스트)
RED = RGBColor(239, 68, 68)          # EF4444 (반려/에러/취소/회차)
GREEN = RGBColor(16, 185, 129)       # 10B981 (정산완료/하차완료/정상)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6] # 빈 레이아웃

# 유틸리티 함수: 텍스트 상자 추가
def add_textbox(slide, left, top, width, height, text, font_size=11, bold=False, color=CHARCOAL, align=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = 'Malgun Gothic'
    if align:
        p.alignment = align
    return txBox

# 유틸리티 함수: 직사각형 도형 추가
def add_rectangle(slide, left, top, width, height, fill_color, line_color=None, line_width=1):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

# 유틸리티 함수: 둥근 직사각형 도형 추가
def add_rounded_rect(slide, left, top, width, height, fill_color, line_color=None, line_width=1):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

# 유틸리티 함수: 모바일 폰 프레임 생성
def add_phone_frame(slide, left, top, width, height, title_text="덤프링"):
    bezel = add_rounded_rect(slide, left, top, width, height, SAND_BG, CHARCOAL, 3)
    screen_left = left + Inches(0.15)
    screen_top = top + Inches(0.4)
    screen_width = width - Inches(0.3)
    screen_height = height - Inches(0.6)
    screen = add_rectangle(slide, screen_left, screen_top, screen_width, screen_height, WHITE, BORDER_GRAY, 1)
    
    # 상단 헤더
    header = add_rectangle(slide, screen_left, screen_top, screen_width, Inches(0.5), SAND_BG, BORDER_GRAY, 1)
    add_textbox(slide, screen_left + Inches(0.15), screen_top + Inches(0.12), screen_width - Inches(0.3), Inches(0.3), title_text, 13, True, CHARCOAL)
    
    return screen_left, screen_top + Inches(0.5), screen_width, screen_height - Inches(0.5)


# =============================================================
# SLIDE 1: 커버 페이지 (Cover Page)
# =============================================================
slide1 = prs.slides.add_slide(blank_layout)
add_rectangle(slide1, Inches(0), Inches(0), Inches(13.33), Inches(7.5), SAND_BG)
add_rectangle(slide1, Inches(1.5), Inches(2.2), Inches(0.15), Inches(2.5), MUSTARD)
add_textbox(slide1, Inches(2.0), Inches(2.1), Inches(9.5), Inches(1.2), "덤프링 (Dumpring)", 44, True, CHARCOAL)
add_textbox(slide1, Inches(2.0), Inches(3.2), Inches(9.5), Inches(0.8), "06. 하차지 이용자 (Landowner) 업무 화면 상세설계서", 28, True, SLATE_GRAY)
add_textbox(slide1, Inches(2.0), Inches(4.1), Inches(9.5), Inches(0.5), "최종 테마: 라이트 머스터드 옐로우 (Light Mustard Yellow)", 16, False, AMBER)
add_textbox(slide1, Inches(2.0), Inches(5.8), Inches(6.0), Inches(0.8), "작성일: 2026. 06. 05\n작성자: 덤프링 플랫폼 기획팀", 12, False, SLATE_GRAY)


# =============================================================
# SLIDE 2: Screen 1. 하차지 홈 (사토장 대시보드)
# =============================================================
slide2 = prs.slides.add_slide(blank_layout)
add_rectangle(slide2, Inches(0), Inches(0), Inches(13.33), Inches(7.5), WHITE)

# 좌측 설명 영역
add_textbox(slide2, Inches(0.8), Inches(0.8), Inches(5.5), Inches(0.8), "01. 하차지 홈 (사토장 대시보드)", 24, True, CHARCOAL)
add_textbox(slide2, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0), 
            "■ 주요 기능\n"
            "- 사토장의 전체 토사 수용 한도 및 현재 매립 잔여 상태를 확인하는 대시보드\n"
            "- 금일 반입 트럭 수 및 당월 누적 정산금액 실시간 모니터링\n"
            "- 지오펜싱(200m 이내) 진입 신호가 감지된 실시간 차량 목록 노출\n\n"
            "■ 주요 UI 구성\n"
            "1. 사토 수용 한도 게이지:\n"
            "   - '6,500㎥ 수용 중 / 총 10,000㎥ 수용 한도' 가로 프로그레스 바\n"
            "2. 금일 현황판:\n"
            "   - 오늘 반입 차량 15대, 누적 정산금: 2,250,000원\n"
            "3. 실시간 진입 관제 (GPS 지오펜싱):\n"
            "   - 게이트 인근 200m 이내로 도달한 트럭의 실시간 상태 리스트", 
            12, False, CHARCOAL)

# 우측 폰 프레임 생성
sc_l, sc_t, sc_w, sc_h = add_phone_frame(slide2, Inches(7.5), Inches(0.5), Inches(4.5), Inches(6.5), "하차지 대시보드")

# 수용량 게이지 카드
cap_rect = add_rounded_rect(slide2, sc_l + Inches(0.15), sc_t + Inches(0.15), sc_w - Inches(0.3), Inches(1.1), WHITE, BORDER_GRAY)
add_textbox(slide2, sc_l + Inches(0.3), sc_t + Inches(0.22), sc_w - Inches(0.6), Inches(0.3), "사토장 잔여 수용량 (65%)", 10, True, CHARCOAL)
# 프로그레스바 배경
add_rounded_rect(slide2, sc_l + Inches(0.3), sc_t + Inches(0.55), sc_w - Inches(0.6), Inches(0.2), LIGHT_GRAY, BORDER_GRAY)
# 채워진바
add_rounded_rect(slide2, sc_l + Inches(0.3), sc_t + Inches(0.55), (sc_w - Inches(0.6)) * 0.65, Inches(0.2), MUSTARD)
# 텍스트
add_textbox(slide2, sc_l + Inches(0.3), sc_t + Inches(0.8), sc_w - Inches(0.6), Inches(0.35), "6,500㎥ 반입 완료 / 총 10,000㎥", 9, False, SLATE_GRAY)

# 실시간 현황 요약
summary_y = sc_t + Inches(1.4)
add_rounded_rect(slide2, sc_l + Inches(0.15), summary_y, sc_w - Inches(0.3), Inches(0.9), WHITE, BORDER_GRAY)
add_textbox(slide2, sc_l + Inches(0.3), summary_y + Inches(0.1), sc_w - Inches(0.6), Inches(0.35), "금일 반입 대수:  15 대", 11, True, CHARCOAL)
add_textbox(slide2, sc_l + Inches(0.3), summary_y + Inches(0.45), sc_w - Inches(0.6), Inches(0.35), "금일 수금/정산액: 2,250,000 원", 11, True, AMBER)

# 실시간 진입 차량 리스트
list_y = sc_t + Inches(2.4)
add_textbox(slide2, sc_l + Inches(0.2), list_y, sc_w - Inches(0.4), Inches(0.3), "실시간 진입 차량 (200m 내)", 11, True, SLATE_GRAY)

# 트럭 1
i1_y = list_y + Inches(0.3)
add_rounded_rect(slide2, sc_l + Inches(0.15), i1_y, sc_w - Inches(0.3), Inches(0.75), WHITE, BORDER_GRAY)
add_textbox(slide2, sc_l + Inches(0.3), i1_y + Inches(0.08), sc_w - Inches(0.6), Inches(0.35), "경기 80사 1234 (김덤프 기사)", 11, True, CHARCOAL)
add_textbox(slide2, sc_l + Inches(0.3), i1_y + Inches(0.38), sc_w - Inches(0.6), Inches(0.3), "양질토  |  진입 완료 (50m 전)", 9, False, GREEN)

# 트럭 2
i2_y = list_y + Inches(1.15)
add_rounded_rect(slide2, sc_l + Inches(0.15), i2_y, sc_w - Inches(0.3), Inches(0.75), WHITE, BORDER_GRAY)
add_textbox(slide2, sc_l + Inches(0.3), i2_y + Inches(0.08), sc_w - Inches(0.6), Inches(0.35), "서울 80사 5678 (이배차 기사)", 11, True, CHARCOAL)
add_textbox(slide2, sc_l + Inches(0.3), i2_y + Inches(0.38), sc_w - Inches(0.6), Inches(0.3), "양질토  |  접근 중 (150m 전)", 9, False, SLATE_GRAY)


# =============================================================
# SLIDE 3: Screen 2. 하차지 반입 승인 및 토질 검사
# =============================================================
slide3 = prs.slides.add_slide(blank_layout)
add_rectangle(slide3, Inches(0), Inches(0), Inches(13.33), Inches(7.5), WHITE)

# 좌측 설명 영역
add_textbox(slide3, Inches(0.8), Inches(0.8), Inches(5.5), Inches(0.8), "02. 반입 승인 및 토질 검증 통제", 24, True, CHARCOAL)
add_textbox(slide3, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0), 
            "■ 주요 기능\n"
            "- 하차지에 도달하여 반입을 기다리는 개별 차량의 승인/회차를 처리하는 화면\n"
            "- 불량 토사(뻘흙 등) 반입을 필터링하기 위한 육안 검증 확인 장치\n"
            "- 반려/거부 시 회차증 발급을 위한 즉각적 현장 증빙 사진 업로드 기능\n\n"
            "■ 주요 UI 구성\n"
            "1. 진입 차량 요약 정보:\n"
            "   - 차량번호, 운송장 정보 및 적재하고 온 선언 토사 명시\n"
            "2. 육안 토질 검증 체크리스트:\n"
            "   - 양질토 여부 및 혼입물 없음 확인 체크 토글\n"
            "3. 반입 제어 버튼 세트:\n"
            "   - [최종 반입 승인] 녹색 버튼 (정산 확정 트리거)\n"
            "   - [반입 거부 및 회차] 적색 버튼 (분쟁 중재 접수)", 
            12, False, CHARCOAL)

# 우측 폰 프레임 생성
sc_l, sc_t, sc_w, sc_h = add_phone_frame(slide3, Inches(7.5), Inches(0.5), Inches(4.5), Inches(6.5), "반입 승인 통제")

# 대상 차량 상세 카드
add_textbox(slide3, sc_l + Inches(0.2), sc_t + Inches(0.2), sc_w - Inches(0.4), Inches(0.3), "진입 차량 상세 정보", 10, True, SLATE_GRAY)
v_card = add_rounded_rect(slide3, sc_l + Inches(0.15), sc_t + Inches(0.5), sc_w - Inches(0.3), Inches(1.0), WHITE, BORDER_GRAY)
add_textbox(slide3, sc_l + Inches(0.3), sc_t + Inches(0.58), sc_w - Inches(0.6), Inches(0.35), "경기 80사 1234 (25톤)", 12, True, CHARCOAL)
add_textbox(slide3, sc_l + Inches(0.3), sc_t + Inches(0.88), sc_w - Inches(0.6), Inches(0.5), 
            "• 상차지: 송도 1공구 아파트 현장\n• 자재 속성: 양질토 선언", 9, False, SLATE_GRAY)

# 육안 검사 확인란
add_textbox(slide3, sc_l + Inches(0.2), sc_t + Inches(1.6), sc_w - Inches(0.4), Inches(0.3), "현장 육안 검토 사항", 10, True, SLATE_GRAY)
add_rounded_rect(slide3, sc_l + Inches(0.15), sc_t + Inches(1.9), sc_w - Inches(0.3), Inches(0.9), WHITE, BORDER_GRAY)
add_textbox(slide3, sc_l + Inches(0.3), sc_t + Inches(2.0), sc_w - Inches(0.6), Inches(0.7), 
            "☑ 토질 적합 검토 (양질토 확인)\n☑ 환경 오염 유발 혼합 토사 없음 확인", 10, False, CHARCOAL)

# 거부 시 증빙 사진 업로드 묘사
add_textbox(slide3, sc_l + Inches(0.2), sc_t + Inches(2.9), sc_w - Inches(0.4), Inches(0.3), "반려 및 회차 시 증빙 자료 첨부", 10, True, SLATE_GRAY)
# 썸네일 박스
add_rounded_rect(slide3, sc_l + Inches(0.15), sc_t + Inches(3.2), Inches(1.4), Inches(1.0), LIGHT_GRAY, BORDER_GRAY)
add_textbox(slide3, sc_l + Inches(0.15), sc_t + Inches(3.55), Inches(1.4), Inches(0.3), "[사진 촬영/첨부]", 8, False, SLATE_GRAY, PP_ALIGN.CENTER)

# 하단 승인/반려 제어 버튼
app_btn = add_rounded_rect(slide3, sc_l + Inches(0.2), sc_t + Inches(4.5), sc_w - Inches(0.4), Inches(0.5), GREEN)
add_textbox(slide3, sc_l + Inches(0.2), sc_t + Inches(4.62), sc_w - Inches(0.4), Inches(0.35), "최종 반입 승인 (정산 완료)", 12, True, WHITE, PP_ALIGN.CENTER)

rej_btn = add_rounded_rect(slide3, sc_l + Inches(0.2), sc_t + Inches(5.1), sc_w - Inches(0.4), Inches(0.5), RED)
add_textbox(slide3, sc_l + Inches(0.2), sc_t + Inches(5.22), sc_w - Inches(0.4), Inches(0.35), "반입 거부 및 회차 지시", 12, True, WHITE, PP_ALIGN.CENTER)


# =============================================================
# SLIDE 4: Screen 3. 하차지 설정 및 수용 공고 등록
# =============================================================
slide4 = prs.slides.add_slide(blank_layout)
add_rectangle(slide4, Inches(0), Inches(0), Inches(13.33), Inches(7.5), WHITE)

# 좌측 설명 영역
add_textbox(slide4, Inches(0.8), Inches(0.8), Inches(5.5), Inches(0.8), "03. 하차지 설정 및 단가 등록", 24, True, CHARCOAL)
add_textbox(slide4, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0), 
            "■ 주요 기능\n"
            "- 신규 사토장을 시스템에 설정하고 수용 정책을 게시하여 매칭을 트리거하는 화면\n"
            "- 지오펜싱 입차 감지를 위한 하차지 정밀 중심 좌표 좌표 매핑\n"
            "- 수용 가능한 토종 성질(양질토/뻘흙/암버럭) 및 회당 비용 정산 정책 기입\n\n"
            "■ 주요 UI 구성\n"
            "1. 사토장 명칭 및 주소 입력:\n"
            "   - '영종도 남단 사토 매립지' 등 행정 명칭 및 주소 입력 필드\n"
            "2. 지오펜싱 입차 반경 설정 바:\n"
            "   - 반입 차량 자동 체크를 위한 지오펜싱 기본 200m 세팅\n"
            "3. 수용 자재 토글:\n"
            "   - 양질토(수용), 뻘흙(반려), 암버럭(수용) 등 자재 제한 조건 선택\n"
            "4. 단가 및 정산 정책:\n"
            "   - 회당 150,000원 수금 조건 설정 필드", 
            12, False, CHARCOAL)

# 우측 폰 프레임 생성
sc_l, sc_t, sc_w, sc_h = add_phone_frame(slide4, Inches(7.5), Inches(0.5), Inches(4.5), Inches(6.5), "사토장 설정")

# 사토장 명칭
add_textbox(slide4, sc_l + Inches(0.2), sc_t + Inches(0.2), sc_w - Inches(0.4), Inches(0.3), "사토장 주소지 및 명칭", 10, True, SLATE_GRAY)
add_rounded_rect(slide4, sc_l + Inches(0.2), sc_t + Inches(0.45), sc_w - Inches(0.4), Inches(0.45), WHITE, BORDER_GRAY)
add_textbox(slide4, sc_l + Inches(0.35), sc_t + Inches(0.53), sc_w - Inches(0.7), Inches(0.3), "인천시 중구 영종남로 11-1 (영종 사토장)", 11, True, CHARCOAL)

# 수용 자재 토글 3종
add_textbox(slide4, sc_l + Inches(0.2), sc_t + Inches(1.0), sc_w - Inches(0.4), Inches(0.3), "수용 가능 자재 필터", 10, True, SLATE_GRAY)
# 1
add_rounded_rect(slide4, sc_l + Inches(0.2), sc_t + Inches(1.3), Inches(1.2), Inches(0.4), MUSTARD)
add_textbox(slide4, sc_l + Inches(0.2), sc_t + Inches(1.38), Inches(1.2), Inches(0.3), "양질토 수용", 9, True, CHARCOAL, PP_ALIGN.CENTER)
# 2
add_rounded_rect(slide4, sc_l + Inches(1.5), sc_t + Inches(1.3), Inches(1.2), Inches(0.4), LIGHT_GRAY, BORDER_GRAY)
add_textbox(slide4, sc_l + Inches(1.5), sc_t + Inches(1.38), Inches(1.2), Inches(0.3), "뻘흙 제한", 9, False, SLATE_GRAY, PP_ALIGN.CENTER)
# 3
add_rounded_rect(slide4, sc_l + Inches(2.8), sc_t + Inches(1.3), Inches(1.2), Inches(0.4), MUSTARD)
add_textbox(slide4, sc_l + Inches(2.8), sc_t + Inches(1.38), Inches(1.2), Inches(0.3), "암버럭 수용", 9, True, CHARCOAL, PP_ALIGN.CENTER)

# 단가 정책 설정
add_textbox(slide4, sc_l + Inches(0.2), sc_t + Inches(1.85), sc_w - Inches(0.4), Inches(0.3), "반입 비용 수금 정책", 10, True, SLATE_GRAY)
add_rounded_rect(slide4, sc_l + Inches(0.2), sc_t + Inches(2.15), sc_w - Inches(0.4), Inches(0.45), WHITE, BORDER_GRAY)
add_textbox(slide4, sc_l + Inches(0.35), sc_t + Inches(2.23), sc_w - Inches(0.7), Inches(0.3), "회당 단가 수금:  150,000 원", 11, True, AMBER)

# 지오펜싱 슬라이더
add_textbox(slide4, sc_l + Inches(0.2), sc_t + Inches(2.75), sc_w - Inches(0.4), Inches(0.3), "하차지 인식 지오펜싱 반경: 200m", 10, True, SLATE_GRAY)
add_rectangle(slide4, sc_l + Inches(0.2), sc_t + Inches(3.1), sc_w - Inches(0.4), Inches(0.06), BORDER_GRAY)
add_rounded_rect(slide4, sc_l + Inches(1.5), sc_t + Inches(3.0), Inches(0.25), Inches(0.25), MUSTARD)

# 완료 단추
save_btn = add_rounded_rect(slide4, sc_l + Inches(0.2), sc_t + Inches(4.7), sc_w - Inches(0.4), Inches(0.6), CHARCOAL)
add_textbox(slide4, sc_l + Inches(0.2), sc_t + Inches(4.85), sc_w - Inches(0.4), Inches(0.4), "수용 공고 설정 및 배포", 14, True, WHITE, PP_ALIGN.CENTER)


# =============================================================
# SLIDE 5: Screen 4. 월별 반입 정산 및 명세서 조회
# =============================================================
slide5 = prs.slides.add_slide(blank_layout)
add_rectangle(slide5, Inches(0), Inches(0), Inches(13.33), Inches(7.5), WHITE)

# 좌측 설명 영역
add_textbox(slide5, Inches(0.8), Inches(0.8), Inches(5.5), Inches(0.8), "04. 하차지 정산 마감 및 세금 거래", 24, True, CHARCOAL)
add_textbox(slide5, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0), 
            "■ 주요 기능\n"
            "- 하차지에서 정상 수입 완료된 전표에 대한 월별 마감 및 정산 대장 확인 화면\n"
            "- 총 수금 완료된 단가 합산 및 수수료 명세 자동 디스플레이\n"
            "- 엑셀 세부 운송장 파일 다운로드 및 국세청 연동 세금계산서 청구 지원\n\n"
            "■ 주요 UI 구성\n"
            "1. 6월 정산 거래 요약:\n"
            "   - 총 반입 건수: 240회, 총 정산 금액: 36,000,000원\n"
            "2. 일별 반입 이력 그리드:\n"
            "   - 날짜별 입차 성공 대수 및 일일 누계 정산액 일람 피드\n"
            "3. 마감 다운로드 액션 단추:\n"
            "   - [정산서 Excel 다운로드] 차콜 단추\n"
            "   - [세금계산서 청구 발행] 녹색 버튼", 
            12, False, CHARCOAL)

# 우측 폰 프레임 생성
sc_l, sc_t, sc_w, sc_h = add_phone_frame(slide5, Inches(7.5), Inches(0.5), Inches(4.5), Inches(6.5), "하차지 마감 정산")

# 마감 요약 카드
summary_rect = add_rounded_rect(slide5, sc_l + Inches(0.15), sc_t + Inches(0.15), sc_w - Inches(0.3), Inches(1.3), WHITE, BORDER_GRAY)
add_textbox(slide5, sc_l + Inches(0.3), sc_t + Inches(0.22), sc_w - Inches(0.6), Inches(0.3), "2026년 6월 반입 마감 명세", 12, True, CHARCOAL)
add_textbox(slide5, sc_l + Inches(0.3), sc_t + Inches(0.55), sc_w - Inches(0.6), Inches(0.7), 
            "• 총 반입 완료 횟수:  240 회\n• 누적 수금/정산액: 36,000,000 원", 10, False, SLATE_GRAY)
badge_c = add_rounded_rect(slide5, sc_l + sc_w - Inches(1.4), sc_t + Inches(0.2), Inches(1.1), Inches(0.32), GREEN)
add_textbox(slide5, sc_l + sc_w - Inches(1.4), sc_t + Inches(0.24), Inches(1.1), Inches(0.25), "정산 완료", 9, True, WHITE, PP_ALIGN.CENTER)

# 일별 리스트
add_textbox(slide5, sc_l + Inches(0.2), sc_t + Inches(1.6), sc_w - Inches(0.4), Inches(0.3), "일별 반입 통계 리스트", 11, True, SLATE_GRAY)

# 일별 1
d1_y = sc_t + Inches(1.9)
add_rounded_rect(slide5, sc_l + Inches(0.15), d1_y, sc_w - Inches(0.3), Inches(0.65), WHITE, BORDER_GRAY)
add_textbox(slide5, sc_l + Inches(0.3), d1_y + Inches(0.12), sc_w - Inches(0.6), Inches(0.35), "06월 05일  |  반입 15회  |  2,250,000원", 11, False, CHARCOAL)

# 일별 2
d2_y = sc_t + Inches(2.6)
add_rounded_rect(slide5, sc_l + Inches(0.15), d2_y, sc_w - Inches(0.3), Inches(0.65), WHITE, BORDER_GRAY)
add_textbox(slide5, sc_l + Inches(0.3), d2_y + Inches(0.12), sc_w - Inches(0.6), Inches(0.35), "06월 04일  |  반입 20회  |  3,000,000원", 11, False, CHARCOAL)

# 일별 3
d3_y = sc_t + Inches(3.3)
add_rounded_rect(slide5, sc_l + Inches(0.15), d3_y, sc_w - Inches(0.3), Inches(0.65), WHITE, BORDER_GRAY)
add_textbox(slide5, sc_l + Inches(0.3), d3_y + Inches(0.12), sc_w - Inches(0.6), Inches(0.35), "06월 03일  |  반입 18회  |  2,700,000원", 11, False, CHARCOAL)

# 하단 액션 버튼
excel_btn = add_rounded_rect(slide5, sc_l + Inches(0.2), sc_t + Inches(4.3), sc_w - Inches(0.4), Inches(0.5), CHARCOAL)
add_textbox(slide5, sc_l + Inches(0.2), sc_t + Inches(4.42), sc_w - Inches(0.4), Inches(0.35), "정산서 Excel 파일 발송", 11, True, WHITE, PP_ALIGN.CENTER)

tax_btn = add_rounded_rect(slide5, sc_l + Inches(0.2), sc_t + Inches(4.9), sc_w - Inches(0.4), Inches(0.5), GREEN)
add_textbox(slide5, sc_l + Inches(0.2), sc_t + Inches(5.02), sc_w - Inches(0.4), Inches(0.35), "세금계산서 발행 완료 요청", 11, True, WHITE, PP_ALIGN.CENTER)


# 파워포인트 파일 저장
output_path = "D:\\Projects\\dumpring\\dumpring-platform-backend\\덤프링_하차지이용자_화면상세설계.pptx"
prs.save(output_path)
print(f"PPTX 파일 생성 성공: {output_path}")
