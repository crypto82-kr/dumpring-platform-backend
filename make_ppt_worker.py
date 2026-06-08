import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# 컬러 팔레트 정의 (라이트 머스터드 옐로우 테마)
WHITE = RGBColor(255, 255, 255)
SAND_BG = RGBColor(249, 248, 246)
MUSTARD = RGBColor(245, 158, 11)     # F59E0B (주요 테마 색상)
AMBER = RGBColor(217, 119, 6)        # D97706 (포인트 색상)
CHARCOAL = RGBColor(31, 41, 55)      # 1F2937 (메인 텍스트)
LIGHT_GRAY = RGBColor(243, 244, 246)  # F3F4F6 (영역 구분 배경)
BORDER_GRAY = RGBColor(229, 231, 235) # E5E7EB (카드 테두리)
SLATE_GRAY = RGBColor(75, 85, 99)    # 4B5563 (서브 텍스트)
RED = RGBColor(239, 68, 68)          # EF4444 (반려/에러/취소/회차)
GREEN = RGBColor(16, 185, 129)       # 10B981 (정산완료/하차완료)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6] # 빈 레이아웃

# 유틸리티 함수: 텍스트 상자 추가
def add_textbox(slide, left, top, width, height, text, font_size=11, bold=False, color=CHARCOAL, align=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = 'Malgun Gothic'
    if align:
        p.alignment = align
    return txBox

# 유틸리티 함수: 직사각형 도형 추가
def add_rectangle(slide, left, top, width, height, fill_color, line_color=None, line_width=1):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

# 유틸리티 함수: 둥근 직사각형 도형 추가
def add_rounded_rect(slide, left, top, width, height, fill_color, line_color=None, line_width=1):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

# 유틸리티 함수: 모바일 폰 프레임 생성
def add_phone_frame(slide, left, top, width, height, title_text="덤프링"):
    bezel = add_rounded_rect(slide, left, top, width, height, SAND_BG, CHARCOAL, 3)
    screen_left = left + Inches(0.15)
    screen_top = top + Inches(0.4)
    screen_width = width - Inches(0.3)
    screen_height = height - Inches(0.6)
    screen = add_rectangle(slide, screen_left, screen_top, screen_width, screen_height, WHITE, BORDER_GRAY, 1)
    
    # 상단 헤더
    header = add_rectangle(slide, screen_left, screen_top, screen_width, Inches(0.5), SAND_BG, BORDER_GRAY, 1)
    add_textbox(slide, screen_left + Inches(0.15), screen_top + Inches(0.12), screen_width - Inches(0.3), Inches(0.3), title_text, 13, True, CHARCOAL)
    
    return screen_left, screen_top + Inches(0.5), screen_width, screen_height - Inches(0.5)


# =============================================================
# SLIDE 1: 커버 페이지 (Cover Page)
# =============================================================
slide1 = prs.slides.add_slide(blank_layout)
add_rectangle(slide1, Inches(0), Inches(0), Inches(13.33), Inches(7.5), SAND_BG)
add_rectangle(slide1, Inches(1.5), Inches(2.2), Inches(0.15), Inches(2.5), MUSTARD)
add_textbox(slide1, Inches(2.0), Inches(2.1), Inches(9.5), Inches(1.2), "덤프링 (Dumpring)", 44, True, CHARCOAL)
add_textbox(slide1, Inches(2.0), Inches(3.2), Inches(9.5), Inches(0.8), "05. 현장 담당자 (Site Worker) 업무 화면 상세설계서", 28, True, SLATE_GRAY)
add_textbox(slide1, Inches(2.0), Inches(4.1), Inches(9.5), Inches(0.5), "최종 테마: 라이트 머스터드 옐로우 (Light Mustard Yellow)", 16, False, AMBER)
add_textbox(slide1, Inches(2.0), Inches(5.8), Inches(6.0), Inches(0.8), "작성일: 2026. 06. 05\n작성자: 덤프링 플랫폼 기획팀", 12, False, SLATE_GRAY)


# =============================================================
# SLIDE 2: Screen 1. 현장 담당자 홈 (게이트 관제)
# =============================================================
slide2 = prs.slides.add_slide(blank_layout)
add_rectangle(slide2, Inches(0), Inches(0), Inches(13.33), Inches(7.5), WHITE)

# 좌측 설명 영역
add_textbox(slide2, Inches(0.8), Inches(0.8), Inches(5.5), Inches(0.8), "01. 현장 담당자 홈 (게이트 관제)", 24, True, CHARCOAL)
add_textbox(slide2, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0), 
            "■ 주요 기능\n"
            "- 현장에 매핑된 담당 직원이 실시간 입차 트럭 대수 및 통계를 조회하는 메인 대시보드\n"
            "- 상차 대기 상태인 덤프 트럭의 실시간 큐(Queue) 리스트 모니터링\n"
            "- 수동 진입 체크인을 통해 GPS 오류 차량도 즉시 처리 가능\n\n"
            "■ 주요 UI 구성\n"
            "1. 금일 현장 출입 요약 (3단 메트릭 카드):\n"
            "   - 입차대기: 2대, 상차중: 3대, 상차완료: 18대 표시\n"
            "2. 실시간 입차 대기 큐:\n"
            "   - 트럭 번호판, 기사명, 대기 시간 및 입차 방식(지오펜싱 자동/수동) 노출\n"
            "3. [수동 체크인 / QR 스캔] 퀵 버튼:\n"
            "   - 하단에 상시 노출하여 신속한 현장 통제 지원", 
            12, False, CHARCOAL)

# 우측 폰 프레임 생성
sc_l, sc_t, sc_w, sc_h = add_phone_frame(slide2, Inches(7.5), Inches(0.5), Inches(4.5), Inches(6.5), "송도 1공구 - 게이트 관제")

# 3단 메트릭 카드
card_w = (sc_w - Inches(0.5)) / 3
card_h = Inches(0.8)

# 카드 1
add_rounded_rect(slide2, sc_l + Inches(0.15), sc_t + Inches(0.15), card_w, card_h, WHITE, BORDER_GRAY)
add_textbox(slide2, sc_l + Inches(0.15), sc_t + Inches(0.22), card_w, Inches(0.25), "입차대기", 9, False, SLATE_GRAY, PP_ALIGN.CENTER)
add_textbox(slide2, sc_l + Inches(0.15), sc_t + Inches(0.48), card_w, Inches(0.35), "2 대", 13, True, MUSTARD, PP_ALIGN.CENTER)

# 카드 2
add_rounded_rect(slide2, sc_l + Inches(0.15) + card_w + Inches(0.1), sc_t + Inches(0.15), card_w, card_h, WHITE, BORDER_GRAY)
add_textbox(slide2, sc_l + Inches(0.15) + card_w + Inches(0.1), sc_t + Inches(0.22), card_w, Inches(0.25), "상차중", 9, False, SLATE_GRAY, PP_ALIGN.CENTER)
add_textbox(slide2, sc_l + Inches(0.15) + card_w + Inches(0.1), sc_t + Inches(0.48), card_w, Inches(0.35), "3 대", 13, True, AMBER, PP_ALIGN.CENTER)

# 카드 3
add_rounded_rect(slide2, sc_l + Inches(0.15) + (card_w * 2) + Inches(0.2), sc_t + Inches(0.15), card_w, card_h, WHITE, BORDER_GRAY)
add_textbox(slide2, sc_l + Inches(0.15) + (card_w * 2) + Inches(0.2), sc_t + Inches(0.22), card_w, Inches(0.25), "상차완료", 9, False, SLATE_GRAY, PP_ALIGN.CENTER)
add_textbox(slide2, sc_l + Inches(0.15) + (card_w * 2) + Inches(0.2), sc_t + Inches(0.48), card_w, Inches(0.35), "18 대", 13, True, CHARCOAL, PP_ALIGN.CENTER)

# 실시간 입차 대기 큐 리스트
list_y = sc_t + Inches(1.1)
add_textbox(slide2, sc_l + Inches(0.2), list_y, sc_w - Inches(0.4), Inches(0.3), "실시간 진입 대기 트럭 (2)", 11, True, SLATE_GRAY)

# 리스트 아이템 1
i1_y = list_y + Inches(0.3)
add_rounded_rect(slide2, sc_l + Inches(0.15), i1_y, sc_w - Inches(0.3), Inches(0.8), WHITE, BORDER_GRAY)
add_textbox(slide2, sc_l + Inches(0.3), i1_y + Inches(0.1), sc_w - Inches(0.6), Inches(0.35), "경기 80사 1234 (25톤)", 11, True, CHARCOAL)
add_textbox(slide2, sc_l + Inches(0.3), i1_y + Inches(0.42), sc_w - Inches(0.6), Inches(0.3), "김덤프 기사  |  지오펜싱 입차 (5분 전)", 9, False, GREEN)
badge1 = add_rounded_rect(slide2, sc_l + sc_w - Inches(1.3), i1_y + Inches(0.1), Inches(1.0), Inches(0.3), MUSTARD)
add_textbox(slide2, sc_l + sc_w - Inches(1.3), i1_y + Inches(0.12), Inches(1.0), Inches(0.25), "대기중", 8, True, CHARCOAL, PP_ALIGN.CENTER)

# 리스트 아이템 2
i2_y = list_y + Inches(1.2)
add_rounded_rect(slide2, sc_l + Inches(0.15), i2_y, sc_w - Inches(0.3), Inches(0.8), WHITE, BORDER_GRAY)
add_textbox(slide2, sc_l + Inches(0.3), i2_y + Inches(0.1), sc_w - Inches(0.6), Inches(0.35), "서울 80사 5678 (25톤)", 11, True, CHARCOAL)
add_textbox(slide2, sc_l + Inches(0.3), i2_y + Inches(0.42), sc_w - Inches(0.6), Inches(0.3), "이배차 기사  |  수동 체크인 (방금 전)", 9, False, SLATE_GRAY)
badge2 = add_rounded_rect(slide2, sc_l + sc_w - Inches(1.3), i2_y + Inches(0.1), Inches(1.0), Inches(0.3), MUSTARD)
add_textbox(slide2, sc_l + sc_w - Inches(1.3), i2_y + Inches(0.12), Inches(1.0), Inches(0.25), "대기중", 8, True, CHARCOAL, PP_ALIGN.CENTER)

# 하단 수동 체크인 버튼
chk_btn = add_rounded_rect(slide2, sc_l + Inches(0.2), sc_t + Inches(4.7), sc_w - Inches(0.4), Inches(0.6), MUSTARD)
add_textbox(slide2, sc_l + Inches(0.2), sc_t + Inches(4.85), sc_w - Inches(0.4), Inches(0.4), "수동 체크인 / QR 스캔", 14, True, CHARCOAL, PP_ALIGN.CENTER)


# =============================================================
# SLIDE 3: Screen 2. 입/출차 수동 통제 및 QR 승인
# =============================================================
slide3 = prs.slides.add_slide(blank_layout)
add_rectangle(slide3, Inches(0), Inches(0), Inches(13.33), Inches(7.5), WHITE)

# 좌측 설명 영역
add_textbox(slide3, Inches(0.8), Inches(0.8), Inches(5.5), Inches(0.8), "02. 입/출차 수동 통제 및 상차 승인", 24, True, CHARCOAL)
add_textbox(slide3, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0), 
            "■ 주요 기능\n"
            "- 현장에 도달한 트럭의 상차 개시 및 반출 상태를 수동으로 승인하는 상세 폼 화면\n"
            "- 토종(자재 속성) 선택 및 상차 게이트 지정을 통해 운송장의 무결성 보장\n"
            "- [상차 승인] 터치 시, 기사 앱에 실시간으로 운송 단가와 GPS 미터기 주행이 즉시 작동됨\n\n"
            "■ 주요 UI 구성\n"
            "1. 대상 차량 정보:\n"
            "   - 차량 번호, 매핑 기사명, 톤수 정보 자동 매핑 출력\n"
            "2. 반출 토사 속성 토글:\n"
            "   - 양질토(선택), 뻘흙, 암버럭 등 자재 물리 정보 매핑 단추\n"
            "3. 담당 게이트 선택:\n"
            "   - 다수 게이트 운영 시 적정한 진입 통제 게이트 설정 드롭다운\n"
            "4. [상차 승인 (출발 처리)] 하단 골드 메인 버튼", 
            12, False, CHARCOAL)

# 우측 폰 프레임 생성
sc_l, sc_t, sc_w, sc_h = add_phone_frame(slide3, Inches(7.5), Inches(0.5), Inches(4.5), Inches(6.5), "상차 승인 통제")

# 대상 차량 정보 카드
add_textbox(slide3, sc_l + Inches(0.2), sc_t + Inches(0.2), sc_w - Inches(0.4), Inches(0.3), "승인 대상 차량", 10, True, SLATE_GRAY)
v_info = add_rounded_rect(slide3, sc_l + Inches(0.15), sc_t + Inches(0.5), sc_w - Inches(0.3), Inches(0.9), WHITE, BORDER_GRAY)
add_textbox(slide3, sc_l + Inches(0.3), sc_t + Inches(0.58), sc_w - Inches(0.6), Inches(0.3), "경기 80사 1234 (25톤)", 12, True, CHARCOAL)
add_textbox(slide3, sc_l + Inches(0.3), sc_t + Inches(0.88), sc_w - Inches(0.6), Inches(0.3), "소속: 대박운수  |  기사: 김덤프 기사", 10, False, SLATE_GRAY)

# 토사 속성 선택
add_textbox(slide3, sc_l + Inches(0.2), sc_t + Inches(1.6), sc_w - Inches(0.4), Inches(0.3), "반출 토사 종류", 10, True, SLATE_GRAY)
add_rounded_rect(slide3, sc_l + Inches(0.2), sc_t + Inches(1.9), Inches(1.2), Inches(0.4), MUSTARD)
add_textbox(slide3, sc_l + Inches(0.2), sc_t + Inches(1.98), Inches(1.2), Inches(0.3), "양질토", 10, True, CHARCOAL, PP_ALIGN.CENTER)

add_rounded_rect(slide3, sc_l + Inches(1.5), sc_t + Inches(1.9), Inches(1.2), Inches(0.4), WHITE, BORDER_GRAY)
add_textbox(slide3, sc_l + Inches(1.5), sc_t + Inches(1.98), Inches(1.2), Inches(0.3), "뻘흙", 10, False, SLATE_GRAY, PP_ALIGN.CENTER)

add_rounded_rect(slide3, sc_l + Inches(2.8), sc_t + Inches(1.9), Inches(1.2), Inches(0.4), WHITE, BORDER_GRAY)
add_textbox(slide3, sc_l + Inches(2.8), sc_t + Inches(1.98), Inches(1.2), Inches(0.3), "암버럭", 10, False, SLATE_GRAY, PP_ALIGN.CENTER)

# 상차 게이트 선택
add_textbox(slide3, sc_l + Inches(0.2), sc_t + Inches(2.5), sc_w - Inches(0.4), Inches(0.3), "상차 통제 게이트", 10, True, SLATE_GRAY)
gate_dd = add_rounded_rect(slide3, sc_l + Inches(0.2), sc_t + Inches(2.8), sc_w - Inches(0.4), Inches(0.45), LIGHT_GRAY, BORDER_GRAY)
add_textbox(slide3, sc_l + Inches(0.35), sc_t + Inches(2.88), sc_w - Inches(0.7), Inches(0.3), "제 1 상차 게이트  ▼", 11, False, CHARCOAL)

# 출발 승인 버튼
start_btn = add_rounded_rect(slide3, sc_l + Inches(0.2), sc_t + Inches(4.7), sc_w - Inches(0.4), Inches(0.6), MUSTARD)
add_textbox(slide3, sc_l + Inches(0.2), sc_t + Inches(4.85), sc_w - Inches(0.4), Inches(0.4), "상차 승인 (출발 처리)", 14, True, CHARCOAL, PP_ALIGN.CENTER)


# =============================================================
# SLIDE 4: Screen 3. 당일 반출 티켓(운송장) 목록
# =============================================================
slide4 = prs.slides.add_slide(blank_layout)
add_rectangle(slide4, Inches(0), Inches(0), Inches(13.33), Inches(7.5), WHITE)

# 좌측 설명 영역
add_textbox(slide4, Inches(0.8), Inches(0.8), Inches(5.5), Inches(0.8), "03. 당일 반출 티켓(운송장) 목록", 24, True, CHARCOAL)
add_textbox(slide4, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0), 
            "■ 주요 기능\n"
            "- 당일 해당 게이트를 통과하여 출발한 모든 운송 전표의 상태를 추적하는 화면\n"
            "- 각 티켓의 현재 실시간 배차 운행 단계(운행중 / 하차완료 / 회차거부)를 한눈에 모니터링\n"
            "- 개별 티켓 클릭 시 상세 궤적 정보 및 모바일 영수증 이력 열람\n\n"
            "■ 주요 UI 구성\n"
            "1. 목록 필터링:\n"
            "   - 상차 완료 시각 및 차량번호 간편 검색 입력바\n"
            "2. 티켓 정보 카드:\n"
            "   - 고유 전표 번호 (#2045), 차량번호, 자재종류, 상차완료 시각 표시\n"
            "   - 운행 상태를 나타내는 색상 뱃지 구분 (운행중: 주황, 하차완료: 초록, 회차거부: 빨강)", 
            12, False, CHARCOAL)

# 우측 폰 프레임 생성
sc_l, sc_t, sc_w, sc_h = add_phone_frame(slide4, Inches(7.5), Inches(0.5), Inches(4.5), Inches(6.5), "당일 반출 티켓")

# 검색바
add_rounded_rect(slide4, sc_l + Inches(0.15), sc_t + Inches(0.15), sc_w - Inches(0.3), Inches(0.45), WHITE, BORDER_GRAY)
add_textbox(slide4, sc_l + Inches(0.35), sc_t + Inches(0.23), sc_w - Inches(0.7), Inches(0.3), "차량번호 또는 기사명 검색", 10, False, SLATE_GRAY)

# 티켓 1 (운행중)
t1_y = sc_t + Inches(0.75)
add_rounded_rect(slide4, sc_l + Inches(0.15), t1_y, sc_w - Inches(0.3), Inches(0.85), WHITE, BORDER_GRAY)
add_textbox(slide4, sc_l + Inches(0.3), t1_y + Inches(0.1), sc_w - Inches(0.6), Inches(0.35), "티켓 #2045  |  경기 80사 1234", 11, True, CHARCOAL)
add_textbox(slide4, sc_l + Inches(0.3), t1_y + Inches(0.45), sc_w - Inches(0.6), Inches(0.3), "김덤프 기사  |  양질토  |  15:42 상차", 9, False, SLATE_GRAY)
badge_t1 = add_rounded_rect(slide4, sc_l + sc_w - Inches(1.3), t1_y + Inches(0.1), Inches(1.0), Inches(0.3), MUSTARD)
add_textbox(slide4, sc_l + sc_w - Inches(1.3), t1_y + Inches(0.12), Inches(1.0), Inches(0.25), "운행중", 8, True, CHARCOAL, PP_ALIGN.CENTER)

# 티켓 2 (하차완료)
t2_y = sc_t + Inches(1.7)
add_rounded_rect(slide4, sc_l + Inches(0.15), t2_y, sc_w - Inches(0.3), Inches(0.85), WHITE, BORDER_GRAY)
add_textbox(slide4, sc_l + Inches(0.3), t2_y + Inches(0.1), sc_w - Inches(0.6), Inches(0.35), "티켓 #2044  |  서울 80사 5678", 11, True, CHARCOAL)
add_textbox(slide4, sc_l + Inches(0.3), t2_y + Inches(0.45), sc_w - Inches(0.6), Inches(0.3), "이배차 기사  |  양질토  |  15:30 상차", 9, False, SLATE_GRAY)
badge_t2 = add_rounded_rect(slide4, sc_l + sc_w - Inches(1.3), t2_y + Inches(0.1), Inches(1.0), Inches(0.3), GREEN)
add_textbox(slide4, sc_l + sc_w - Inches(1.3), t2_y + Inches(0.12), Inches(1.0), Inches(0.25), "하차완료", 8, True, WHITE, PP_ALIGN.CENTER)

# 티켓 3 (회차거부)
t3_y = sc_t + Inches(2.65)
add_rounded_rect(slide4, sc_l + Inches(0.15), t3_y, sc_w - Inches(0.3), Inches(0.85), WHITE, BORDER_GRAY)
add_textbox(slide4, sc_l + Inches(0.3), t3_y + Inches(0.1), sc_w - Inches(0.6), Inches(0.35), "티켓 #2043  |  인천 80사 9012", 11, True, CHARCOAL)
add_textbox(slide4, sc_l + Inches(0.3), t3_y + Inches(0.45), sc_w - Inches(0.6), Inches(0.3), "박사토 기사  |  뻘흙    |  15:10 상차", 9, False, SLATE_GRAY)
badge_t3 = add_rounded_rect(slide4, sc_l + sc_w - Inches(1.3), t3_y + Inches(0.1), Inches(1.0), Inches(0.3), RED)
add_textbox(slide4, sc_l + sc_w - Inches(1.3), t3_y + Inches(0.12), Inches(1.0), Inches(0.25), "회차(거부)", 8, True, WHITE, PP_ALIGN.CENTER)


# 파워포인트 파일 저장
output_path = "D:\\Projects\\dumpring\\dumpring-platform-backend\\덤프링_현장담당자_화면상세설계.pptx"
prs.save(output_path)
print(f"PPTX 파일 생성 성공: {output_path}")
